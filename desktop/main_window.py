from __future__ import annotations

import os
import statistics
from datetime import datetime
from typing import Any, Dict, List, Optional

import openpyxl

from PySide6.QtCharts import QChart, QChartView, QLineSeries, QValueAxis
from PySide6.QtCore import QMargins, QRectF, Qt, QTimer, Signal, Slot
from PySide6.QtGui import QColor, QFont, QPainter, QPen
from PySide6.QtWidgets import (
    QApplication,
    QCheckBox,
    QColorDialog,
    QComboBox,
    QFileDialog,
    QFormLayout,
    QGroupBox,
    QHBoxLayout,
    QLabel,
    QListWidget,
    QMainWindow,
    QTreeWidget,
    QTreeWidgetItem,
    QMessageBox,
    QProgressBar,
    QPushButton,
    QSizePolicy,
    QSplitter,
    QSpinBox,
    QStackedWidget,
    QStatusBar,
    QTabWidget,
    QTableWidget,
    QTableWidgetItem,
    QVBoxLayout,
    QWidget,
    QLineEdit,
)

from core.engine import MeasurementEngine
from core.equipment.base import MeasureResult
from core.export import ExcelExporter
from core.gamut_utils import DCI_P3_UV, BT2020_UV, calc_gamut_stats
from .worker import AutoAllWorker, ConnectWorker, LgTvReadWorker, MeasurementWorker, wire_worker_cleanup
from .module_panel import ModulePanel, GammaSubPanel, ColorSubPanel, _CalmanSweepPanel


def _save_all_session(engine: MeasurementEngine) -> str:
    """모든 세션 데이터를 {brand}_{model}_all.xlsx 한 파일에 자동 저장."""
    if not engine.auto_save_dir:
        return ""
    brand = engine.brand or "brand"
    model = engine.model_name or "model"
    path = os.path.join(engine.auto_save_dir, f"{brand}_{model}_all.xlsx")
    try:
        ExcelExporter().export_all_session(
            brand=brand, model=model,
            session_swing=engine.session_swing,
            session_loading=engine.session_loading,
            session_gamut=engine.session_gamut,
            session_contrast=engine.session_contrast,
            file_path=path,
            serial_number=engine.lg_serial_number,
            sw_version=engine.lg_sw_version,
            sw_codename=engine.lg_sw_codename,
        )
        engine.session_key = f"{brand}_{model}"
    except Exception as e:
        print(f"[_save_all_session] 오류: {e}")
    return path

def _resolve_style(tmpl: str, light: bool) -> str:
    import os as _os
    _d = _os.path.dirname(_os.path.abspath(__file__)).replace("\\", "/")
    suffix = "_light" if light else ""
    return (tmpl
            .replace("__ARR_DOWN__", f"{_d}/arr_down{suffix}.svg")
            .replace("__ARR_UP__",   f"{_d}/arr_up{suffix}.svg"))

def _dark_style()  -> str: return _resolve_style(_DARK_STYLE_TMPL,  light=False)
def _light_style() -> str: return _resolve_style(_LIGHT_STYLE_TMPL, light=True)

_DARK_STYLE_TMPL = """
/* ════════════════════════════════════════════════════════════════════
   LED Panel Analyzer — Design System
   Color tokens:
     --deep   : #0d1120   sidebar / deepest
     --base   : #141828   app background
     --card   : #1a2038   card / panel surface
     --raised : #202840   slightly raised surface
     --border : #263058   separator line
     --bdr-a  : #3a52a0   active/focus border
     --txt    : #dde5ff   primary text
     --txt-s  : #7888b8   secondary text
     --txt-m  : #445070   muted / disabled text
     --accent : #4e8df8   primary blue
     --green  : #1dd9a0   success green
     --red    : #f05050   danger red
     --yellow : #f0b040   warning yellow
   ════════════════════════════════════════════════════════════════════ */

/* ── Global ───────────────────────────────────────────────────────── */
QMainWindow, QWidget {
    background: #141828;
    color: #dde5ff;
    font-family: 'Segoe UI', 'Apple SD Gothic Neo', 'Malgun Gothic', sans-serif;
    font-size: 12px;
}

/* ── Cards (GroupBox) ─────────────────────────────────────────────── */
QGroupBox {
    background: #1a2038;
    border: 1px solid #263058;
    border-radius: 10px;
    margin-top: 16px;
    padding: 14px 12px 10px 12px;
}
QGroupBox::title {
    subcontrol-origin: margin;
    left: 14px;
    padding: 0 6px;
    color: #4e8df8;
    font-size: 9px;
    font-weight: bold;
    text-transform: uppercase;
    letter-spacing: 0.12em;
    background: #1a2038;
}

/* ── Buttons ──────────────────────────────────────────────────────── */
QPushButton {
    background: #202840;
    border: 1px solid #3a52a0;
    border-radius: 6px;
    padding: 6px 18px;
    color: #a8b8e8;
    font-weight: 600;
    font-size: 11px;
    min-height: 26px;
}
QPushButton:hover {
    background: #2a3860;
    border-color: #4e8df8;
    color: #dde5ff;
}
QPushButton:pressed { background: #4e8df8; border-color: #3a7ae8; color: #fff; }
QPushButton:disabled { background: #141e30; border-color: #1e2840; color: #3a4868; }

QPushButton#primary {
    background: #4e8df8;
    border: 1px solid #3a7ae8;
    color: #ffffff;
    font-weight: 700;
    letter-spacing: 0.02em;
}
QPushButton#primary:hover   { background: #5e9dff; border-color: #4e8df8; }
QPushButton#primary:pressed { background: #3a7ae8; }
QPushButton#primary:disabled { background: #203060; border-color: #1a2850; color: #506090; }

QPushButton#danger {
    background: #f05050;
    border: 1px solid #d03838;
    color: #ffffff;
    font-weight: 700;
}
QPushButton#danger:hover   { background: #f86060; }
QPushButton#danger:pressed { background: #d03838; }

QPushButton#warning {
    background: #e8a030;
    border: 1px solid #c88018;
    color: #ffffff;
    font-weight: 700;
}
QPushButton#warning:hover { background: #f0b040; }

QPushButton#success {
    background: #1dd9a0;
    border: 1px solid #14b880;
    color: #0a1020;
    font-weight: 700;
}
QPushButton#success:hover { background: #28e8b0; }

/* ── Inputs ───────────────────────────────────────────────────────── */
QComboBox, QLineEdit, QSpinBox, QDoubleSpinBox {
    background: #111828;
    border: 1px solid #2e4080;
    border-radius: 6px;
    padding: 5px 10px;
    color: #dde5ff;
    selection-background-color: #4e8df8;
    min-height: 22px;
}
QComboBox:hover, QLineEdit:hover, QSpinBox:hover, QDoubleSpinBox:hover {
    border-color: #4e8df8;
    background: #151e38;
}
QComboBox:focus, QLineEdit:focus, QSpinBox:focus, QDoubleSpinBox:focus {
    border: 1.5px solid #4e8df8;
    background: #18244a;
    outline: none;
}
QComboBox::drop-down {
    border: none;
    width: 26px;
    border-left: 1px solid #3a52a0;
    background: #253050;
    border-radius: 0 6px 6px 0;
}
QComboBox::drop-down:hover { background: #2e3a68; }
QComboBox::down-arrow {
    image: url(__ARR_DOWN__);
    width: 12px;
    height: 8px;
}
QComboBox QAbstractItemView {
    background: #1a2038;
    border: 1px solid #3a52a0;
    border-radius: 6px;
    selection-background-color: #4e8df8;
    selection-color: #ffffff;
    outline: none;
    padding: 2px;
}
QComboBox QAbstractItemView::item { padding: 5px 10px; border-radius: 4px; }
QSpinBox::up-button, QSpinBox::down-button,
QDoubleSpinBox::up-button, QDoubleSpinBox::down-button {
    background: #253050;
    border: none;
    border-left: 1px solid #3a52a0;
    width: 20px;
    border-radius: 0;
}
QSpinBox::up-button:hover, QSpinBox::down-button:hover,
QDoubleSpinBox::up-button:hover, QDoubleSpinBox::down-button:hover {
    background: #2e3a68;
}
QSpinBox::up-button { border-radius: 0 6px 0 0; }
QSpinBox::down-button { border-radius: 0 0 6px 0; }
QSpinBox::up-arrow, QDoubleSpinBox::up-arrow {
    image: url(__ARR_UP__);
    width: 10px;
    height: 7px;
}
QSpinBox::down-arrow, QDoubleSpinBox::down-arrow {
    image: url(__ARR_DOWN__);
    width: 10px;
    height: 7px;
}

/* ── Progress Bar ─────────────────────────────────────────────────── */
QProgressBar {
    background: #111828;
    border: none;
    border-radius: 10px;
    min-height: 6px;
    max-height: 8px;
    text-align: center;
    color: transparent;
}
QProgressBar::chunk {
    background: qlineargradient(x1:0,y1:0,x2:1,y2:0,
                stop:0.0 #4e8df8,
                stop:0.5 #28d0c8,
                stop:1.0 #1dd9a0);
    border-radius: 10px;
}

/* ── Table ────────────────────────────────────────────────────────── */
QTableWidget {
    background: #141828;
    gridline-color: #1e2a48;
    border: 1px solid #263058;
    border-radius: 6px;
    alternate-background-color: #192238;
    outline: none;
    color: #c8d4f8;
}
QTableWidget::item {
    padding: 5px 10px;
    border: none;
    color: #c0d0f0;
}
QTableWidget::item:hover   { background: #1e3060; }
QTableWidget::item:selected { background: #2a50c8; color: #ffffff; }
QHeaderView { background: transparent; border: none; }
QHeaderView::section {
    background: #192238;
    color: #6878b0;
    border: none;
    border-bottom: 2px solid #263058;
    border-right: 1px solid #1e2a48;
    padding: 6px 10px;
    font-size: 9px;
    font-weight: bold;
    text-transform: uppercase;
    letter-spacing: 0.08em;
}
QHeaderView::section:first { border-radius: 6px 0 0 0; }
QHeaderView::section:last  { border-radius: 0 6px 0 0; border-right: none; }
QHeaderView::section:hover { background: #202e52; color: #a0b8f0; }
QTableCornerButton::section { background: #192238; border: none; }

/* ── CheckBox ─────────────────────────────────────────────────────── */
QCheckBox { spacing: 7px; color: #a8b8e0; }
QCheckBox::indicator {
    width: 15px; height: 15px;
    border-radius: 4px;
    border: 1.5px solid #3a52a0;
    background: #111828;
}
QCheckBox::indicator:checked {
    background: #4e8df8;
    border-color: #4e8df8;
    image: none;
}
QCheckBox::indicator:hover { border-color: #4e8df8; background: #18244a; }

/* ── Splitter ─────────────────────────────────────────────────────── */
QSplitter::handle { background: #1e2a48; }
QSplitter::handle:horizontal { width: 1px; }
QSplitter::handle:vertical   { height: 1px; }
QSplitter::handle:hover { background: #4e8df8; }

/* ── Scrollbar ─────────────────────────────────────────────────────  */
QScrollBar:vertical   { background: transparent; width: 6px; margin: 2px; }
QScrollBar:horizontal { background: transparent; height: 6px; margin: 2px; }
QScrollBar::handle:vertical, QScrollBar::handle:horizontal {
    background: #2e4080;
    border-radius: 3px;
    min-height: 20px; min-width: 20px;
}
QScrollBar::handle:vertical:hover, QScrollBar::handle:horizontal:hover {
    background: #4e8df8;
}
QScrollBar::add-line, QScrollBar::sub-line,
QScrollBar::add-page, QScrollBar::sub-page { background: none; border: none; }

/* ── Tab ──────────────────────────────────────────────────────────── */
QTabWidget::pane {
    border: 1px solid #263058;
    border-radius: 0 8px 8px 8px;
    background: #1a2038;
}
QTabBar::tab {
    background: #111828;
    color: #6878b0;
    border: 1px solid #263058;
    border-bottom: none;
    padding: 6px 16px;
    margin-right: 2px;
    border-radius: 6px 6px 0 0;
    font-size: 11px;
}
QTabBar::tab:selected {
    background: #1a2038;
    color: #dde5ff;
    border-top: 2px solid #4e8df8;
    font-weight: bold;
}
QTabBar::tab:hover { background: #1a2540; color: #b0c0f0; }

/* ── List ─────────────────────────────────────────────────────────── */
QListWidget {
    background: #141828;
    border: 1px solid #263058;
    border-radius: 6px;
    outline: none;
}
QListWidget::item {
    padding: 7px 12px;
    border-bottom: 1px solid #192238;
    color: #a8b8e0;
    font-size: 12px;
}
QListWidget::item:hover  { background: #1e3060; color: #dde5ff; }
QListWidget::item:selected { background: #4e8df8; color: #fff; font-weight: 600; }

/* ── Sidebar Tree ─────────────────────────────────────────────────── */
QTreeWidget {
    background: #0d1120;
    border: none;
    outline: none;
}
QTreeWidget::item {
    padding: 8px 10px;
    border-radius: 0;
    color: #8090c0;
    font-size: 12px;
}
QTreeWidget::item:hover  { background: #18243c; color: #c8d8ff; }
QTreeWidget::item:selected {
    background: #182a50;
    color: #dde5ff;
    border-left: 3px solid #1dd9a0;
    font-weight: 600;
}
QTreeWidget::branch { background: #0d1120; }
QTreeWidget::branch:has-children { border-image: none; image: none; }

/* ── Terminal ─────────────────────────────────────────────────────── */
QPlainTextEdit {
    background: #090e1a;
    color: #3de878;
    border: 1px solid #1e2a48;
    border-radius: 6px;
    font-family: 'Consolas', 'JetBrains Mono', 'Courier New', monospace;
    font-size: 10px;
    selection-background-color: #4e8df8;
    line-height: 1.4;
}

/* ── Tooltip ──────────────────────────────────────────────────────── */
QToolTip {
    background: #1a2038;
    color: #dde5ff;
    border: 1px solid #4e8df8;
    border-radius: 6px;
    padding: 6px 10px;
    font-size: 11px;
}

/* ── Status bar ───────────────────────────────────────────────────── */
QStatusBar {
    background: #0d1120;
    color: #4a5878;
    font-size: 10px;
    border-top: 1px solid #1e2a48;
}

/* ── Named labels ─────────────────────────────────────────────────── */
QLabel#status_ok  { color: #1dd9a0; font-weight: bold; }
QLabel#status_err { color: #f05050; font-weight: bold; }
QLabel#muted      { color: #4a5878; font-size: 11px; }

/* ── Chart views ──────────────────────────────────────────────────── */
QChartView {
    background: #1c2039;
    border: 1px solid #2e3857;
    border-radius: 4px;
}
"""

# ── Light (Soft-Neumorphic) Theme ─────────────────────────────────────────
_LIGHT_STYLE_TMPL = """
/* ════════════════════════════════════════════════════════════════════
   LED Panel Analyzer — Light Soft-UI Theme
   Color tokens:
     --bg     : #e8ecf4   app background
     --card   : #eef1f8   card surface
     --raised : #f3f5fb   top-lit surface
     --inset  : #e0e4ef   recessed / input
     --border : #d0d5e8   separator
     --bdr-a  : #00c9b5   active / teal
     --txt    : #2c3350   primary text
     --txt-s  : #6b7898   secondary text
     --txt-m  : #9ba8c8   muted text
     --accent : #00c9b5   teal
     --red    : #e05050   danger
     --yellow : #e8a030   warning
   ════════════════════════════════════════════════════════════════════ */

/* ── Global ───────────────────────────────────────────────────────── */
QMainWindow, QWidget {
    background: #e8ecf4;
    color: #2c3350;
    font-family: 'Segoe UI', 'Apple SD Gothic Neo', 'Malgun Gothic', sans-serif;
    font-size: 12px;
}

/* ── Cards (GroupBox) ─────────────────────────────────────────────── */
QGroupBox {
    background: #eef1f8;
    border: 1px solid #d0d5e8;
    border-radius: 12px;
    margin-top: 16px;
    padding: 14px 12px 10px 12px;
}
QGroupBox::title {
    subcontrol-origin: margin;
    left: 14px;
    padding: 0 6px;
    color: #00a898;
    font-size: 9px;
    font-weight: bold;
    text-transform: uppercase;
    letter-spacing: 0.12em;
    background: #eef1f8;
}

/* ── Buttons ──────────────────────────────────────────────────────── */
QPushButton {
    background: #eef1f8;
    border: 1px solid #c8cedf;
    border-radius: 8px;
    padding: 6px 18px;
    color: #4a5580;
    font-weight: 600;
    font-size: 11px;
    min-height: 26px;
}
QPushButton:hover {
    background: #e4e8f4;
    border-color: #00c9b5;
    color: #2c3350;
}
QPushButton:pressed { background: #d8dff0; border-color: #00a898; color: #2c3350; }
QPushButton:disabled { background: #eef1f8; border-color: #dde1ee; color: #b0bace; }

QPushButton#primary {
    background: #00c9b5;
    border: 1px solid #00a898;
    color: #ffffff;
    font-weight: 700;
    letter-spacing: 0.02em;
}
QPushButton#primary:hover   { background: #00d9c4; border-color: #00c9b5; }
QPushButton#primary:pressed { background: #00a898; }
QPushButton#primary:disabled { background: #b0e8e2; border-color: #9cdcd6; color: #ffffff; }

QPushButton#danger {
    background: #f05050;
    border: 1px solid #d03838;
    color: #ffffff;
    font-weight: 700;
}
QPushButton#danger:hover   { background: #f86060; }
QPushButton#danger:pressed { background: #d03838; }

QPushButton#warning {
    background: #e8a030;
    border: 1px solid #c88018;
    color: #ffffff;
    font-weight: 700;
}
QPushButton#warning:hover { background: #f0b040; }

QPushButton#success {
    background: #00c9b5;
    border: 1px solid #00a898;
    color: #ffffff;
    font-weight: 700;
}
QPushButton#success:hover { background: #00d9c4; }

/* ── Inputs ───────────────────────────────────────────────────────── */
QComboBox, QLineEdit, QSpinBox, QDoubleSpinBox {
    background: #e0e4ef;
    border: 1px solid #c8cedf;
    border-radius: 8px;
    padding: 5px 10px;
    color: #2c3350;
    selection-background-color: #00c9b5;
    min-height: 22px;
}
QComboBox:hover, QLineEdit:hover, QSpinBox:hover, QDoubleSpinBox:hover {
    border-color: #00c9b5;
    background: #d8dcea;
}
QComboBox:focus, QLineEdit:focus, QSpinBox:focus, QDoubleSpinBox:focus {
    border: 1.5px solid #00c9b5;
    background: #d4d9e8;
    outline: none;
}
QComboBox::drop-down {
    border: none;
    width: 26px;
    border-left: 1px solid #c0c8dc;
    background: #d4d9e8;
    border-radius: 0 8px 8px 0;
}
QComboBox::drop-down:hover { background: #c8d0e4; }
QComboBox::down-arrow {
    image: url(__ARR_DOWN__);
    width: 12px;
    height: 8px;
}
QComboBox QAbstractItemView {
    background: #f0f3fa;
    border: 1px solid #c8cedf;
    border-radius: 8px;
    selection-background-color: #00c9b5;
    selection-color: #ffffff;
    outline: none;
    padding: 2px;
    color: #2c3350;
}
QComboBox QAbstractItemView::item { padding: 5px 10px; border-radius: 4px; }
QSpinBox::up-button, QSpinBox::down-button,
QDoubleSpinBox::up-button, QDoubleSpinBox::down-button {
    background: #d4d9e8;
    border: none;
    border-left: 1px solid #c0c8dc;
    width: 20px;
    border-radius: 0;
}
QSpinBox::up-button:hover, QSpinBox::down-button:hover,
QDoubleSpinBox::up-button:hover, QDoubleSpinBox::down-button:hover {
    background: #c4cce0;
}
QSpinBox::up-button { border-radius: 0 8px 0 0; }
QSpinBox::down-button { border-radius: 0 0 8px 0; }
QSpinBox::up-arrow, QDoubleSpinBox::up-arrow {
    image: url(__ARR_UP__);
    width: 10px;
    height: 7px;
}
QSpinBox::down-arrow, QDoubleSpinBox::down-arrow {
    image: url(__ARR_DOWN__);
    width: 10px;
    height: 7px;
}

/* ── Progress Bar ─────────────────────────────────────────────────── */
QProgressBar {
    background: #d8dcea;
    border: none;
    border-radius: 10px;
    min-height: 6px;
    max-height: 8px;
    text-align: center;
    color: transparent;
}
QProgressBar::chunk {
    background: qlineargradient(x1:0,y1:0,x2:1,y2:0,
                stop:0.0 #00c9b5,
                stop:0.5 #00dfc8,
                stop:1.0 #20e8d0);
    border-radius: 10px;
}

/* ── Table ────────────────────────────────────────────────────────── */
QTableWidget {
    background: #eef1f8;
    gridline-color: #d4d9ea;
    border: 1px solid #d0d5e8;
    border-radius: 8px;
    alternate-background-color: #e8ecf4;
    outline: none;
    color: #2c3350;
}
QTableWidget::item {
    padding: 5px 10px;
    border: none;
    color: #3a4468;
}
QTableWidget::item:hover   { background: #dce8f4; }
QTableWidget::item:selected { background: #00c9b5; color: #ffffff; }
QHeaderView { background: transparent; border: none; }
QHeaderView::section {
    background: #e4e8f4;
    color: #7888a8;
    border: none;
    border-bottom: 2px solid #d0d5e8;
    border-right: 1px solid #d4d9ea;
    padding: 6px 10px;
    font-size: 9px;
    font-weight: bold;
    text-transform: uppercase;
    letter-spacing: 0.08em;
}
QHeaderView::section:first { border-radius: 8px 0 0 0; }
QHeaderView::section:last  { border-radius: 0 8px 0 0; border-right: none; }
QHeaderView::section:hover { background: #d8ddf0; color: #2c3350; }
QTableCornerButton::section { background: #e4e8f4; border: none; }

/* ── CheckBox ─────────────────────────────────────────────────────── */
QCheckBox { spacing: 7px; color: #4a5580; }
QCheckBox::indicator {
    width: 15px; height: 15px;
    border-radius: 4px;
    border: 1.5px solid #c0c8dc;
    background: #e0e4ef;
}
QCheckBox::indicator:checked {
    background: #00c9b5;
    border-color: #00a898;
    image: none;
}
QCheckBox::indicator:hover { border-color: #00c9b5; background: #d4d9e8; }

/* ── Splitter ─────────────────────────────────────────────────────── */
QSplitter::handle { background: #d0d5e8; }
QSplitter::handle:horizontal { width: 1px; }
QSplitter::handle:vertical   { height: 1px; }
QSplitter::handle:hover { background: #00c9b5; }

/* ── Scrollbar ─────────────────────────────────────────────────────  */
QScrollBar:vertical   { background: transparent; width: 6px; margin: 2px; }
QScrollBar:horizontal { background: transparent; height: 6px; margin: 2px; }
QScrollBar::handle:vertical, QScrollBar::handle:horizontal {
    background: #c0c8dc;
    border-radius: 3px;
    min-height: 20px; min-width: 20px;
}
QScrollBar::handle:vertical:hover, QScrollBar::handle:horizontal:hover {
    background: #00c9b5;
}
QScrollBar::add-line, QScrollBar::sub-line,
QScrollBar::add-page, QScrollBar::sub-page { background: none; border: none; }

/* ── Tab ──────────────────────────────────────────────────────────── */
QTabWidget::pane {
    border: 1px solid #d0d5e8;
    border-radius: 0 8px 8px 8px;
    background: #eef1f8;
}
QTabBar::tab {
    background: #e4e8f4;
    color: #7888a8;
    border: 1px solid #d0d5e8;
    border-bottom: none;
    padding: 6px 16px;
    margin-right: 2px;
    border-radius: 6px 6px 0 0;
    font-size: 11px;
}
QTabBar::tab:selected {
    background: #eef1f8;
    color: #2c3350;
    border-top: 2px solid #00c9b5;
    font-weight: bold;
}
QTabBar::tab:hover { background: #e8ecf4; color: #3a4468; }

/* ── List ─────────────────────────────────────────────────────────── */
QListWidget {
    background: #eef1f8;
    border: 1px solid #d0d5e8;
    border-radius: 8px;
    outline: none;
}
QListWidget::item {
    padding: 7px 12px;
    border-bottom: 1px solid #e0e4ef;
    color: #4a5580;
    font-size: 12px;
}
QListWidget::item:hover  { background: #dce8f4; color: #2c3350; }
QListWidget::item:selected { background: #00c9b5; color: #fff; font-weight: 600; }

/* ── Sidebar Tree ─────────────────────────────────────────────────── */
QTreeWidget {
    background: #dde1ee;
    border: none;
    outline: none;
}
QTreeWidget::item {
    padding: 8px 10px;
    border-radius: 0;
    color: #5a6480;
    font-size: 12px;
}
QTreeWidget::item:hover  { background: #d0d8ea; color: #2c3350; }
QTreeWidget::item:selected {
    background: #c8e8e4;
    color: #007a70;
    border-left: 3px solid #00c9b5;
    font-weight: 600;
}
QTreeWidget::branch { background: #dde1ee; }
QTreeWidget::branch:has-children { border-image: none; image: none; }

/* ── Terminal (LG serial log) ────────────────────────────────────── */
QPlainTextEdit {
    background: #f5f8ff;
    color: #007a70;
    border: 1px solid #d0d5e8;
    border-radius: 8px;
    font-family: 'Consolas', 'JetBrains Mono', 'Courier New', monospace;
    font-size: 10px;
    selection-background-color: #00c9b5;
    line-height: 1.4;
}

/* ── Tooltip ──────────────────────────────────────────────────────── */
QToolTip {
    background: #eef1f8;
    color: #2c3350;
    border: 1px solid #00c9b5;
    border-radius: 6px;
    padding: 6px 10px;
    font-size: 11px;
}

/* ── Status bar ───────────────────────────────────────────────────── */
QStatusBar {
    background: #dde1ee;
    color: #7888a8;
    font-size: 10px;
    border-top: 1px solid #d0d5e8;
}

/* ── Named labels ─────────────────────────────────────────────────── */
QLabel#status_ok  { color: #00a898; font-weight: bold; }
QLabel#status_err { color: #e05050; font-weight: bold; }
QLabel#muted      { color: #8898b8; font-size: 11px; }

/* ── Chart views ──────────────────────────────────────────────────── */
QChartView {
    background: #eef1f8;
    border: 1px solid #c8d0e0;
    border-radius: 4px;
}
"""

_DOT_STYLES = {
    "disconnected": "background:#e74c3c;",
    "connecting":   "background:#f39c12;",
    "connected":    "background:#1a9e50;",
    "error":        "background:#c0392b;",
}
_DOT_BASE = (
    "border-radius:6px;"
    "border:1px solid rgba(0,0,0,0.35);"
    "min-width:12px;max-width:12px;min-height:12px;max-height:12px;"
)


class _StatusDot(QLabel):
    """12×12 colored circle indicator. Replaces verbose status text labels."""

    def __init__(self, parent: "QWidget | None" = None) -> None:
        super().__init__(parent)
        self.setFixedSize(12, 12)
        self.set_state("disconnected")

    def set_state(self, state: str, tooltip: str = "") -> None:
        color_css = _DOT_STYLES.get(state, _DOT_STYLES["disconnected"])
        self.setStyleSheet(_DOT_BASE + color_css)
        labels = {
            "disconnected": "Disconnected",
            "connecting":   "Connecting…",
            "connected":    "Connected",
            "error":        "Error",
        }
        tip = labels.get(state, state)
        if tooltip:
            tip += f"  ({tooltip})"
        self.setToolTip(tip)


class ConnectionPanel(QGroupBox):
    def __init__(self, engine: MeasurementEngine, parent: QWidget | None = None) -> None:
        super().__init__("Equipment", parent)
        self._engine = engine
        self._connect_worker: Optional[ConnectWorker] = None
        self._reset_worker: Optional[ConnectWorker] = None
        self._lg_read_worker: Optional[LgTvReadWorker] = None
        self._build_ui()
        self._engine.on_lg_device_info = self._apply_lg_device_info

    def _build_ui(self) -> None:
        root = QVBoxLayout(self)
        root.setSpacing(4)
        root.setContentsMargins(6, 6, 6, 4)

        # ── 상단 행: 기본 정보 / CA 색채휘도계 / 패턴 제너레이터 ──────────────
        top_row = QHBoxLayout()
        top_row.setSpacing(6)

        # Brand / Model
        info_box = QGroupBox("Basic Info")
        info_form = QFormLayout(info_box)
        info_form.setContentsMargins(8, 6, 8, 6)
        info_form.setSpacing(8)
        info_form.setHorizontalSpacing(10)
        info_form.setFieldGrowthPolicy(QFormLayout.FieldGrowthPolicy.ExpandingFieldsGrow)
        self._brand_edit = QLineEdit("AAA")
        self._model_edit = QLineEdit("BBB")
        self._brand_edit.setMinimumWidth(160)
        self._model_edit.setMinimumWidth(160)
        self._brand_edit.textChanged.connect(self._sync_info)
        self._model_edit.textChanged.connect(self._sync_info)
        self._brand_edit.editingFinished.connect(self._check_session_key_change)
        self._model_edit.editingFinished.connect(self._check_session_key_change)
        info_form.addRow("Brand:", self._brand_edit)
        info_form.addRow("Model:", self._model_edit)
        top_row.addWidget(info_box, stretch=2)

        # 포트 스캔 공통 버튼
        btn_scan = QPushButton("Scan Ports")
        btn_scan.setToolTip("Re-scan connected serial ports")
        btn_scan.clicked.connect(self._scan_ports)

        # Meter — 좌: 포트/장비 상하, 우: 버튼 원위치
        meter_box = QGroupBox("CA Colorimeter")
        meter_main = QHBoxLayout(meter_box)
        meter_main.setContentsMargins(8, 6, 8, 6)
        meter_main.setSpacing(10)

        meter_form = QFormLayout()
        meter_form.setContentsMargins(0, 0, 0, 0)
        meter_form.setSpacing(4)
        meter_form.setHorizontalSpacing(8)
        meter_form.setFieldGrowthPolicy(QFormLayout.FieldGrowthPolicy.ExpandingFieldsGrow)
        self._meter_port = QComboBox()
        self._meter_port.setEditable(True)
        self._meter_model = QComboBox()
        self._meter_model.addItems(["CA-410", "CA-310"])
        self._meter_model.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Fixed)
        self._meter_model.currentTextChanged.connect(self._on_meter_model_changed)
        meter_form.addRow("Port:", self._meter_port)
        meter_form.addRow("Device:", self._meter_model)
        meter_main.addLayout(meter_form, stretch=1)

        self._btn_meter = QPushButton("Connect")
        self._btn_meter.setObjectName("primary")
        self._btn_meter.clicked.connect(self._connect_meter)
        self._btn_meter_dis = QPushButton("Disconnect")
        self._btn_meter_dis.setObjectName("danger")
        self._btn_meter_dis.setEnabled(False)
        self._btn_meter_dis.clicked.connect(self._disconnect_meter)
        self._meter_status = _StatusDot()
        meter_main.addWidget(self._btn_meter)
        meter_main.addWidget(self._btn_meter_dis)
        meter_main.addWidget(self._meter_status)
        top_row.addWidget(meter_box, stretch=2)

        # Generator — 좌: 포트/장비 상하, 우: 버튼 원위치
        gen_box = QGroupBox("Pattern Generator")
        gen_main = QHBoxLayout(gen_box)
        gen_main.setContentsMargins(8, 6, 8, 6)
        gen_main.setSpacing(10)

        gen_form = QFormLayout()
        gen_form.setContentsMargins(0, 0, 0, 0)
        gen_form.setSpacing(4)
        gen_form.setHorizontalSpacing(8)
        gen_form.setFieldGrowthPolicy(QFormLayout.FieldGrowthPolicy.ExpandingFieldsGrow)
        self._gen_port = QComboBox()
        self._gen_port.setEditable(True)
        self._gen_model = QComboBox()
        self._gen_model.addItems(["VG-879", "VG-876"])
        self._gen_model.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Fixed)
        gen_form.addRow("Port:", self._gen_port)
        gen_form.addRow("Device:", self._gen_model)
        gen_main.addLayout(gen_form, stretch=1)

        self._btn_gen = QPushButton("Connect")
        self._btn_gen.setObjectName("primary")
        self._btn_gen.clicked.connect(self._connect_generator)
        self._btn_gen_dis = QPushButton("Disconnect")
        self._btn_gen_dis.setObjectName("danger")
        self._btn_gen_dis.setEnabled(False)
        self._btn_gen_dis.clicked.connect(self._disconnect_generator)
        self._btn_gen_reset = QPushButton("Reset Device")
        self._btn_gen_reset.setObjectName("warning")
        self._btn_gen_reset.setEnabled(False)
        self._btn_gen_reset.setToolTip("Recover from freeze: re-enter ENQ and reload colorbar")
        self._btn_gen_reset.clicked.connect(self._reset_generator)
        self._gen_status = _StatusDot()
        gen_main.addWidget(self._btn_gen)
        gen_main.addWidget(self._btn_gen_dis)
        gen_main.addWidget(self._btn_gen_reset)
        gen_main.addWidget(self._gen_status)
        top_row.addWidget(gen_box, stretch=2)

        root.addLayout(top_row)

        # ── 하단 행: 포트 스캔 / 자동 저장 폴더 / 전체 해제 ──────────────
        bottom_row = QHBoxLayout()
        bottom_row.setSpacing(4)
        self._btn_dis_all = QPushButton("Disconnect All")
        self._btn_dis_all.setObjectName("danger")
        self._btn_dis_all.clicked.connect(self._disconnect_all)
        bottom_row.addWidget(btn_scan)
        bottom_row.addSpacing(12)
        bottom_row.addWidget(QLabel("Auto-save Folder:"))
        self._save_dir_edit = QLineEdit()
        self._save_dir_edit.setPlaceholderText("Folder to auto-save completed measurements")
        self._save_dir_edit.setReadOnly(True)
        self._save_dir_edit.setMinimumWidth(220)
        bottom_row.addWidget(self._save_dir_edit, stretch=1)
        btn_folder = QPushButton("📁")
        btn_folder.setFixedWidth(32)
        btn_folder.clicked.connect(self._pick_save_dir)
        bottom_row.addWidget(btn_folder)
        bottom_row.addSpacing(12)
        bottom_row.addWidget(self._btn_dis_all)
        root.addLayout(bottom_row)

        # 시작 시 포트 목록 채우기 + 자동 저장 폴더 초기값 = 실행 폴더
        self._scan_ports()
        default_dir = os.path.dirname(os.path.abspath(__file__))
        self._engine.auto_save_dir = default_dir
        self._save_dir_edit.setText(default_dir)

    def _on_meter_model_changed(self, model: str) -> None:
        """CA-310 선택 시 포트 목록을 'usb' 하나만 표시."""
        if model == "CA-310":
            self._meter_port.clear()
            self._meter_port.addItem("usb")
            self._meter_port.setCurrentText("usb")
        else:
            self._meter_port.clear()
            self._repopulate_meter_ports()

    def _repopulate_meter_ports(self) -> None:
        import serial.tools.list_ports
        ports = sorted(p.device for p in serial.tools.list_ports.comports())
        self._meter_port.addItems(ports)

    def _scan_ports(self) -> None:
        """시리얼 포트를 스캔해 콤보박스들을 갱신한다."""
        import serial.tools.list_ports
        ports = sorted(p.device for p in serial.tools.list_ports.comports())

        # CA-310 선택 중이면 미터 포트는 'usb' 고정 — 스캔 목록으로 덮지 않음
        if self._meter_model.currentText() != "CA-310":
            current = self._meter_port.currentText()
            self._meter_port.clear()
            self._meter_port.addItems(ports)
            idx = self._meter_port.findText(current)
            if idx >= 0:
                self._meter_port.setCurrentIndex(idx)

        current_gen = self._gen_port.currentText()
        self._gen_port.clear()
        self._gen_port.addItems(ports)
        idx = self._gen_port.findText(current_gen)
        if idx >= 0:
            self._gen_port.setCurrentIndex(idx)

    def _pick_save_dir(self) -> None:
        folder = QFileDialog.getExistingDirectory(
            self, "Select Auto-save Folder",
            self._engine.auto_save_dir or os.path.expanduser("~"),
        )
        if folder:
            self._engine.auto_save_dir = folder
            self._save_dir_edit.setText(folder)

    def _sync_info(self) -> None:
        self._engine.brand = self._brand_edit.text().strip()
        self._engine.model_name = self._model_edit.text().strip()

    def _check_session_key_change(self) -> None:
        """Brand/Model 입력 확정 시 세션 키가 바뀌었으면 기존 데이터를 초기화."""
        current_key = f"{self._engine.brand}_{self._engine.model_name}"
        saved_key = self._engine.session_key
        if not saved_key or current_key == saved_key:
            return
        has_data = any([
            self._engine.session_swing,
            self._engine.session_loading,
            self._engine.session_gamut,
            self._engine.session_contrast,
        ])
        if not has_data:
            return
        ret = QMessageBox.question(
            self, "새 모델 감지",
            f"기본 정보가 변경되었습니다 ({saved_key} → {current_key}).\n"
            "기존 세션 데이터를 초기화하고 새로 시작하시겠습니까?",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No,
        )
        if ret == QMessageBox.StandardButton.Yes:
            self._engine.session_swing.clear()
            self._engine.session_loading.clear()
            self._engine.session_gamut.clear()
            self._engine.session_contrast.clear()
            self._engine.session_key = ""

    def _apply_lg_device_info(self, brand: str, model: str) -> None:
        """luna 응답에서 파싱된 Brand/Model을 UI에 반영."""
        self._brand_edit.setText(brand)
        self._model_edit.setText(model)

    def _connect_meter(self) -> None:
        port = self._meter_port.currentText()
        model = self._meter_model.currentText()
        self._meter_status.set_state("connecting")
        self._btn_meter.setEnabled(False)
        self._connect_worker = ConnectWorker(
            lambda: self._engine.connect_meter(port, model)
        )
        self._connect_worker.succeeded.connect(self._on_meter_connected)
        self._connect_worker.error.connect(lambda msg: self._on_connect_error("CA Connection Error", msg, self._btn_meter))
        wire_worker_cleanup(self._connect_worker, self, '_connect_worker')
        self._connect_worker.start()

    def _on_meter_connected(self) -> None:
        self._sync_info()
        ident = getattr(self._engine.meter, "ident", None)
        self._meter_status.set_state("connected", tooltip=ident or "")
        self._btn_meter.setEnabled(False)
        self._btn_meter_dis.setEnabled(True)

    def _disconnect_meter(self) -> None:
        try:
            if self._engine.meter and self._engine.meter.is_connected:
                self._engine.meter.disconnect()
            self._engine.meter = None
        except Exception:
            pass
        self._meter_status.set_state("disconnected")
        self._btn_meter.setEnabled(True)
        self._btn_meter_dis.setEnabled(False)

    def _connect_generator(self) -> None:
        port = self._gen_port.currentText()
        model = self._gen_model.currentText()
        self._gen_status.set_state("connecting")
        self._btn_gen.setEnabled(False)
        self._connect_worker = ConnectWorker(lambda: self._engine.connect_generator(port, model))
        self._connect_worker.succeeded.connect(self._on_gen_connected)
        self._connect_worker.error.connect(lambda msg: self._on_connect_error("VG Connection Error", msg, self._btn_gen))
        wire_worker_cleanup(self._connect_worker, self, '_connect_worker')
        self._connect_worker.start()

    def _on_gen_connected(self) -> None:
        self._gen_status.set_state("connected")
        self._btn_gen.setEnabled(False)
        self._btn_gen_dis.setEnabled(True)
        self._btn_gen_reset.setEnabled(True)

    def _on_connect_error(self, title: str, msg: str, retry_btn: "QPushButton") -> None:
        retry_btn.setEnabled(True)
        QMessageBox.critical(self, title, msg)

    def _reset_generator(self) -> None:
        """장비 freeze 시 ENQ 재진입 + EXPDN4(2286,0) 컬러바 복귀."""
        gen = self._engine.generator
        if gen is None or not gen.is_connected:
            return
        self._btn_gen_reset.setEnabled(False)
        self._gen_status.set_state("connecting")
        worker = ConnectWorker(lambda: gen.reset())
        worker.succeeded.connect(self._on_gen_reset_done)
        worker.error.connect(lambda msg: (
            QMessageBox.critical(self, "VG Reset Error", msg),
            self._btn_gen_reset.setEnabled(True),
            self._gen_status.set_state("connected"),
        ))
        wire_worker_cleanup(worker, self, '_reset_worker')
        worker.start()
        self._reset_worker = worker

    def _on_gen_reset_done(self) -> None:
        self._gen_status.set_state("connected")
        self._btn_gen_reset.setEnabled(True)

    def _disconnect_generator(self) -> None:
        try:
            if self._engine.generator and self._engine.generator.is_connected:
                self._engine.generator.disconnect()
            self._engine.generator = None
        except Exception:
            pass
        self._gen_status.set_state("disconnected")
        self._btn_gen.setEnabled(True)
        self._btn_gen_dis.setEnabled(False)
        self._btn_gen_reset.setEnabled(False)

    def _disconnect_all(self) -> None:
        self._disconnect_meter()
        self._disconnect_generator()
        self._disconnect_lg_serial()

    def _disconnect_lg_serial(self) -> None:
        """LG TV 시리얼 연결 해제 (engine.lg_tv_serial 닫기)."""
        ser = self._engine.lg_tv_serial
        if ser:
            try:
                if ser.is_open:
                    ser.write(b"exit\r\n")
                    import time; time.sleep(0.2)
                    ser.close()
            except Exception:
                pass
            self._engine.lg_tv_serial = None
            self._engine.lg_log_tx = None


# ---------------------------------------------------------------------------
# Center Alignment Panel
# ---------------------------------------------------------------------------

class CenterAlignPanel(QWidget):
    def __init__(self, engine: MeasurementEngine, parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self._engine = engine
        self._worker: Optional[MeasurementWorker] = None
        self._just_scan_worker: Optional[ConnectWorker] = None
        layout = QVBoxLayout(self)

        title = QLabel("🎯 Center Align")
        title.setStyleSheet("font-size:15px;font-weight:bold;")
        layout.addWidget(title)
        desc = QLabel("Outputs the ABC center alignment pattern from the Pattern Generator.\n"
                      "Point the measurement lens at the screen center, then click [OK].")
        desc.setObjectName("muted")
        desc.setWordWrap(True)
        layout.addWidget(desc)

        notice = QLabel(
            "⚠  Check TV settings before outputting pattern\n"
            "   •  Aspect Ratio :  Original\n"
            "   •  Just Scan      :  On"
        )
        notice.setStyleSheet(
            "background:#fff8e1; border:1px solid #f0c040; border-radius:5px;"
            "padding:8px 12px; color:#7a5800; font-size:12px;"
        )
        layout.addWidget(notice)

        # ── 참고 이미지 (x.png) ──────────────────────────────────────────
        import os as _os
        import sys as _sys
        _base = _sys._MEIPASS if getattr(_sys, "frozen", False) else _os.path.join(_os.path.dirname(__file__), "..")  # type: ignore[attr-defined]
        _img_path = _os.path.join(_base, "x.png")
        if _os.path.isfile(_img_path):
            from PySide6.QtGui import QPixmap as _QPixmap
            _pix = _QPixmap(_img_path)
            if not _pix.isNull():
                img_label = QLabel()
                img_label.setPixmap(
                    _pix.scaledToWidth(480, Qt.TransformationMode.SmoothTransformation)
                )
                img_label.setAlignment(Qt.AlignmentFlag.AlignLeft)
                layout.addWidget(img_label)

        # ── 버튼 행 ──────────────────────────────────────────────────────
        btn_row = QHBoxLayout()
        self._btn_start = QPushButton("▶  Output Pattern")
        self._btn_start.setObjectName("primary")
        self._btn_start.clicked.connect(self._start)
        self._btn_ok = QPushButton("✔  OK — Confirm Center")
        self._btn_ok.setObjectName("success")
        self._btn_ok.setEnabled(False)
        self._btn_ok.clicked.connect(self._confirm)
        btn_row.addWidget(self._btn_start)
        btn_row.addWidget(self._btn_ok)
        btn_row.addStretch()
        layout.addLayout(btn_row)

        # ── Just Scan ON ──────────────────────────────────────────────────
        js_row = QHBoxLayout()
        self._btn_just_scan = QPushButton("📺  Just Scan  ON")
        self._btn_just_scan.setObjectName("success")
        self._btn_just_scan.setToolTip(
            "LG TV 시리얼 연결 후 Just Scan을 ON으로 설정합니다\n"
            "(LG TV Serial 패널에서 먼저 연결하세요)"
        )
        self._btn_just_scan.clicked.connect(self._just_scan_on)
        js_row.addWidget(self._btn_just_scan)
        js_row.addStretch()
        layout.addLayout(js_row)

        self._status = QLabel("Idle")
        self._status.setObjectName("muted")
        layout.addWidget(self._status)
        layout.addStretch()

    def _start(self) -> None:
        if self._engine.generator is None or not self._engine.generator.is_connected:
            QMessageBox.critical(self, "Generator Not Connected",
                                 "Pattern Generator is not connected.\n"
                                 "Select the VG port in the Connection panel and click [Connect].")
            return
        self._btn_start.setEnabled(False)
        self._status.setText("Displaying pattern…")
        self._worker = MeasurementWorker(self._engine, "center_align")
        self._worker.succeeded.connect(self._on_ready)
        self._worker.error.connect(self._on_error)
        wire_worker_cleanup(self._worker, self, '_worker')
        self._worker.start()

    def _on_ready(self, _result: Any) -> None:
        self._status.setText("Pattern displayed — confirm center, then click OK")
        self._btn_ok.setEnabled(True)
        self._btn_start.setEnabled(True)

    def _on_error(self, msg: str) -> None:
        self._status.setText(f"Error: {msg}")
        self._btn_start.setEnabled(True)
        QMessageBox.critical(self, "Center Alignment Error", msg)

    def _confirm(self) -> None:
        self._status.setText("✔  Center confirmed. Proceed to the next step.")
        self._btn_ok.setEnabled(False)

    def _just_scan_on(self) -> None:
        ser = self._engine.lg_tv_serial
        if not ser or not ser.is_open:
            QMessageBox.warning(self, "LG TV Not Connected",
                                "LG TV Serial 패널에서 TV에 먼저 연결하세요.")
            return
        cmds: list[str] = [
            """luna-send -n 1 -f luna://com.webos.service.settingsservice/setSystemSettings '{"category":"picture","settings":{"pictureMode":"vivid"}}'""",
            """luna-send -n 1 -f luna://com.webos.service.networkinput/sendSpecialKey '{"key":"ENTER"}'""",
            """luna-send -n 1 -f luna://com.webos.service.settingsservice/setSystemSettings '{"category":"aspectRatio","settings":{"arcPerApp":"original"}}'""",
            """luna-send -n 1 -f luna://com.webos.service.settingsservice/setSystemSettings '{"category":"aspectRatio","settings":{"justScan":"on"}}'""",
        ]
        self._btn_just_scan.setEnabled(False)
        self._status.setText("Sending Just Scan ON commands…")

        def _send_all() -> None:
            import time
            for cmd in cmds:
                ser.write((cmd + "\n").encode("utf-8"))
                if self._engine.lg_log_tx:
                    self._engine.lg_log_tx(cmd)
                time.sleep(0.8)

        def _on_done() -> None:
            self._btn_just_scan.setEnabled(True)
            self._status.setText("Just Scan ON commands sent.")
            QMessageBox.information(
                self, "Manual Energy Saving Setting Required",
                "⚠️  Energy Saving may be automatically restored due to TV policy.\n\n"
                "Please check and set it to OFF directly from the TV menu.\n\n"
                "  Setting → System → Energy Saving → Off",
            )

        def _on_err(msg: str) -> None:
            self._btn_just_scan.setEnabled(True)
            self._status.setText(f"Error: {msg}")

        worker = ConnectWorker(_send_all)
        worker.succeeded.connect(_on_done)
        worker.error.connect(_on_err)
        wire_worker_cleanup(worker, self, '_just_scan_worker')
        self._just_scan_worker = worker
        worker.start()


# ---------------------------------------------------------------------------
# LG TV Serial Panel
# ---------------------------------------------------------------------------

class LgTvPanel(QWidget):
    """LG TV 시리얼 연결 · 터미널 · luna 명령 전송 패널."""

    def __init__(self, engine: MeasurementEngine, parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self._engine = engine
        self._lg_read_worker:   Optional[LgTvReadWorker] = None
        self._just_scan_worker: Optional[ConnectWorker]  = None
        self._fetch_info_worker: Optional[ConnectWorker] = None
        self._lg_json_buf: str = ""

        from PySide6.QtWidgets import QPlainTextEdit
        layout = QVBoxLayout(self)

        title = QLabel("📺 LG TV Serial")
        title.setStyleSheet("font-size:15px;font-weight:bold;")
        layout.addWidget(title)

        # ── 포트 / Baud / 연결 버튼 ──────────────────────────────────────
        conn_row = QHBoxLayout()
        conn_row.setSpacing(4)
        conn_row.addWidget(QLabel("Port:"))
        self._lg_port = QComboBox()
        self._lg_port.setEditable(True)
        self._lg_port.setFixedWidth(130)
        conn_row.addWidget(self._lg_port)
        btn_scan = QPushButton("Scan")
        btn_scan.setMinimumWidth(52)
        btn_scan.setToolTip("Scan ports")
        btn_scan.clicked.connect(self._lg_scan_ports)
        conn_row.addWidget(btn_scan)
        conn_row.addWidget(QLabel("Baud:"))
        self._lg_baud = QComboBox()
        self._lg_baud.addItems(["9600", "19200", "38400", "57600", "115200"])
        self._lg_baud.setCurrentText("115200")
        self._lg_baud.setMinimumWidth(88)
        conn_row.addWidget(self._lg_baud)
        self._btn_lg = QPushButton("Connect")
        self._btn_lg.setObjectName("primary")
        self._btn_lg.setMinimumWidth(88)
        self._btn_lg.clicked.connect(self._connect_lg)
        conn_row.addWidget(self._btn_lg)
        self._btn_lg_dis = QPushButton("Disconnect")
        self._btn_lg_dis.setObjectName("danger")
        self._btn_lg_dis.setMinimumWidth(96)
        self._btn_lg_dis.setEnabled(False)
        self._btn_lg_dis.clicked.connect(self._disconnect_lg)
        conn_row.addWidget(self._btn_lg_dis)
        self._lg_status = _StatusDot()
        conn_row.addWidget(self._lg_status)
        conn_row.addStretch()
        layout.addLayout(conn_row)

        # ── 수신 터미널 ──────────────────────────────────────────────────
        term_row = QHBoxLayout()
        term_row.setSpacing(4)
        self._lg_terminal = QPlainTextEdit()
        self._lg_terminal.setReadOnly(True)
        self._lg_terminal.setMinimumHeight(200)
        self._lg_terminal.setPlaceholderText("LG TV received data")
        self._lg_terminal.setStyleSheet(
            "background:#111828;color:#78d878;"
            "font-family:'Consolas','Courier New',monospace;"
            "font-size:11px;border:1px solid #2e3857;border-radius:3px;"
        )
        btn_clear = QPushButton("Clear")
        btn_clear.setFixedWidth(72)
        btn_clear.clicked.connect(self._lg_terminal.clear)
        term_row.addWidget(self._lg_terminal, stretch=1)
        term_row.addWidget(btn_clear)
        layout.addLayout(term_row)

        # ── 명령어 전송 ──────────────────────────────────────────────────
        send_row = QHBoxLayout()
        send_row.setSpacing(4)
        send_row.addWidget(QLabel("TX:"))
        self._lg_cmd_edit = QLineEdit()
        self._lg_cmd_edit.setPlaceholderText("Enter luna:// command and press Enter")
        self._lg_cmd_edit.returnPressed.connect(self._lg_send_cmd)
        send_row.addWidget(self._lg_cmd_edit, stretch=1)
        self._lg_eol_combo = QComboBox()
        self._lg_eol_combo.addItems(["\\n", "\\r\\n", "\\r", "None"])
        self._lg_eol_combo.setFixedWidth(62)
        self._lg_eol_combo.setToolTip("Line ending (EOL)")
        send_row.addWidget(self._lg_eol_combo)
        btn_send = QPushButton("Send")
        btn_send.setObjectName("primary")
        btn_send.setFixedWidth(52)
        btn_send.clicked.connect(self._lg_send_cmd)
        send_row.addWidget(btn_send)
        layout.addLayout(send_row)

        # ── TV 설정 바로가기 버튼 ────────────────────────────────────────
        preset_row = QHBoxLayout()
        btn_just_scan = QPushButton("Just Scan  ON")
        btn_just_scan.setObjectName("success")
        btn_just_scan.setToolTip("Sets Just Scan to ON")
        btn_just_scan.clicked.connect(self._lg_just_scan_on)
        preset_row.addWidget(btn_just_scan)
        btn_ui_reset = QPushButton("UI Reset")
        btn_ui_reset.setToolTip("resetPictureSettings — Resets PQ settings")
        btn_ui_reset.clicked.connect(self._lg_ui_reset)
        preset_row.addWidget(btn_ui_reset)
        btn_set_info = QPushButton("Get Info")
        btn_set_info.setToolTip("Reads model, serial, and SW version from the TV")
        btn_set_info.clicked.connect(self._lg_fetch_device_info)
        preset_row.addWidget(btn_set_info)
        preset_row.addStretch()
        layout.addLayout(preset_row)
        layout.addStretch()

    # ── LG 메서드 (기존 CenterAlignPanel에서 이동) ──────────────────────

    def _lg_ui_reset(self) -> None:
        ser = self._engine.lg_tv_serial
        if not ser or not ser.is_open:
            self._lg_terminal.appendPlainText("[Error] LG TV is not connected.")
            return
        cmd = "luna-send -n 1 -f luna://com.webos.service.pqcontroller/resetPictureSettings '{}'"
        try:
            ser.write((cmd + "\n").encode("utf-8"))
            self._lg_terminal.appendPlainText(f"[TX] {cmd}")
        except Exception as exc:
            self._lg_terminal.appendPlainText(f"[TX Error] {exc}")

    def _lg_just_scan_on(self) -> None:
        ser = self._engine.lg_tv_serial
        if not ser or not ser.is_open:
            self._lg_terminal.appendPlainText("[Error] LG TV is not connected.")
            return
        cmds: list[str] = [
            """luna-send -n 1 -f luna://com.webos.service.settingsservice/setSystemSettings '{"category":"picture","settings":{"pictureMode":"vivid"}}'""",
            """luna-send -n 1 -f luna://com.webos.service.networkinput/sendSpecialKey '{"key":"ENTER"}'""",
            """luna-send -n 1 -f luna://com.webos.service.settingsservice/setSystemSettings '{"category":"aspectRatio","settings":{"arcPerApp":"original"}}'""",
            """luna-send -n 1 -f luna://com.webos.service.settingsservice/setSystemSettings '{"category":"aspectRatio","settings":{"justScan":"on"}}'""",
        ]
        def _send_all() -> None:
            import time
            for cmd in cmds:
                ser.write((cmd + "\n").encode("utf-8"))
                time.sleep(0.8)
        def _on_done() -> None:
            for cmd in cmds:
                self._lg_terminal.appendPlainText(f"[TX] {cmd}")
            QMessageBox.information(
                self, "Manual Energy Saving Setting Required",
                "⚠️  Energy Saving may be automatically restored due to TV policy.\n\n"
                "Please check and set it to OFF directly from the TV menu.\n\n"
                "  Setting → System → Energy Saving → Off",
            )
        worker = ConnectWorker(_send_all)
        worker.succeeded.connect(_on_done)
        worker.error.connect(lambda msg: self._lg_terminal.appendPlainText(f"[TX Error] {msg}"))
        wire_worker_cleanup(worker, self, '_just_scan_worker')
        self._just_scan_worker = worker
        worker.start()

    def _lg_fetch_device_info(self) -> None:
        ser = self._engine.lg_tv_serial
        if not ser or not ser.is_open:
            self._lg_terminal.appendPlainText("[Error] LG TV is not connected.")
            return
        cmds = [
            'luna-send -n 1 -f luna://com.webos.service.tv.systemproperty/getSystemInfo \'{"keys":["modelName","serialNumber"]}\'',
            'luna-send -n 1 -f luna://com.palm.systemservice/osInfo/query \'{"parameters" : ["core_os_release","core_os_release_codename"]}\'',
        ]
        def _send_all() -> None:
            import time
            for cmd in cmds:
                ser.write((cmd + "\n").encode("utf-8"))
                time.sleep(1.0)
        def _on_done() -> None:
            for cmd in cmds:
                self._lg_terminal.appendPlainText(f"[TX] {cmd}")
        worker = ConnectWorker(_send_all)
        worker.succeeded.connect(_on_done)
        worker.error.connect(lambda msg: self._lg_terminal.appendPlainText(f"[TX Error] {msg}"))
        wire_worker_cleanup(worker, self, '_fetch_info_worker')
        self._fetch_info_worker = worker
        worker.start()

    def _lg_scan_ports(self) -> None:
        import serial.tools.list_ports
        ports = sorted(p.device for p in serial.tools.list_ports.comports())
        current = self._lg_port.currentText()
        self._lg_port.clear()
        self._lg_port.addItems(ports)
        idx = self._lg_port.findText(current)
        if idx >= 0:
            self._lg_port.setCurrentIndex(idx)

    def _connect_lg(self) -> None:
        import serial, time
        port = self._lg_port.currentText().strip()
        baud = int(self._lg_baud.currentText())
        if not port:
            return
        self._lg_status.set_state("connecting")
        self._btn_lg.setEnabled(False)
        def _init():
            ser = serial.Serial(port, baudrate=baud, timeout=1.0)
            self._engine.lg_tv_serial = ser
            time.sleep(0.3)
            ser.write(b"debug\r\n")
            time.sleep(0.5)
            ser.write(b"s\r\n")
            time.sleep(0.5)
        worker = ConnectWorker(_init)
        worker.succeeded.connect(self._on_lg_connected)
        worker.error.connect(self._on_lg_connect_error)
        wire_worker_cleanup(worker, self, '_lg_init_worker')
        self._lg_init_worker = worker
        worker.start()

    def _on_lg_connected(self) -> None:
        port = self._lg_port.currentText().strip()
        self._lg_status.set_state("connected", tooltip=port)
        self._btn_lg.setEnabled(False)
        self._btn_lg_dis.setEnabled(True)
        self._lg_terminal.appendPlainText("[TX] debug")
        self._lg_terminal.appendPlainText("[TX] s")
        self._lg_terminal.appendPlainText("── Init complete, ready for luna commands ──")
        self._engine.lg_log_tx = lambda msg: self._lg_terminal.appendPlainText(f"[TX] {msg}")
        self._lg_read_worker = LgTvReadWorker(self._engine.lg_tv_serial)
        self._lg_read_worker.data_received.connect(self._on_lg_data)
        wire_worker_cleanup(self._lg_read_worker, self, '_lg_read_worker')
        self._lg_read_worker.start()
        self._lg_auto_init()

    def _lg_auto_init(self) -> None:
        ser = self._engine.lg_tv_serial
        if not ser or not ser.is_open:
            return
        cmds = [
            'luna-send -n 1 -f luna://com.webos.service.tv.systemproperty/getSystemInfo \'{"keys":["modelName","serialNumber"]}\'',
            'luna-send -n 1 -f luna://com.palm.systemservice/osInfo/query \'{"parameters" : ["core_os_release","core_os_release_codename"]}\'',
        ]
        def _send_next(idx: int) -> None:
            if idx >= len(cmds):
                return
            cmd = cmds[idx]
            try:
                ser.write((cmd + "\n").encode("utf-8"))
                self._lg_terminal.appendPlainText(f"[TX] {cmd}")
            except Exception as exc:
                self._lg_terminal.appendPlainText(f"[auto-init error] {exc}")
                return
            QTimer.singleShot(1200, lambda: _send_next(idx + 1))
        QTimer.singleShot(800, lambda: _send_next(0))

    def _on_lg_connect_error(self, msg: str) -> None:
        self._lg_status.set_state("error")
        self._btn_lg.setEnabled(True)
        QMessageBox.critical(self, "LG TV Connection Error", msg)

    def _disconnect_lg(self) -> None:
        if self._lg_read_worker:
            self._lg_read_worker.stop()
            self._lg_read_worker = None
        try:
            ser = self._engine.lg_tv_serial
            if ser and ser.is_open:
                ser.write(b"exit\r\n")
                import time; time.sleep(0.2)
                ser.close()
        except Exception:
            pass
        self._engine.lg_tv_serial = None
        self._engine.lg_log_tx = None
        self._lg_terminal.appendPlainText("[TX] exit")
        self._lg_terminal.appendPlainText("── Disconnected ──")
        self._lg_status.set_state("disconnected")
        self._btn_lg.setEnabled(True)
        self._btn_lg_dis.setEnabled(False)

    def _on_lg_data(self, text: str) -> None:
        import json as _json
        clean = text.replace("\r\n", "\n").replace("\r", "\n").strip("\n")
        if clean:
            self._lg_terminal.appendPlainText(f"[RX] {clean}")
            self._lg_terminal.ensureCursorVisible()
        self._lg_json_buf += "\n" + clean
        buf = self._lg_json_buf
        while True:
            start = buf.find("{")
            if start == -1:
                break
            depth = 0
            end = -1
            for i, ch in enumerate(buf[start:], start):
                if ch == "{":
                    depth += 1
                elif ch == "}":
                    depth -= 1
                    if depth == 0:
                        end = i + 1
                        break
            if end == -1:
                break
            chunk = buf[start:end]
            buf = buf[end:]
            try:
                obj = _json.loads(chunk)
            except _json.JSONDecodeError:
                continue
            self._handle_lg_json(obj)
        self._lg_json_buf = buf

    def _handle_lg_json(self, obj: dict) -> None:
        if "modelName" in obj:
            raw_model = str(obj["modelName"])
            model_part = raw_model.split(".")[0]
            self._engine.brand = "LG"
            self._engine.model_name = model_part
            if self._engine.on_lg_device_info:
                self._engine.on_lg_device_info("LG", model_part)
            self._lg_terminal.appendPlainText(f"── Auto-fill: Brand=LG  Model={model_part} ──")
        if "serialNumber" in obj:
            self._engine.lg_serial_number = str(obj["serialNumber"])
            self._lg_terminal.appendPlainText(f"── Serial No.: {obj['serialNumber']} ──")
        if "core_os_release" in obj:
            self._engine.lg_sw_version = str(obj["core_os_release"])
            self._lg_terminal.appendPlainText(f"── SW Version: {obj['core_os_release']} ──")
        if "core_os_release_codename" in obj:
            self._engine.lg_sw_codename = str(obj["core_os_release_codename"])
            self._lg_terminal.appendPlainText(f"── SW Codename: {obj['core_os_release_codename']} ──")

    def _lg_send_cmd(self) -> None:
        ser = self._engine.lg_tv_serial
        if not ser or not ser.is_open:
            self._lg_terminal.appendPlainText("[Error] LG TV is not connected.")
            return
        raw = self._lg_cmd_edit.text()
        if not raw:
            return
        eol_map = {"\\n": b"\n", "\\r\\n": b"\r\n", "\\r": b"\r", "None": b""}
        eol = eol_map.get(self._lg_eol_combo.currentText(), b"\n")
        try:
            ser.write(raw.encode("utf-8") + eol)
            self._lg_terminal.appendPlainText(f"[TX] {raw}")
        except Exception as exc:
            self._lg_terminal.appendPlainText(f"[TX Error] {exc}")
        self._lg_cmd_edit.clear()


# ---------------------------------------------------------------------------
# Setting Panel  (Center Align 1 : LG TV Serial 3)
# ---------------------------------------------------------------------------

class SettingPanel(QWidget):
    """Center Align (좌 1/4) + LG TV Serial (우 3/4) 통합 패널."""

    def __init__(self, engine: MeasurementEngine, parent: QWidget | None = None) -> None:
        super().__init__(parent)
        from PySide6.QtWidgets import QFrame as _QFrame

        self._center = CenterAlignPanel(engine)
        self._lg     = LgTvPanel(engine)

        layout = QHBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)

        layout.addWidget(self._center, stretch=1)

        div = _QFrame()
        div.setFrameShape(_QFrame.Shape.VLine)
        div.setFixedWidth(1)
        div.setStyleSheet("color: #d0d5e8;")
        layout.addWidget(div)

        layout.addWidget(self._lg, stretch=3)

    # ── LG 속성 프록시 (run_demo.py 호환) ─────────────────────────────
    @property
    def _lg_status(self):     return self._lg._lg_status
    @property
    def _btn_lg(self):        return self._lg._btn_lg
    @property
    def _btn_lg_dis(self):    return self._lg._btn_lg_dis
    @property
    def _lg_terminal(self):   return self._lg._lg_terminal

    @property
    def _lg_read_worker(self):        return self._lg._lg_read_worker
    @_lg_read_worker.setter
    def _lg_read_worker(self, v):     self._lg._lg_read_worker = v

    def _on_lg_data(self, text: str) -> None:
        self._lg._on_lg_data(text)


# ---------------------------------------------------------------------------
# Luminance Swing Panel
# ---------------------------------------------------------------------------

_SWING_CASE_COLORS = {"Vivid": "#e74c3c", "Standard": "#4f8ef7", "Cinema": "#27ae60"}


_SWING_CASE_ORDER = ["Vivid", "Standard", "Cinema"]


class LumSwingPanel(QWidget):
    def __init__(self, engine: MeasurementEngine, parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self._engine = engine
        self._worker: Optional[MeasurementWorker] = None
        # "SDR_Vivid", "HDR_Standard" 등 키로 결과 누적
        self._all_data: Dict[str, List[MeasureResult]] = {}
        # 현재 측정 중인 키
        self._current_key: str = ""
        # 범례 순서 고정: 초기화 시 Vivid→Standard→Cinema 순으로 미리 생성
        # {case: QLineSeries}
        self._sdr_series: Dict[str, QLineSeries] = {}
        self._hdr_series: Dict[str, QLineSeries] = {}

        layout = QVBoxLayout(self)

        # ── 컨트롤 한 줄 ────────────────────────────────────────────────
        ctrl = QHBoxLayout()
        title = QLabel("📈 Lum. Swing")
        title.setStyleSheet("font-size:15px;font-weight:bold;")
        ctrl.addWidget(title)
        ctrl.addWidget(QLabel("PSM:"))
        self._case_combo = QComboBox()
        self._case_combo.addItems(["Vivid", "Standard", "Cinema"])
        ctrl.addWidget(self._case_combo)
        self._hdr_check = QCheckBox("HDR")
        self._hdr_check.setChecked(False)
        self._hdr_check.stateChanged.connect(self._on_hdr_toggled)
        ctrl.addWidget(self._hdr_check)
        self._btn_run = QPushButton("▶  Start")
        self._btn_run.setObjectName("primary")
        self._btn_run.clicked.connect(self._run)
        ctrl.addWidget(self._btn_run)
        self._btn_stop = QPushButton("■  Stop")
        self._btn_stop.setObjectName("danger")
        self._btn_stop.setEnabled(False)
        self._btn_stop.clicked.connect(self._stop)
        ctrl.addWidget(self._btn_stop)
        self._btn_clear = QPushButton("🗑 Clear")
        self._btn_clear.clicked.connect(self._clear)
        ctrl.addWidget(self._btn_clear)
        self._btn_export = QPushButton("💾  Save Excel")
        self._btn_export.clicked.connect(self._export)
        ctrl.addWidget(self._btn_export)
        ctrl.addStretch()
        layout.addLayout(ctrl)

        # ── 시간 설정 한 줄 (초 단위, 1초 간격 고정) ─────────────────────
        time_row = QHBoxLayout()
        time_row.setSpacing(6)
        time_row.addWidget(QLabel("Meas. Time:"))
        self._total_sec = QSpinBox()
        self._total_sec.setRange(1, 7200)   # 최대 2시간
        self._total_sec.setValue(301)        # 기본 301초 (≈5분)
        self._total_sec.setSuffix(" s")
        self._total_sec.setFixedWidth(100)
        self._total_sec.valueChanged.connect(self._update_time_hint)
        time_row.addWidget(self._total_sec)
        self._time_hint = QLabel()
        self._time_hint.setObjectName("muted")
        time_row.addWidget(self._time_hint)
        time_row.addStretch()
        layout.addLayout(time_row)
        self._update_time_hint()

        # ── 프로그레스 + 상태 (다음 줄) ─────────────────────────────────
        prog_row = QHBoxLayout()
        self._progress = QProgressBar()
        self._progress.setFixedHeight(10)
        prog_row.addWidget(self._progress)
        self._status_label = QLabel("Idle")
        self._status_label.setObjectName("muted")
        prog_row.addWidget(self._status_label)
        prog_row.addStretch()
        layout.addLayout(prog_row)

        # ── SDR 차트 | HDR 차트 (좌우 분할) ─────────────────────────────
        def _make_chart(title: str):
            chart = QChart()
            chart.setTitle(title)
            chart.setBackgroundBrush(QColor("#ffffff"))

            # ── 그래프 글자 크기 조정 ──────────────────────────────────────
            title_font = chart.titleFont()
            title_font.setPointSize(12)  # 차트 제목 폰트 크기 (pt)
            chart.setTitleFont(title_font)
            chart.setTitleBrush(QColor("#1a1d2e"))

            legend = chart.legend()
            legend.setVisible(True)
            legend.setLabelColor(QColor("#1a1d2e"))
            legend_font = legend.font()
            legend_font.setPointSize(11)  # 범례 폰트 크기 (pt)
            legend.setFont(legend_font)

            # 여백 최소화
            chart.setMargins(QMargins(2, 2, 2, 2))  # 차트 외부 여백 (px): 상하좌우

            # X축
            ax = QValueAxis()
            ax.setTitleText("Sample #")
            ax.setLabelFormat("%d")
            ax.setLabelsBrush(QColor("#6b7080"))
            ax.setTitleBrush(QColor("#6b7080"))
            ax_font = ax.labelsFont()
            ax_font.setPointSize(10)  # X축 눈금 레이블 폰트 (pt)
            ax.setLabelsFont(ax_font)
            ax_title_font = ax.titleFont()
            ax_title_font.setPointSize(10)  # X/Y축 타이틀 폰트 (pt)
            ax.setTitleFont(ax_title_font)

            # Y축
            ay = QValueAxis()
            ay.setTitleText("Lv (cd/m²)")
            ay.setLabelFormat("%d")
            ay.setLabelsBrush(QColor("#6b7080"))
            ay.setTitleBrush(QColor("#6b7080"))
            ay.setLabelsFont(ax_font)   # Y축 눈금 레이블
            ay.setTitleFont(ax_title_font)  # Y축 타이틀

            chart.addAxis(ax, Qt.AlignmentFlag.AlignBottom)
            chart.addAxis(ay, Qt.AlignmentFlag.AlignLeft)
            view = QChartView(chart)
            view.setRenderHint(QPainter.RenderHint.Antialiasing)
            view.setMinimumHeight(320)  # 차트 뷰 최소 높이 (px) — 화면 분할 기준
            view.setStyleSheet("border:1px solid #c8d0e0;border-radius:4px;")
            return chart, ax, ay, view

        (self._chart_sdr, self._ax_x_sdr, self._ax_y_sdr, view_sdr) = _make_chart("SDR")
        (self._chart_hdr, self._ax_x_hdr, self._ax_y_hdr, view_hdr) = _make_chart("HDR")

        # ── 시리즈를 Vivid→Standard→Cinema 순서로 미리 생성 (범례 순서 고정) ──
        def _make_series(case: str, chart: QChart, ax_x: QValueAxis, ax_y: QValueAxis) -> QLineSeries:
            s = QLineSeries()
            s.setName(case)
            pen = s.pen()
            pen.setColor(QColor(_SWING_CASE_COLORS[case]))
            pen.setWidth(2)
            s.setPen(pen)
            chart.addSeries(s)
            s.attachAxis(ax_x)
            s.attachAxis(ax_y)
            return s

        for _case in _SWING_CASE_ORDER:
            self._sdr_series[_case] = _make_series(
                _case, self._chart_sdr, self._ax_x_sdr, self._ax_y_sdr)
            self._hdr_series[_case] = _make_series(
                _case, self._chart_hdr, self._ax_x_hdr, self._ax_y_hdr)

        chart_split = QSplitter(Qt.Orientation.Horizontal)
        chart_split.addWidget(view_sdr)
        chart_split.addWidget(view_hdr)
        chart_split.setStretchFactor(0, 1)
        chart_split.setStretchFactor(1, 1)
        layout.addWidget(chart_split, stretch=3)

        # ── Lv 피벗 테이블: 행=#, 열=모드(SDR_Vivid 등) ──────────────────
        self._table = QTableWidget(0, 1)
        self._table.setAlternatingRowColors(True)
        self._table.setHorizontalHeaderLabels(["#"])
        self._table.setMinimumHeight(120)
        layout.addWidget(self._table, stretch=1)

    # ── 차트 헬퍼 ────────────────────────────────────────────────────────

    def _get_axes(self, is_hdr: bool):
        if is_hdr:
            return self._ax_x_hdr, self._ax_y_hdr
        return self._ax_x_sdr, self._ax_y_sdr

    def _get_series(self, case: str, is_hdr: bool) -> QLineSeries:
        return (self._hdr_series if is_hdr else self._sdr_series)[case]

    def _clear(self) -> None:
        self._all_data.clear()
        self._current_key = ""
        # 시리즈 데이터만 삭제 — 차트에서 제거하지 않아 범례 순서 유지
        for s in list(self._sdr_series.values()) + list(self._hdr_series.values()):
            s.clear()
        # 축 범위 초기화
        for ax in (self._ax_x_sdr, self._ax_x_hdr):
            ax.setRange(0, 10)
        for ay in (self._ax_y_sdr, self._ax_y_hdr):
            ay.setRange(0, 100)
        self._table.setRowCount(0)
        self._table.setColumnCount(1)
        self._table.setHorizontalHeaderLabels(["#"])
        self._status_label.setText("Cleared")

    # ── 시간 힌트 ────────────────────────────────────────────────────────

    def _update_time_hint(self) -> None:
        total_sec = self._total_sec.value()
        m, s = divmod(total_sec, 60)
        self._time_hint.setText(f"→ {total_sec} samples  ({m}m {s}s)")

    # ── 실행 / 중지 ──────────────────────────────────────────────────────

    def _run(self) -> None:
        case = self._case_combo.currentText()
        is_hdr = self._hdr_check.isChecked()
        mode = "HDR" if is_hdr else "SDR"
        self._current_key = f"{mode}_{case}"

        sample_count = self._total_sec.value()   # 1초 간격 고정 → 샘플 수 = 총 초

        # 재측정: 해당 시리즈 데이터만 초기화 (범례 순서 그대로 유지)
        self._get_series(case, is_hdr).clear()
        self._all_data[self._current_key] = []

        self._worker = MeasurementWorker(self._engine, "lum_swing",
                                          case=case, is_hdr=is_hdr,
                                          sample_count=sample_count,
                                          interval_sec=1.0)
        self._worker.progress.connect(self._on_progress)
        self._worker.succeeded.connect(self._on_finished)
        self._worker.error.connect(self._on_error)
        wire_worker_cleanup(self._worker, self, '_worker')
        self._worker.start()
        self._btn_run.setEnabled(False)
        self._btn_stop.setEnabled(True)
        self._status_label.setText(f"{mode} {case} measuring… ({sample_count} samples / {sample_count} s)")

    def _stop(self) -> None:
        self._engine.stop_sequence()
        if self._worker is not None:
            self._worker.requestInterruption()

    @Slot(int)
    def _on_hdr_toggled(self, state: int) -> None:
        gen = self._engine.generator
        if gen is None or not gen.is_connected:
            self._hdr_check.blockSignals(True)
            self._hdr_check.setChecked(not bool(state))
            self._hdr_check.blockSignals(False)
            self._status_label.setText("Generator not connected.")
            return
        if getattr(self, '_hdr_worker', None) is not None:
            return
        enabled = bool(state)
        self._hdr_check.setEnabled(False)
        self._status_label.setText("Switching to HDR…" if enabled else "Switching to SDR…")

        def _done():
            self._hdr_check.setEnabled(True)
            self._status_label.setText("HDR mode" if enabled else "SDR mode")
            setattr(self, '_hdr_worker', None)

        worker = ConnectWorker(lambda: gen.set_hdr(enabled))
        worker.error.connect(lambda msg: (
            QMessageBox.critical(self, "HDR Switch Error", msg),
            self._hdr_check.blockSignals(True),
            self._hdr_check.setChecked(not enabled),
            self._hdr_check.blockSignals(False),
        ))
        wire_worker_cleanup(worker, self, '_hdr_worker', extra_cb=_done)
        worker.start()
        self._hdr_worker = worker

    def _auto_save(self) -> None:
        if not self._engine.auto_save_dir or not self._all_data:
            return
        brand = self._engine.brand or "brand"
        model = self._engine.model_name or "model"
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"lum_swing_{brand}_{model}_{ts}.xlsx"
        path = os.path.join(self._engine.auto_save_dir, filename)
        try:
            ExcelExporter().export_lum_swing(
                self._all_data, self._engine.brand, self._engine.model_name,
                file_path=path,
            )
            self._status_label.setText(f"Done  |  Saved: {path}")
        except Exception as e:
            QMessageBox.warning(self, "Auto-save Failed", str(e))

    def _export(self) -> None:
        if not self._all_data:
            QMessageBox.information(self, "Notice", "No data to save.")
            return
        brand = self._engine.brand or "brand"
        model = self._engine.model_name or "model"
        path, _ = QFileDialog.getSaveFileName(self, "Save Excel",
                                               f"lum_swing_{brand}_{model}.xlsx",
                                               "Excel (*.xlsx)")
        if path:
            ExcelExporter().export_lum_swing(
                self._all_data, self._engine.brand, self._engine.model_name,
                file_path=path,
            )
            QMessageBox.information(self, "Saved", f"Saved: {path}")

    @Slot(str, float, object)
    def _on_progress(self, _step: str, pct: float, data: Any) -> None:
        self._progress.setValue(int(pct * 100))
        if not isinstance(data, MeasureResult):
            return
        key = self._current_key
        self._all_data.setdefault(key, []).append(data)
        rows = self._all_data[key]
        n = len(rows)

        is_hdr = key.startswith("HDR")
        case = key.split("_", 1)[1]
        series = self._get_series(case, is_hdr)
        ax_x, ax_y = self._get_axes(is_hdr)
        series.append(float(n), data.Lv)
        ax_x.setMax(max(ax_x.max(), float(n) + 2))
        ax_y.setMax(max(ax_y.max(), data.Lv * 1.15))

        self._update_table(key, data, n)
        self._status_label.setText(f"{key} #{n}  Lv={data.Lv:.3f} cd/m²")

    @Slot(object)
    def _on_finished(self, _result: Any) -> None:
        self._progress.setValue(100)
        key = self._current_key
        n = len(self._all_data.get(key, []))
        self._status_label.setText(f"{key} done — {n} measurements")
        self._btn_run.setEnabled(True)
        self._btn_stop.setEnabled(False)
        self._engine.session_swing[key] = list(self._all_data.get(key, []))
        path = _save_all_session(self._engine)
        if path:
            self._status_label.setText(f"{key} done — {n} measurements  |  Saved: {path}")

    @Slot(str)
    def _on_error(self, msg: str) -> None:
        QMessageBox.critical(self, "Error", msg)
        self._btn_run.setEnabled(True)
        self._btn_stop.setEnabled(False)

    def _get_col_for_key(self, key: str) -> int:
        for c in range(1, self._table.columnCount()):
            h = self._table.horizontalHeaderItem(c)
            if h and h.text() == key:
                return c
        col = self._table.columnCount()
        self._table.setColumnCount(col + 1)
        self._table.setHorizontalHeaderItem(col, QTableWidgetItem(key))
        return col

    def _update_table(self, key: str, r: MeasureResult, n: int) -> None:
        col = self._get_col_for_key(key)
        # 행이 부족하면 추가 (# 셀 포함)
        while n > self._table.rowCount():
            row_idx = self._table.rowCount()
            self._table.insertRow(row_idx)
            idx_item = QTableWidgetItem(str(row_idx + 1))
            idx_item.setTextAlignment(Qt.AlignmentFlag.AlignCenter)
            self._table.setItem(row_idx, 0, idx_item)
        item = QTableWidgetItem(f"{r.Lv:.3f}")
        item.setTextAlignment(Qt.AlignmentFlag.AlignCenter)
        self._table.setItem(n - 1, col, item)
        self._table.scrollToBottom()

    def begin_auto_stream(self, key: str) -> None:
        """AutoAll step 시작 시 호출 — 해당 key의 시리즈를 초기화하고 live 수신 준비."""
        self._current_key = key
        self._all_data[key] = []
        is_hdr = key.startswith("HDR")
        case   = key.split("_", 1)[1]
        self._get_series(case, is_hdr).clear()
        self._status_label.setText(f"{key} measuring (Auto All)…")

    def refresh_from_engine(self, key: str) -> None:
        """Reload swing chart/table from engine.session_swing[key] (called by AutoAll)."""
        data = self._engine.session_swing.get(key, [])
        if not data:
            return
        is_hdr = key.startswith("HDR")
        case   = key.split("_", 1)[1]
        series = self._get_series(case, is_hdr)
        ax_x, ax_y = self._get_axes(is_hdr)
        series.clear()
        for n, r in enumerate(data, 1):
            series.append(float(n), r.Lv)
        if data:
            ax_x.setMax(max(ax_x.max(), float(len(data)) + 2))
            ax_y.setMax(max(ax_y.max(), max(r.Lv for r in data) * 1.15))
        self._all_data[key] = list(data)
        self._status_label.setText(f"{key} — {len(data)} samples loaded")


# ---------------------------------------------------------------------------
# Luminance Loading Panel
# ---------------------------------------------------------------------------

_BRAND_COLOR_QT = {"samsung": "#0070C0", "lg": "#FF0000", "sony": "#00B050"}


def _qt_brand_color(brand: str) -> str:
    return _BRAND_COLOR_QT.get(brand.lower().strip(), "#FF8800")


_CASE_COLORS = {"Vivid": "#e74c3c", "Standard": "#4f8ef7", "Cinema": "#27ae60"}


class LumLoadingPanel(QWidget):
    def __init__(self, engine: MeasurementEngine, parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self._engine = engine
        self._results: Dict[str, Any] = {}
        self._raw_data: Dict[int, List[MeasureResult]] = {}
        # mode("SDR"/"HDR") → case → apl → results
        self._all_data: Dict[str, Dict[str, Dict[int, List[MeasureResult]]]] = {
            "SDR": {}, "HDR": {}
        }
        layout = QVBoxLayout(self)

        title = QLabel("📊 APL Loading (Luminance Loading)")
        title.setStyleSheet("font-size:15px;font-weight:bold;")
        layout.addWidget(title)

        # ── 설정 한 줄 ────────────────────────────────────────────────
        cfg_row = QHBoxLayout()
        cfg_row.addWidget(QLabel("Version:"))
        self._version_combo = QComboBox()
        self._version_combo.addItems(["10-step", "11-step", "37-step", "2-step"])
        self._version_combo.setFixedWidth(80)
        cfg_row.addWidget(self._version_combo)
        cfg_row.addSpacing(12)
        cfg_row.addWidget(QLabel("Case:"))
        self._case_combo = QComboBox()
        self._case_combo.addItems(["Vivid", "Standard", "Cinema"])
        self._case_combo.setFixedWidth(90)
        cfg_row.addWidget(self._case_combo)
        cfg_row.addSpacing(12)
        cfg_row.addWidget(QLabel("Meas. Count:"))
        self._meas_count = QSpinBox()
        self._meas_count.setRange(1, 10)
        self._meas_count.setValue(1)
        self._meas_count.setSuffix(" ×")
        self._meas_count.setFixedWidth(100)
        cfg_row.addWidget(self._meas_count)
        cfg_row.addSpacing(12)
        self._hdr_check = QCheckBox("HDR")
        self._hdr_check.stateChanged.connect(self._on_hdr_toggled)
        cfg_row.addWidget(self._hdr_check)
        cfg_row.addStretch()
        layout.addLayout(cfg_row)

        # ── 쿨링 조건 한 줄 ──────────────────────────────────────────
        cool_row = QHBoxLayout()
        self._cooling_check = QCheckBox("Cooling")
        cool_row.addWidget(self._cooling_check)
        cool_row.addWidget(QLabel("APL <"))
        self._cool_apl_spin = QSpinBox()
        self._cool_apl_spin.setRange(1, 100)
        self._cool_apl_spin.setValue(10)
        self._cool_apl_spin.setSuffix(" %")
        self._cool_apl_spin.setFixedWidth(100)
        cool_row.addWidget(self._cool_apl_spin)
        cool_row.addWidget(QLabel("at"))
        self._cool_sec_spin = QSpinBox()
        self._cool_sec_spin.setRange(1, 60)
        self._cool_sec_spin.setValue(5)
        self._cool_sec_spin.setSuffix(" s")
        self._cool_sec_spin.setFixedWidth(100)
        cool_row.addWidget(self._cool_sec_spin)
        cool_row.addWidget(QLabel("Black output"))
        cool_row.addStretch()
        layout.addLayout(cool_row)

        # ── 버튼 행 ──────────────────────────────────────────────────
        btn_row = QHBoxLayout()
        self._btn_run = QPushButton("▶  Start")
        self._btn_run.setObjectName("primary")
        self._btn_run.clicked.connect(self._run)
        btn_row.addWidget(self._btn_run)
        self._btn_stop = QPushButton("■  Stop")
        self._btn_stop.setObjectName("danger")
        self._btn_stop.setEnabled(False)
        self._btn_stop.clicked.connect(self._stop)
        btn_row.addWidget(self._btn_stop)
        self._btn_export = QPushButton("💾  Save Excel")
        self._btn_export.clicked.connect(self._export)
        btn_row.addWidget(self._btn_export)
        self._btn_clear = QPushButton("🗑  Clear")
        self._btn_clear.clicked.connect(self._clear_chart)
        btn_row.addWidget(self._btn_clear)
        btn_row.addStretch()
        layout.addLayout(btn_row)

        self._progress = QProgressBar()
        layout.addWidget(self._progress)
        self._status_label = QLabel("Idle")
        self._status_label.setObjectName("muted")
        layout.addWidget(self._status_label)

        # ── SDR / HDR 차트 좌우 분할 ─────────────────────────────────
        def _make_apl_chart(title: str):
            chart = QChart()
            chart.setTitle(title)
            chart.setBackgroundBrush(QColor("#ffffff"))
            chart.setMargins(QMargins(2, 2, 2, 2))
            chart.setContentsMargins(0, 0, 0, 0)

            title_font = chart.titleFont()
            title_font.setPointSize(9)
            chart.setTitleFont(title_font)
            chart.setTitleBrush(QColor("#1a1d2e"))

            legend = chart.legend()
            legend.setLabelColor(QColor("#1a1d2e"))
            legend.show()
            leg_font = legend.font()
            leg_font.setPointSize(8)
            legend.setFont(leg_font)

            ax = QValueAxis()
            ax.setTitleText("APL (%)")
            ax.setRange(0, 100)
            ax.setTickCount(11)
            ax.setLabelsBrush(QColor("#6b7080"))
            ax.setTitleBrush(QColor("#6b7080"))
            ax.setLabelFormat("%d")
            ax_font = ax.labelsFont()
            ax_font.setPointSize(8)
            ax.setLabelsFont(ax_font)
            ax_title_font = ax.titleFont()
            ax_title_font.setPointSize(8)
            ax.setTitleFont(ax_title_font)

            ay = QValueAxis()
            ay.setTitleText("Lv (cd/m²)")
            ay.setLabelsBrush(QColor("#6b7080"))
            ay.setTitleBrush(QColor("#6b7080"))
            ay.setLabelFormat("%d")
            ay.setLabelsFont(ax_font)
            ay.setTitleFont(ax_title_font)

            chart.addAxis(ax, Qt.AlignmentFlag.AlignBottom)
            chart.addAxis(ay, Qt.AlignmentFlag.AlignLeft)
            view = QChartView(chart)
            view.setRenderHint(QPainter.RenderHint.Antialiasing)
            view.setMinimumHeight(60)
            view.setStyleSheet("border:1px solid #c8d0e0;border-radius:4px;")
            return chart, ax, ay, view

        (self._chart_sdr, self._ax_x_sdr,
         self._ax_y_sdr, view_sdr) = _make_apl_chart("SDR")
        (self._chart_hdr, self._ax_x_hdr,
         self._ax_y_hdr, view_hdr) = _make_apl_chart("HDR")

        # ── 시리즈를 Vivid→Standard→Cinema 순서로 미리 생성 (범례 순서 고정) ──
        _CASE_ORDER = ["Vivid", "Standard", "Cinema"]
        self._sdr_apl_series: Dict[str, QLineSeries] = {}
        self._hdr_apl_series: Dict[str, QLineSeries] = {}

        def _make_apl_series(case: str, chart: QChart,
                             ax_x: QValueAxis, ax_y: QValueAxis) -> QLineSeries:
            s = QLineSeries()
            s.setName(case)
            pen = s.pen()
            pen.setColor(QColor(_CASE_COLORS.get(case, "#888899")))
            pen.setWidth(2)
            s.setPen(pen)
            chart.addSeries(s)
            s.attachAxis(ax_x)
            s.attachAxis(ax_y)
            return s

        for _case in _CASE_ORDER:
            self._sdr_apl_series[_case] = _make_apl_series(
                _case, self._chart_sdr, self._ax_x_sdr, self._ax_y_sdr)
            self._hdr_apl_series[_case] = _make_apl_series(
                _case, self._chart_hdr, self._ax_x_hdr, self._ax_y_hdr)

        # SDR / HDR 차트 좌우 분할
        chart_split = QSplitter(Qt.Orientation.Horizontal)
        chart_split.addWidget(view_sdr)
        chart_split.addWidget(view_hdr)
        chart_split.setStretchFactor(0, 1)
        chart_split.setStretchFactor(1, 1)

        self._table = QTableWidget(0, 9)
        self._table.setAlternatingRowColors(True)
        self._table.setHorizontalHeaderLabels(["APL%", "#", "Lv (cd/m²)", "x", "y", "u'", "v'", "CCT (K)", "Duv"])
        self._table.setMinimumHeight(80)

        # 차트-테이블 상하 분할 (드래그로 비율 조정 가능)
        v_split = QSplitter(Qt.Orientation.Vertical)
        v_split.addWidget(chart_split)
        v_split.addWidget(self._table)
        v_split.setSizes([110, 140])   # ← 초기 비율 (픽셀), 여기를 바꾸면 됨
        layout.addWidget(v_split)

    @Slot(int)
    def _on_hdr_toggled(self, state: int) -> None:
        gen = self._engine.generator
        if gen is None or not gen.is_connected:
            self._hdr_check.blockSignals(True)
            self._hdr_check.setChecked(not bool(state))
            self._hdr_check.blockSignals(False)
            self._status_label.setText("Generator not connected.")
            return
        if getattr(self, '_hdr_worker', None) is not None:
            return
        enabled = bool(state)
        self._hdr_check.setEnabled(False)
        self._status_label.setText("Switching to HDR…" if enabled else "Switching to SDR…")

        def _done():
            self._hdr_check.setEnabled(True)
            self._status_label.setText("HDR mode" if enabled else "SDR mode")
            setattr(self, '_hdr_worker', None)

        worker = ConnectWorker(lambda: gen.set_hdr(enabled))
        worker.error.connect(lambda msg: (
            QMessageBox.critical(self, "HDR Switch Error", msg),
            self._hdr_check.blockSignals(True),
            self._hdr_check.setChecked(not enabled),
            self._hdr_check.blockSignals(False),
        ))
        wire_worker_cleanup(worker, self, '_hdr_worker', extra_cb=_done)
        worker.start()
        self._hdr_worker = worker

    def _run(self) -> None:
        version_map = {"37-step": "37", "11-step": "11", "10-step": "10", "2-step": "2"}
        version = version_map[self._version_combo.currentText()]
        self._raw_data.clear()
        self._table.setRowCount(0)
        self._live_series: Optional["QLineSeries"] = None   # 현재 측정 중인 시리즈
        self._btn_run.setEnabled(False)
        self._worker = MeasurementWorker(
            self._engine, "lum_loading",
            version=version,
            case=self._case_combo.currentText(),
            is_hdr=self._hdr_check.isChecked(),
            cooling_enabled=self._cooling_check.isChecked(),
            cooling_apl_threshold=self._cool_apl_spin.value(),
            cooling_duration_sec=float(self._cool_sec_spin.value()),
            measurements_per_step=self._meas_count.value(),
        )
        self._worker.progress.connect(self._on_progress)
        self._worker.succeeded.connect(self._on_finished)
        self._worker.error.connect(lambda m: (
            QMessageBox.critical(self, "Error", m),
            self._btn_run.setEnabled(True),
            self._btn_stop.setEnabled(False),
        ))
        wire_worker_cleanup(self._worker, self, '_worker')
        self._worker.start()
        self._btn_run.setEnabled(False)
        self._btn_stop.setEnabled(True)

    def _stop(self) -> None:
        self._engine.stop_sequence()
        if self._worker is not None:
            self._worker.requestInterruption()

    def _all_data_flat(self) -> Dict[str, Dict[int, List[MeasureResult]]]:
        """_all_data 전체를 {"SDR_Vivid": {...}, "HDR_Standard": {...}, ...} 로 평탄화.

        export_lum_loading 에 넘길 results_by_case 형태.
        """
        flat: Dict[str, Dict[int, List[MeasureResult]]] = {}
        for mode, cases in self._all_data.items():
            for case, apl_dict in cases.items():
                if apl_dict:
                    flat[f"{mode}_{case}"] = apl_dict
        return flat

    def _auto_save(self) -> None:
        flat = self._all_data_flat()
        if not self._engine.auto_save_dir or not flat:
            return
        brand = self._engine.brand or "brand"
        model = self._engine.model_name or "model"
        # 고정 파일명 — 케이스가 추가될 때마다 덮어씌워 항상 최신 전체 데이터 유지
        filename = f"lum_loading_{brand}_{model}.xlsx"
        path = os.path.join(self._engine.auto_save_dir, filename)
        try:
            ExcelExporter().export_lum_loading(
                flat, self._engine.brand, self._engine.model_name,
                file_path=path,
            )
            cases_str = ", ".join(flat.keys())
            self._status_label.setText(
                f"Done — {len(self._raw_data)} APL steps  |  Auto-saved ({cases_str}): {path}"
            )
        except Exception as e:
            QMessageBox.warning(self, "Auto-save Failed", str(e))

    def _on_finished(self, result: Any) -> None:
        self._results = result or {}
        self._live_series = None  # 라이브 시리즈 참조 해제
        mode = "HDR" if self._hdr_check.isChecked() else "SDR"
        case = self._case_combo.currentText()
        self._all_data[mode][case] = dict(self._raw_data)
        # 세션 업데이트 후 통합 파일 자동 저장
        self._engine.session_loading[f"{mode}_{case}"] = dict(self._raw_data)
        self._status_label.setText(f"Done — {len(self._raw_data)} APL steps measured")
        self._btn_run.setEnabled(True)
        self._btn_stop.setEnabled(False)
        path = _save_all_session(self._engine)
        if path:
            self._status_label.setText(
                f"Done — {len(self._raw_data)} APL steps  |  Saved: {path}"
            )

    @Slot(str, float, object)
    def _on_progress(self, _step: str, pct: float, data: Any) -> None:
        self._progress.setValue(int(pct * 100))
        if not (isinstance(data, dict) and "apl" in data):
            return
        apl = int(data["apl"])
        results: List[MeasureResult] = data.get("results", [])
        self._raw_data[apl] = results

        # ── 테이블: 새 APL 스텝 행만 추가 (전체 재구성 없음) ─────────────
        for idx, r in enumerate(results, start=1):
            self._add_table_row(f"{apl}%", str(idx), r)
        self._table.scrollToBottom()

        lv_avg = sum(r.Lv for r in results) / len(results) if results else 0
        self._status_label.setText(f"APL {apl}% — Lv={lv_avg:.3f} cd/m²  ({int(pct*100)}%)")

        # ── 차트: 현재 케이스 시리즈에 포인트 추가만 (전체 재생성 없음) ──
        mode = "HDR" if self._hdr_check.isChecked() else "SDR"
        case = self._case_combo.currentText()
        self._all_data[mode].setdefault(case, {})[apl] = results

        ax_y  = self._ax_y_sdr  if mode == "SDR" else self._ax_y_hdr

        if self._live_series is None:
            # 새 측정 시작 — 미리 만든 시리즈 데이터만 초기화 후 사용
            series_map = self._sdr_apl_series if mode == "SDR" else self._hdr_apl_series
            series = series_map[case]
            series.clear()
            self._live_series = series

        self._live_series.append(float(apl), lv_avg)

        # Y축 범위 동적 확장 — 해당 mode의 모든 케이스 데이터 기준으로 계산
        all_lv = [
            sum(r.Lv for r in rs) / len(rs)
            for case_data in self._all_data[mode].values()
            for rs in case_data.values() if rs
        ]
        if all_lv:
            ax_y.setRange(0, max(all_lv) * 1.15)

    def _refresh_table(self) -> None:
        self._table.setRowCount(0)
        for apl in sorted(self._raw_data):
            results = self._raw_data[apl]
            for idx, r in enumerate(results, start=1):
                self._add_table_row(f"{apl}%", str(idx), r)

    def _refresh_apl_chart(self) -> None:
        for s in list(self._sdr_apl_series.values()) + list(self._hdr_apl_series.values()):
            s.clear()

        for mode, series_map, ax_y in [
            ("SDR", self._sdr_apl_series, self._ax_y_sdr),
            ("HDR", self._hdr_apl_series, self._ax_y_hdr),
        ]:
            all_lv: List[float] = []
            for case, apl_dict in self._all_data[mode].items():
                series = series_map.get(case)
                if not series or not apl_dict:
                    continue
                for apl in sorted(apl_dict):
                    results = apl_dict[apl]
                    if not results:
                        continue
                    lv = sum(r.Lv for r in results) / len(results)
                    series.append(float(apl), lv)
                    all_lv.append(lv)
            if all_lv:
                ax_y.setRange(0, max(all_lv) * 1.15)

    def _clear_chart(self) -> None:
        self._all_data = {"SDR": {}, "HDR": {}}
        self._raw_data.clear()
        self._live_series = None
        for s in list(self._sdr_apl_series.values()) + list(self._hdr_apl_series.values()):
            s.clear()
        for ay in (self._ax_y_sdr, self._ax_y_hdr):
            ay.setRange(0, 100)
        self._table.setRowCount(0)
        self._status_label.setText("Chart cleared")

    def _add_table_row(self, apl_label: str, idx_label: str, r: MeasureResult) -> None:
        row = self._table.rowCount()
        self._table.insertRow(row)
        cct_str = f"{r.cct:.0f}" if r.cct else "—"
        duv_str = f"{r.duv:.5f}" if r.cct else "—"
        for ci, val in enumerate([apl_label, idx_label,
                                   f"{r.Lv:.3f}", f"{r.x:.4f}", f"{r.y:.4f}",
                                   f"{r.u_prime:.4f}", f"{r.v_prime:.4f}",
                                   cct_str, duv_str]):
            item = QTableWidgetItem(val)
            item.setTextAlignment(Qt.AlignmentFlag.AlignCenter)
            self._table.setItem(row, ci, item)

    def _export(self) -> None:
        if not self._raw_data:
            QMessageBox.information(self, "Notice", "No data to save.")
            return
        case = self._case_combo.currentText()
        mode = "HDR" if self._hdr_check.isChecked() else "SDR"
        brand = self._engine.brand or "brand"
        model = self._engine.model_name or "model"
        default_name = f"lum_loading_{mode}_{case.lower()}_{brand}_{model}.xlsx"
        path, _ = QFileDialog.getSaveFileName(self, "Save Excel", default_name, "Excel (*.xlsx)")
        if path:
            try:
                ExcelExporter().export_lum_loading(
                    {case: self._raw_data},
                    self._engine.brand, self._engine.model_name,
                    file_path=path,
                )
                QMessageBox.information(self, "Saved", f"Saved: {path}")
            except Exception:
                import traceback
                QMessageBox.critical(self, "Save Error", traceback.format_exc())

    def refresh_from_engine(self, key: str) -> None:
        """Reload APL loading chart/table from engine.session_loading[key] (called by AutoAll)."""
        apl_dict = self._engine.session_loading.get(key, {})
        if not apl_dict:
            return
        mode, case = ("HDR", key[4:]) if key.startswith("HDR") else ("SDR", key[4:])
        series_map = self._hdr_apl_series if mode == "HDR" else self._sdr_apl_series
        ax_y = self._ax_y_hdr if mode == "HDR" else self._ax_y_sdr
        series = series_map.get(case)
        if series is None:
            return
        series.clear()
        all_lv: list[float] = []
        for apl in sorted(apl_dict):
            results = apl_dict[apl]
            if not results:
                continue
            lv = sum(r.Lv for r in results) / len(results)
            series.append(float(apl), lv)
            all_lv.append(lv)
        if all_lv:
            ax_y.setRange(0, max(all_lv) * 1.15)
        self._all_data[mode][case] = dict(apl_dict)
        self._raw_data.update(apl_dict)
        self._refresh_table()
        n_apl = len(apl_dict)
        self._status_label.setText(f"{key} — {n_apl} APL steps loaded")


# ---------------------------------------------------------------------------
# Gamut Panel
# ---------------------------------------------------------------------------

class GamutPanel(QWidget):
    # u'v' 색도 좌표 색상 매핑
    _UV_COLOR = {
        "red":   QColor("#e74c3c"),
        "green": QColor("#2ecc71"),
        "blue":  QColor("#4f8ef7"),
        "white": QColor("#888899"),
        "black": QColor("#aab0c0"),
    }

    def __init__(self, engine: MeasurementEngine, parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self._engine = engine
        self._results: Dict[str, MeasureResult] = {}
        self._gamut_stats: Dict[str, float] = {}
        layout = QVBoxLayout(self)
        layout.setSpacing(4)

        # ── 상단 고정 컨트롤 (1행: 타이틀 + 버튼 + 통계) ────────────────
        top_row = QHBoxLayout()
        top_row.setSpacing(6)
        title = QLabel("🎨 Gamut")
        title.setStyleSheet("font-size:15px;font-weight:bold;")
        top_row.addWidget(title)
        self._hdr_check = QCheckBox("HDR")
        self._hdr_check.stateChanged.connect(self._on_hdr_toggled)
        top_row.addWidget(self._hdr_check)
        self._btn_run = QPushButton("▶  Start")
        self._btn_run.setObjectName("primary")
        self._btn_run.clicked.connect(self._run)
        top_row.addWidget(self._btn_run)
        self._btn_stop = QPushButton("■  Stop")
        self._btn_stop.setObjectName("danger")
        self._btn_stop.setEnabled(False)
        self._btn_stop.clicked.connect(self._stop)
        top_row.addWidget(self._btn_stop)
        self._btn_export = QPushButton("💾  Save Excel")
        self._btn_export.clicked.connect(self._export)
        top_row.addWidget(self._btn_export)
        top_row.addSpacing(12)
        self._lbl_dci    = QLabel("DCI-P3: —")
        self._lbl_bt2020 = QLabel("BT.2020: —")
        for lbl in (self._lbl_dci, self._lbl_bt2020):
            lbl.setStyleSheet("font-weight:bold; font-size:13px; padding:2px 8px;")
            top_row.addWidget(lbl)
        top_row.addStretch()

        # 프로그레스 + 상태를 같은 줄 오른쪽에 배치
        self._progress = QProgressBar()
        self._progress.setFixedHeight(8)
        self._progress.setMaximumWidth(160)
        top_row.addWidget(self._progress)
        self._status_label = QLabel("Idle")
        self._status_label.setObjectName("muted")
        top_row.addWidget(self._status_label)

        layout.addLayout(top_row)  # stretch=0 → 상단 고정

        # 저장 경로 (한 줄, 평소에는 숨김)
        self._path_label = QLabel()
        self._path_label.setObjectName("muted")
        self._path_label.setStyleSheet("font-size:11px; padding:0px 2px;")
        self._path_label.setVisible(False)
        layout.addWidget(self._path_label)  # stretch=0 → 고정

        # ── 차트 + 테이블 분할 (stretch=1 → 나머지 공간 모두 차지) ─────
        splitter = QSplitter(Qt.Orientation.Horizontal)
        from PySide6.QtWidgets import QFrame as _QFr
        splitter.setFrameShape(_QFr.Shape.NoFrame)

        # u'v' ScatterChart
        self._chart = QChart()
        self._chart.setTitle("u'v' Chromaticity Diagram")
        self._chart.setBackgroundBrush(QColor("#ffffff"))
        self._chart.setTitleBrush(QColor("#1a1d2e"))
        self._chart.setMargins(QMargins(4, 4, 4, 4))
        self._axis_u = QValueAxis()
        self._axis_u.setTitleText("u'")
        self._axis_u.setRange(0.0, 0.65)
        self._axis_u.setLabelsBrush(QColor("#6b7080"))
        self._axis_u.setTitleBrush(QColor("#6b7080"))
        self._axis_v = QValueAxis()
        self._axis_v.setTitleText("v'")
        self._axis_v.setRange(0.0, 0.65)
        self._axis_v.setLabelsBrush(QColor("#6b7080"))
        self._axis_v.setTitleBrush(QColor("#6b7080"))
        self._chart.addAxis(self._axis_u, Qt.AlignmentFlag.AlignBottom)
        self._chart.addAxis(self._axis_v, Qt.AlignmentFlag.AlignLeft)
        self._chart.legend().setLabelColor(QColor("#1a1d2e"))

        self._add_reference_gamut()

        self._gamut_chart_view = QChartView(self._chart)
        self._gamut_chart_view.setRenderHint(QPainter.RenderHint.Antialiasing)
        self._gamut_chart_view.setMinimumWidth(400)
        self._gamut_chart_view.setMinimumHeight(400)
        self._gamut_chart_view.setSizePolicy(
            QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Expanding)
        self._gamut_chart_view.setStyleSheet("background:#ffffff;border:1px solid #d0d5e8;")
        splitter.addWidget(self._gamut_chart_view)

        self._table = QTableWidget(0, 9)
        self._table.setAlternatingRowColors(True)
        self._table.setHorizontalHeaderLabels(["Color", "Lv", "x", "y", "u'", "v'", "X", "Y", "Z"])
        self._table.setMinimumWidth(240)
        self._table.setSizePolicy(QSizePolicy.Policy.Preferred, QSizePolicy.Policy.Expanding)
        splitter.addWidget(self._table)
        splitter.setStretchFactor(0, 4)  # 차트 : 테이블 = 4 : 1
        splitter.setStretchFactor(1, 1)
        splitter.setSizes([600, 200])

        layout.addWidget(splitter, stretch=1)  # ← 핵심: 남은 공간 전부 차지

    def _add_reference_gamut(self) -> None:
        """DCI-P3(점선 흰색), BT.2020(점선 회색) 기준 삼각형 표시."""
        for ref_pts, color, name in [
            (DCI_P3_UV, QColor("#4f8ef7"), "DCI-P3"),
            (BT2020_UV, QColor("#aab0c0"), "BT.2020"),
        ]:
            series = QLineSeries()
            series.setName(name)
            pen = series.pen()
            pen.setColor(color)
            pen.setWidth(1)
            series.setPen(pen)
            for u, v in ref_pts:
                series.append(u, v)
            # 닫힌 삼각형
            series.append(ref_pts[0][0], ref_pts[0][1])
            self._chart.addSeries(series)
            series.attachAxis(self._axis_u)
            series.attachAxis(self._axis_v)

    @Slot(int)
    def _on_hdr_toggled(self, state: int) -> None:
        gen = self._engine.generator
        if gen is None or not gen.is_connected:
            self._hdr_check.blockSignals(True)
            self._hdr_check.setChecked(not bool(state))
            self._hdr_check.blockSignals(False)
            self._status_label.setText("Generator not connected.")
            return
        if getattr(self, '_hdr_worker', None) is not None:
            return
        enabled = bool(state)
        self._hdr_check.setEnabled(False)
        self._status_label.setText("Switching to HDR…" if enabled else "Switching to SDR…")

        def _done():
            self._hdr_check.setEnabled(True)
            self._status_label.setText("HDR mode" if enabled else "SDR mode")
            setattr(self, '_hdr_worker', None)

        worker = ConnectWorker(lambda: gen.set_hdr(enabled))
        worker.error.connect(lambda msg: (
            QMessageBox.critical(self, "HDR Switch Error", msg),
            self._hdr_check.blockSignals(True),
            self._hdr_check.setChecked(not enabled),
            self._hdr_check.blockSignals(False),
        ))
        wire_worker_cleanup(worker, self, '_hdr_worker', extra_cb=_done)
        worker.start()
        self._hdr_worker = worker

    def _run(self) -> None:
        self._table.setRowCount(0)
        self._results.clear()
        self._set_path("")
        # 기존 측정점 시리즈 제거 후 기준 gamut 재설정
        self._chart.removeAllSeries()
        self._add_reference_gamut()
        self._lbl_dci.setText("DCI-P3: —")
        self._lbl_bt2020.setText("BT.2020: —")
        self._btn_run.setEnabled(False)
        self._btn_stop.setEnabled(True)
        self._worker = MeasurementWorker(self._engine, "gamut", is_hdr=self._hdr_check.isChecked())
        self._worker.progress.connect(self._on_progress)
        self._worker.succeeded.connect(self._on_finished)
        self._worker.error.connect(lambda m: (
            QMessageBox.critical(self, "Error", m),
            self._btn_run.setEnabled(True),
            self._btn_stop.setEnabled(False),
        ))
        wire_worker_cleanup(self._worker, self, '_worker')
        self._worker.start()

    def _stop(self) -> None:
        self._engine.stop_sequence()
        if self._worker is not None:
            self._worker.requestInterruption()

    def _set_path(self, path: str) -> None:
        if path:
            self._path_label.setText(f"Saved: {path}")
            self._path_label.setVisible(True)
        else:
            self._path_label.setVisible(False)

    def _auto_save(self) -> None:
        if not self._engine.auto_save_dir or not self._results:
            return
        mode = "HDR" if self._hdr_check.isChecked() else "SDR"
        brand = self._engine.brand or "brand"
        model = self._engine.model_name or "model"
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"gamut_{mode}_{brand}_{model}_{ts}.xlsx"
        path = os.path.join(self._engine.auto_save_dir, filename)
        try:
            ExcelExporter().export_gamut(self._results, self._engine.brand,
                                         self._engine.model_name, file_path=path,
                                         gamut_stats=self._gamut_stats or None)
            self._status_label.setText("Done  |  Auto-saved")
            self._set_path(path)
        except Exception as e:
            QMessageBox.warning(self, "Auto-save Failed", str(e))

    @Slot(object)
    def _on_finished(self, result: Any) -> None:
        self._results = result or {}
        self._status_label.setText("Done")
        self._btn_run.setEnabled(True)
        self._btn_stop.setEnabled(False)
        self._update_gamut_stats()
        self._draw_meas_triangle()
        # 세션 업데이트 후 통합 파일 자동 저장
        mode = "HDR" if self._hdr_check.isChecked() else "SDR"
        self._engine.session_gamut[mode] = dict(self._results)
        path = _save_all_session(self._engine)
        self._set_path(path)

    def _update_gamut_stats(self) -> None:
        r_r = self._results.get("red")
        r_g = self._results.get("green")
        r_b = self._results.get("blue")
        if r_r and r_g and r_b:
            self._gamut_stats = calc_gamut_stats(
                (r_r.u_prime, r_r.v_prime),
                (r_g.u_prime, r_g.v_prime),
                (r_b.u_prime, r_b.v_prime),
            )
            self._lbl_dci.setText(f"DCI-P3: {self._gamut_stats['dci_overlap']:.1f}%")
            self._lbl_bt2020.setText(f"BT.2020: {self._gamut_stats['bt2020_overlap']:.1f}%")

    def _draw_meas_triangle(self) -> None:
        r_r = self._results.get("red")
        r_g = self._results.get("green")
        r_b = self._results.get("blue")
        if not (r_r and r_g and r_b):
            return
        pts = [
            (r_r.u_prime, r_r.v_prime),
            (r_g.u_prime, r_g.v_prime),
            (r_b.u_prime, r_b.v_prime),
        ]
        tri = QLineSeries()
        tri.setName("Measured")
        pen = tri.pen()
        pen.setColor(QColor("#f7c94f"))
        pen.setWidth(2)
        tri.setPen(pen)
        for u, v in pts:
            tri.append(u, v)
        tri.append(pts[0][0], pts[0][1])
        self._chart.addSeries(tri)
        tri.attachAxis(self._axis_u)
        tri.attachAxis(self._axis_v)

    @Slot(str, float, object)
    def _on_progress(self, _step: str, pct: float, data: Any) -> None:
        self._progress.setValue(int(pct * 100))
        if isinstance(data, dict) and "color" in data:
            color = data["color"]
            if color is None:
                self._status_label.setText("Clearing…")
                return
            r = data.get("result")
            if r:
                # 테이블 행 추가
                row = self._table.rowCount()
                self._table.insertRow(row)
                for ci, val in enumerate([color, f"{r.Lv:.3f}", f"{r.x:.4f}", f"{r.y:.4f}",
                                           f"{r.u_prime:.4f}", f"{r.v_prime:.4f}",
                                           f"{r.X:.3f}", f"{r.Y:.3f}", f"{r.Z:.3f}"]):
                    item = QTableWidgetItem(str(val))
                    item.setTextAlignment(Qt.AlignmentFlag.AlignCenter)
                    self._table.setItem(row, ci, item)
                # 꼭짓점 점 표시 제거 — 삼각형 라인만 사용
            self._status_label.setText(f"{color} measured — {int(pct*100)}%")

    def _export(self) -> None:
        if not self._results:
            QMessageBox.information(self, "Notice", "No data to save.")
            return
        brand = self._engine.brand or "brand"
        model = self._engine.model_name or "model"
        default_name = f"gamut_{brand}_{model}.xlsx"
        path, _ = QFileDialog.getSaveFileName(self, "Save Excel", default_name, "Excel (*.xlsx)")
        if path:
            try:
                ExcelExporter().export_gamut(self._results, self._engine.brand,
                                             self._engine.model_name, file_path=path,
                                             gamut_stats=self._gamut_stats or None)
                QMessageBox.information(self, "Saved", f"Saved: {path}")
            except Exception as e:
                QMessageBox.critical(self, "Save Error", str(e))

    def refresh_from_engine(self, mode: str) -> None:
        """Reload gamut chart/table from engine.session_gamut[mode] (called by AutoAll)."""
        results = self._engine.session_gamut.get(mode, {})
        if not results:
            return
        self._results = dict(results)
        # Rebuild table
        self._table.setRowCount(0)
        for color, r in results.items():
            row = self._table.rowCount()
            self._table.insertRow(row)
            for ci, val in enumerate([color, f"{r.Lv:.3f}", f"{r.x:.4f}", f"{r.y:.4f}",
                                       f"{r.u_prime:.4f}", f"{r.v_prime:.4f}",
                                       f"{r.X:.3f}", f"{r.Y:.3f}", f"{r.Z:.3f}"]):
                item = QTableWidgetItem(str(val))
                item.setTextAlignment(Qt.AlignmentFlag.AlignCenter)
                self._table.setItem(row, ci, item)
        # Rebuild chart triangle
        for s in self._chart.series():
            if s.name() == "Measured":
                self._chart.removeSeries(s)
        self._update_gamut_stats()
        self._draw_meas_triangle()
        self._status_label.setText(f"{mode} gamut loaded")


# ---------------------------------------------------------------------------
# Contrast Panel
# ---------------------------------------------------------------------------

class ContrastPanel(QWidget):
    def __init__(self, engine: MeasurementEngine, parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self._engine = engine
        self._results: Dict[float, MeasureResult] = {}
        layout = QVBoxLayout(self)

        title = QLabel("⬛ Contrast Ratio")
        title.setStyleSheet("font-size:15px;font-weight:bold;")
        layout.addWidget(title)
        desc = QLabel("White Raster + Black Window — measures in order: H/V 100% → 50% → 20% → 14.1%.")
        desc.setObjectName("muted")
        layout.addWidget(desc)

        btn_row = QHBoxLayout()
        self._hdr_check = QCheckBox("HDR")
        self._hdr_check.stateChanged.connect(self._on_hdr_toggled)
        btn_row.addWidget(self._hdr_check)
        self._btn_run = QPushButton("▶  Start")
        self._btn_run.setObjectName("primary")
        self._btn_run.clicked.connect(self._run)
        btn_row.addWidget(self._btn_run)
        self._btn_stop = QPushButton("■  Stop")
        self._btn_stop.setObjectName("danger")
        self._btn_stop.setEnabled(False)
        self._btn_stop.clicked.connect(self._stop)
        btn_row.addWidget(self._btn_stop)
        self._btn_export = QPushButton("💾  Save Excel")
        self._btn_export.clicked.connect(self._export)
        btn_row.addWidget(self._btn_export)
        btn_row.addStretch()
        layout.addLayout(btn_row)

        self._progress = QProgressBar()
        layout.addWidget(self._progress)
        self._status_label = QLabel("Idle")
        self._status_label.setObjectName("muted")
        layout.addWidget(self._status_label)

        self._lv_ref: Optional[float] = None  # 창 100% 기준 흑휘도
        self._table = QTableWidget(0, 5)
        self._table.setAlternatingRowColors(True)
        self._table.setHorizontalHeaderLabels(["Black H/V (%)", "Lv (cd/m²)", "x", "y", "CR (White/Lv)"])
        layout.addWidget(self._table)

    @Slot(int)
    def _on_hdr_toggled(self, state: int) -> None:
        gen = self._engine.generator
        if gen is None or not gen.is_connected:
            self._hdr_check.blockSignals(True)
            self._hdr_check.setChecked(not bool(state))
            self._hdr_check.blockSignals(False)
            self._status_label.setText("Generator not connected.")
            return
        if getattr(self, '_hdr_worker', None) is not None:
            return
        enabled = bool(state)
        self._hdr_check.setEnabled(False)
        self._status_label.setText("Switching to HDR…" if enabled else "Switching to SDR…")

        def _done():
            self._hdr_check.setEnabled(True)
            self._status_label.setText("HDR mode" if enabled else "SDR mode")
            setattr(self, '_hdr_worker', None)

        worker = ConnectWorker(lambda: gen.set_hdr(enabled))
        worker.error.connect(lambda msg: (
            QMessageBox.critical(self, "HDR Switch Error", msg),
            self._hdr_check.blockSignals(True),
            self._hdr_check.setChecked(not enabled),
            self._hdr_check.blockSignals(False),
        ))
        wire_worker_cleanup(worker, self, '_hdr_worker', extra_cb=_done)
        worker.start()
        self._hdr_worker = worker

    def _auto_save(self) -> None:
        if not self._engine.auto_save_dir or not self._results:
            return
        mode = "HDR" if self._hdr_check.isChecked() else "SDR"
        brand = self._engine.brand or "brand"
        model = self._engine.model_name or "model"
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"contrast_{mode}_{brand}_{model}_{ts}.xlsx"
        path = os.path.join(self._engine.auto_save_dir, filename)
        try:
            ExcelExporter().export_contrast(self._results, self._engine.brand,
                                            self._engine.model_name, file_path=path)
            self._status_label.setText(f"Done  |  Auto-saved: {path}")
        except Exception as e:
            QMessageBox.warning(self, "Auto-save Failed", str(e))

    def _on_finished(self, result: Any) -> None:
        self._results = result or {}
        self._status_label.setText("Done")
        self._btn_run.setEnabled(True)
        self._btn_stop.setEnabled(False)
        # 세션 업데이트 후 통합 파일 자동 저장
        mode = "HDR" if self._hdr_check.isChecked() else "SDR"
        self._engine.session_contrast[mode] = dict(self._results)
        path = _save_all_session(self._engine)
        if path:
            self._status_label.setText(f"Done  |  Saved: {path}")

    def _run(self) -> None:
        self._table.setRowCount(0)
        self._lv_ref = None
        self._btn_run.setEnabled(False)
        self._btn_stop.setEnabled(True)
        self._worker = MeasurementWorker(self._engine, "contrast", is_hdr=self._hdr_check.isChecked())
        self._worker.progress.connect(self._on_progress)
        self._worker.succeeded.connect(self._on_finished)
        self._worker.error.connect(lambda m: (
            QMessageBox.critical(self, "Error", m),
            self._btn_run.setEnabled(True),
            self._btn_stop.setEnabled(False),
        ))
        wire_worker_cleanup(self._worker, self, '_worker')
        self._worker.start()

    def _stop(self) -> None:
        self._engine.stop_sequence()
        if self._worker is not None:
            self._worker.requestInterruption()

    @Slot(str, float, object)
    def _on_progress(self, _step: str, pct: float, data: Any) -> None:
        self._progress.setValue(int(pct * 100))
        if not (isinstance(data, dict) and "win_size" in data):
            return
        win_size = data["win_size"]   # 0.0 = full white ref, else black window H/V side %
        r = data.get("result")
        if not r:
            return

        # win_size 0.0 = Full White → 기준 Lv 저장
        if win_size == 0.0:
            self._lv_ref = r.Lv

        hv_label = "Full White" if win_size == 0.0 else f"{win_size:.1f} × {win_size:.1f}"

        # CR = Full White Lv / Black Lv (black window 행에만, Lv > 0 가드)
        if win_size > 0.0 and self._lv_ref and self._lv_ref > 0 and r.Lv > 0:
            cr_str = f"{self._lv_ref / r.Lv:.1f} : 1"
        else:
            cr_str = "—"

        row = self._table.rowCount()
        self._table.insertRow(row)
        for ci, val in enumerate([hv_label, f"{r.Lv:.4f}",
                                   f"{r.x:.4f}", f"{r.y:.4f}", cr_str]):
            item = QTableWidgetItem(val)
            item.setTextAlignment(Qt.AlignmentFlag.AlignCenter)
            self._table.setItem(row, ci, item)

        step_label = "Full White" if win_size == 0.0 else f"Win {win_size:.1f}%"
        self._status_label.setText(f"{step_label} measuring — {int(pct*100)}%")

    def _export(self) -> None:
        if not self._results:
            QMessageBox.information(self, "Notice", "No data to save.")
            return
        brand = self._engine.brand or "brand"
        model = self._engine.model_name or "model"
        default_name = f"contrast_{brand}_{model}.xlsx"
        path, _ = QFileDialog.getSaveFileName(self, "Save Excel", default_name, "Excel (*.xlsx)")
        if path:
            ExcelExporter().export_contrast(self._results, self._engine.brand, self._engine.model_name,
                                            file_path=path)
            QMessageBox.information(self, "Saved", f"Saved:\n{path}")

    def refresh_from_engine(self, mode: str) -> None:
        """Reload contrast table from engine.session_contrast[mode] (called by AutoAll)."""
        results = self._engine.session_contrast.get(mode, {})
        if not results:
            return
        self._results = dict(results)
        self._table.setRowCount(0)
        lv_ref: float | None = None
        for win_size in sorted(results):
            r = results[win_size]
            if win_size == 0.0:
                lv_ref = r.Lv
            hv_label = "Full White" if win_size == 0.0 else f"{win_size:.1f} × {win_size:.1f}"
            cr_str = (f"{lv_ref / r.Lv:.1f} : 1"
                      if win_size > 0.0 and lv_ref and lv_ref > 0 and r.Lv > 0 else "—")
            row = self._table.rowCount()
            self._table.insertRow(row)
            for ci, val in enumerate([hv_label, f"{r.Lv:.4f}",
                                       f"{r.x:.4f}", f"{r.y:.4f}", cr_str]):
                item = QTableWidgetItem(val)
                item.setTextAlignment(Qt.AlignmentFlag.AlignCenter)
                self._table.setItem(row, ci, item)
        self._status_label.setText(f"{mode} contrast loaded — {len(results)} steps")


# ---------------------------------------------------------------------------
# Report Panel
# ---------------------------------------------------------------------------

_DEFAULT_MODEL_COLORS = [
    "#4f8ef7", "#e74c3c", "#27ae60", "#f39c12",
    "#9b59b6", "#1abc9c", "#e67e22", "#34495e",
]


class ReportPanel(QWidget):
    # ── 경쟁사 비교 장표 ──────────────────────────────────────────────────────
    _ROW_LABELS = [
        ("White Luminance [nit]", "HDR 10%"),
        ("White Luminance [nit]", "HDR 100%"),
        ("White Luminance [nit]", "SDR 10%"),
        ("White Luminance [nit]", "SDR 100%"),
        ("White Luminance [nit]", "Contrast Ratio"),
        ("White Luminance [nit]", "Black"),
        ("Color Gamut[%]", "DCI-P3 (%)"),
        ("Color Gamut[%]", "BT.2020 (%)"),
    ]
    _COMP_KEYS = [
        "hdr_10", "hdr_100", "sdr_10", "sdr_100",
        "contrast_ratio", "black_lv", "dci_overlap", "bt2020_overlap",
    ]

    # 정수(cd/m²)로 표시할 키 집합 — black_lv 는 소수 유지
    _LV_INT_KEYS: set = {
        "hdr_10", "hdr_100", "sdr_10", "sdr_100",
        "sdr_vivid_10",    "sdr_vivid_100",
        "sdr_standard_10", "sdr_standard_100",
        "hdr_vivid_10",    "hdr_vivid_100",
        "hdr_standard_10", "hdr_standard_100",
        "hdr_cinema_10",   "hdr_cinema_100",
    }

    # ── 광학 측정 데이터 ──────────────────────────────────────────────────────
    _OPTICAL_ROW_LABELS = [
        ("Lum.", "Vivid SDR 10% / 100%"),
        ("Lum.", "Standard SDR 10% / 100%"),
        ("Lum.", "Vivid HDR 10% / 100%"),
        ("Lum.", "Standard HDR 10% / 100%"),
        ("Lum.", "Cinema HDR 10% / 100%"),
        ("Contrast", "Black (Ratio)"),
        ("Color Gamut", "DCI-P3 (%)"),
        ("Color Gamut", "BT.2020 (%)"),
    ]
    # (key_10, key_100) — key_100=None for single-value rows
    _OPTICAL_KEYS: List[tuple] = [
        ("sdr_vivid_10",    "sdr_vivid_100"),
        ("sdr_standard_10", "sdr_standard_100"),
        ("hdr_vivid_10",    "hdr_vivid_100"),
        ("hdr_standard_10", "hdr_standard_100"),
        ("hdr_cinema_10",   "hdr_cinema_100"),
        ("contrast_ratio",  None),
        ("dci_overlap",     None),
        ("bt2020_overlap",  None),
    ]

    def __init__(
        self,
        engine: MeasurementEngine,
        gamut_panel: "GamutPanel",
        lum_panel: "LumLoadingPanel",
        contrast_panel: "ContrastPanel",
        parent: QWidget | None = None,
    ) -> None:
        super().__init__(parent)
        self._engine = engine
        self._gamut_panel = gamut_panel
        self._lum_panel = lum_panel
        self._contrast_panel = contrast_panel
        self._models: List[Dict] = []
        self._model_colors: Dict[str, str] = {}

        layout = QVBoxLayout(self)

        title = QLabel("📋 Report Template")
        title.setStyleSheet("font-size:15px;font-weight:bold;")
        layout.addWidget(title)

        # 1행: 파일·형식·집계
        top_row = QHBoxLayout()
        top_row.setSpacing(6)
        self._btn_load = QPushButton("📂 Load Files")
        self._btn_load.setObjectName("primary")
        self._btn_load.setMinimumWidth(120)  # 파일 불러오기 버튼 최소 너비 (px)
        self._btn_load.clicked.connect(self._load_files)
        top_row.addWidget(self._btn_load)

        top_row.addWidget(QLabel("Report Format:"))
        self._format_combo = QComboBox()
        self._format_combo.addItems(["Competitor Comparison", "Optical Measurement Data"])
        self._format_combo.setMinimumWidth(140)  # 보고서 형식 콤보박스 최소 너비 (px)
        self._format_combo.currentIndexChanged.connect(self._on_format_changed)
        top_row.addWidget(self._format_combo)

        top_row.addWidget(QLabel("Luminance Agg.:"))
        self._agg_combo = QComboBox()
        self._agg_combo.addItems(["Max", "Median", "Min"])
        self._agg_combo.setMinimumWidth(80)
        top_row.addWidget(self._agg_combo)
        top_row.addStretch()
        layout.addLayout(top_row)

        # 2행: 모델 조작·내보내기
        action_row = QHBoxLayout()
        action_row.setSpacing(6)
        self._btn_color = QPushButton("🎨 Change Color")
        self._btn_color.setMinimumWidth(90)
        self._btn_color.clicked.connect(self._change_model_color)
        action_row.addWidget(self._btn_color)
        self._btn_del = QPushButton("✕ Delete Selected")
        self._btn_del.setMinimumWidth(90)
        self._btn_del.clicked.connect(self._delete_selected)
        action_row.addWidget(self._btn_del)
        self._btn_excel = QPushButton("💾 Save Excel")
        self._btn_excel.setMinimumWidth(90)
        self._btn_excel.clicked.connect(self._export_excel)
        action_row.addWidget(self._btn_excel)
        self._btn_copy = QPushButton("📋 Copy to Clipboard")
        self._btn_copy.setMinimumWidth(110)
        self._btn_copy.clicked.connect(self._copy_clipboard)
        action_row.addWidget(self._btn_copy)
        self._btn_ppt = QPushButton("📊 Save PPT")
        self._btn_ppt.setObjectName("success")
        self._btn_ppt.setMinimumWidth(100)
        self._btn_ppt.clicked.connect(self._export_ppt)
        action_row.addWidget(self._btn_ppt)
        action_row.addStretch()
        layout.addLayout(action_row)

        self._model_list = QListWidget()
        self._model_list.setFixedHeight(72)  # 모델 목록 리스트 고정 높이 (px) — 최대 3행 표시
        layout.addWidget(self._model_list)

        self._report_table = QTableWidget(len(self._ROW_LABELS), 2)
        self._report_table.setHorizontalHeaderLabels(["Category", "Item"])
        self._report_table.setMinimumHeight(50)
        self._report_table.verticalHeader().setDefaultSectionSize(26)  # 보고서 테이블 행 높이 (px)
        self._report_table.setColumnWidth(0, 80)   # 구분 열 너비 (px)
        self._report_table.setColumnWidth(1, 80)   # 항목 열 너비 (px)

        # ── APL 차트 생성 헬퍼 ────────────────────────────────────────
        def _make_apl_chart(title: str):
            chart = QChart()
            chart.setTitle(title)
            chart.setBackgroundBrush(QColor("#ffffff"))
            chart.setMargins(QMargins(8, 8, 8, 8))  # 차트 외부 여백 (px) — matplotlib 스타일
            chart.setDropShadowEnabled(False)
            chart.setBackgroundRoundness(0)

            tf = chart.titleFont(); tf.setPointSize(11); tf.setBold(True)  # 차트 제목: 11pt bold
            chart.setTitleFont(tf)
            chart.setTitleBrush(QColor("#222222"))

            _grid_pen = QPen(QColor("#dddddd"))  # 그리드 선 색상 — 연한 회색
            _grid_pen.setStyle(Qt.PenStyle.DashLine)
            _grid_pen.setWidthF(0.8)  # 그리드 선 두께 (px)
            _label_color = QColor("#333333")
            sf = chart.titleFont(); sf.setPointSize(10); sf.setBold(False)  # 축 레이블·타이틀 폰트 크기 (pt)

            ax = QValueAxis()
            ax.setTitleText("APL (%)"); ax.setRange(0, 100); ax.setTickCount(6)  # X축 눈금 개수 (0,20,40,60,80,100%)
            ax.setLabelFormat("%d")
            ax.setLabelsBrush(_label_color); ax.setTitleBrush(_label_color)
            ax.setLabelsFont(sf); ax.setTitleFont(sf)
            ax.setGridLinePen(_grid_pen)
            ax.setLinePen(QPen(QColor("#aaaaaa")))

            ay = QValueAxis()
            ay.setTitleText("Lv (cd/m²)"); ay.setLabelFormat("%d")
            ay.setLabelsBrush(_label_color); ay.setTitleBrush(_label_color)
            ay.setLabelsFont(sf); ay.setTitleFont(sf)
            ay.setGridLinePen(_grid_pen)
            ay.setLinePen(QPen(QColor("#aaaaaa")))

            chart.addAxis(ax, Qt.AlignmentFlag.AlignBottom)
            chart.addAxis(ay, Qt.AlignmentFlag.AlignLeft)

            legend = chart.legend()
            lf = legend.font(); lf.setPointSize(10); legend.setFont(lf)  # 범례 폰트 크기 (pt)
            legend.setLabelColor(QColor("#222222"))
            legend.setAlignment(Qt.AlignmentFlag.AlignBottom)
            legend.setBackgroundVisible(True)
            legend.setBrush(QColor("#ffffff"))
            legend.setPen(QPen(QColor("#dddddd")))

            view = QChartView(chart)
            view.setRenderHint(QPainter.RenderHint.Antialiasing)
            view.setMinimumHeight(100)  # APL 차트 최소 높이 (px)
            view.setStyleSheet("border:1px solid #c8d0e0;border-radius:4px;")
            return chart, ax, ay, view

        # ── SDR Vivid APL 차트 (위) ────────────────────────────────────
        (self._apl_chart_sdr,
         self._apl_axis_x_sdr,
         self._apl_axis_y_sdr,
         self._apl_view_sdr) = _make_apl_chart("SDR Vivid")
        apl_view_sdr = self._apl_view_sdr

        # ── HDR Vivid APL 차트 (아래) ──────────────────────────────────
        (self._apl_chart_hdr,
         self._apl_axis_x_hdr,
         self._apl_axis_y_hdr,
         self._apl_view_hdr) = _make_apl_chart("HDR Vivid")
        apl_view_hdr = self._apl_view_hdr

        # 좌측 휘도: SDR / HDR 탭으로 나누기
        apl_tabs = QTabWidget()
        apl_tabs.setStyleSheet(
            "QTabBar::tab { padding: 4px 14px; font-size: 11px; font-weight: bold; }"
            "QTabBar::tab:selected { color: #1460c0; border-bottom: 2px solid #1878d0; }"
        )
        apl_tabs.addTab(apl_view_sdr, "SDR")
        apl_tabs.addTab(apl_view_hdr, "HDR")

        # ── Gamut u'v' 차트 (단일) ─────────────────────────────────────
        def _make_gamut_chart(title: str):
            chart = QChart()
            chart.setTitle(title)
            chart.setBackgroundBrush(QColor("#ffffff"))
            chart.setMargins(QMargins(8, 8, 8, 8))
            chart.setDropShadowEnabled(False)
            chart.setBackgroundRoundness(0)
            tf = chart.titleFont(); tf.setPointSize(11); tf.setBold(True)
            chart.setTitleFont(tf); chart.setTitleBrush(QColor("#222222"))

            _gpen = QPen(QColor("#dddddd")); _gpen.setStyle(Qt.PenStyle.DashLine); _gpen.setWidthF(0.8)
            _lc = QColor("#333333")
            _sf = chart.titleFont(); _sf.setPointSize(10); _sf.setBold(False)

            gl = chart.legend()
            gl.setLabelColor(QColor("#222222"))
            glf = gl.font(); glf.setPointSize(10); gl.setFont(glf)
            gl.detachFromChart()
            gl.setBackgroundVisible(True)
            gl.setBrush(QColor("#ffffff"))
            gl.setPen(QPen(QColor("#dddddd")))

            ax_u = QValueAxis()
            ax_u.setTitleText("u'"); ax_u.setRange(0.0, 0.65)
            ax_u.setLabelsBrush(_lc); ax_u.setTitleBrush(_lc)
            ax_u.setLabelsFont(_sf); ax_u.setTitleFont(_sf)
            ax_u.setGridLinePen(_gpen); ax_u.setLinePen(QPen(QColor("#aaaaaa")))

            ax_v = QValueAxis()
            ax_v.setTitleText("v'"); ax_v.setRange(0.0, 0.65)
            ax_v.setLabelsBrush(_lc); ax_v.setTitleBrush(_lc)
            ax_v.setLabelsFont(_sf); ax_v.setTitleFont(_sf)
            ax_v.setGridLinePen(_gpen); ax_v.setLinePen(QPen(QColor("#aaaaaa")))

            chart.addAxis(ax_u, Qt.AlignmentFlag.AlignBottom)
            chart.addAxis(ax_v, Qt.AlignmentFlag.AlignLeft)

            view = QChartView(chart)
            view.setRenderHint(QPainter.RenderHint.Antialiasing)
            view.setMinimumHeight(200); view.setMinimumWidth(200)
            view.setStyleSheet("background:#ffffff;border:1px solid #c8d0e0;border-radius:4px;")
            return chart, ax_u, ax_v, view

        (self._gamut_chart,
         self._gamut_axis_u,
         self._gamut_axis_v,
         self._gamut_chart_view) = _make_gamut_chart("u'v' Chromaticity")

        # 기준 삼각형(DCI-P3, BT.2020) 추가
        self._add_ref_gamuts_to(self._gamut_chart, self._gamut_axis_u, self._gamut_axis_v)

        # plotAreaChanged 콜백
        self._gamut_chart.plotAreaChanged.connect(
            lambda _: self._position_gamut_legend_for(self._gamut_chart)
        )
        QTimer.singleShot(0, lambda: self._position_gamut_legend_for(self._gamut_chart))

        # ── 메인 수평 스플리터: 좌=표, 우=3탭(SDR Vivid / HDR Vivid / Gamut) ──
        self._report_table.setMinimumHeight(45)

        right_tabs = QTabWidget()
        right_tabs.setStyleSheet(
            "QTabBar::tab { padding: 5px 16px; font-size: 11px; font-weight: bold; }"
            "QTabBar::tab:selected { color: #1460c0; border-bottom: 2px solid #1878d0; }"
        )
        right_tabs.addTab(apl_view_sdr, "SDR Vivid")
        right_tabs.addTab(apl_view_hdr, "HDR Vivid")
        right_tabs.addTab(self._gamut_chart_view, "Gamut")

        main_splitter = QSplitter(Qt.Orientation.Horizontal)
        main_splitter.addWidget(self._report_table)
        main_splitter.addWidget(right_tabs)
        main_splitter.setStretchFactor(0, 5)
        main_splitter.setStretchFactor(1, 7)
        layout.addWidget(main_splitter, 1)

        self._refresh_report_table()

    def _change_model_color(self) -> None:
        row = self._model_list.currentRow()
        if row < 0 or row >= len(self._models):
            QMessageBox.information(self, "Notice", "Please select a model to change color.")
            return
        entry = self._models[row]
        key = f"{entry['brand']}_{entry['model']}"
        current = QColor(self._model_colors.get(key, _DEFAULT_MODEL_COLORS[0]))
        color = QColorDialog.getColor(current, self, "Select Model Color")
        if color.isValid():
            self._model_colors[key] = color.name()
            self._refresh_apl_chart()
            self._refresh_gamut_chart()

    def _delete_selected(self) -> None:
        row = self._model_list.currentRow()
        if row < 0 or row >= len(self._models):
            return
        self._models.pop(row)
        self._model_list.takeItem(row)
        self._refresh_report_table()

    def _is_optical_format(self) -> bool:
        return self._format_combo.currentIndex() == 1

    def _on_format_changed(self) -> None:
        self._refresh_report_table()

    def _current_row_labels(self):
        return self._OPTICAL_ROW_LABELS if self._is_optical_format() else self._ROW_LABELS

    @staticmethod
    def _fmt(key: str, val) -> str:
        """키 종류에 따라 숫자 포맷 결정."""
        if val is None:
            return "—"
        if key in ReportPanel._LV_INT_KEYS:
            try:
                return str(round(float(val)))
            except (TypeError, ValueError):
                return str(val)
        try:
            return f"{float(val):.1f}"
        except (TypeError, ValueError):
            return str(val)

    def _cell_value(self, entry: Dict, ri: int) -> str:
        if self._is_optical_format():
            k10, k100 = self._OPTICAL_KEYS[ri]
            v10  = entry.get(k10)
            v100 = entry.get(k100) if k100 else None
            if v10 is None and v100 is None:
                return "—"
            if k100 is None:
                return self._fmt(k10, v10)
            return f"{self._fmt(k10, v10)} / {self._fmt(k100, v100)}"
        else:
            key = self._COMP_KEYS[ri]
            return self._fmt(key, entry.get(key))

    def _refresh_report_table(self) -> None:
        row_labels = self._current_row_labels()
        self._report_table.setRowCount(len(row_labels))
        self._report_table.setColumnCount(2 + len(self._models))
        headers = ["Category", "Item"] + [f"{e['brand']}_{e['model']}" for e in self._models]
        self._report_table.setHorizontalHeaderLabels(headers)

        for ri, (section, item) in enumerate(row_labels):
            for ci, val in enumerate([section, item]):
                c = QTableWidgetItem(val)
                c.setTextAlignment(Qt.AlignmentFlag.AlignCenter)
                self._report_table.setItem(ri, ci, c)
            for mi, entry in enumerate(self._models):
                c = QTableWidgetItem(self._cell_value(entry, ri))
                c.setTextAlignment(Qt.AlignmentFlag.AlignCenter)
                self._report_table.setItem(ri, 2 + mi, c)

    def _table_to_text(self) -> str:
        row_labels = self._current_row_labels()
        lines = []
        headers = ["Category", "Item"] + [f"{e['brand']}_{e['model']}" for e in self._models]
        lines.append("\t".join(headers))
        for ri, (section, item) in enumerate(row_labels):
            row = [section, item]
            for entry in self._models:
                row.append(self._cell_value(entry, ri))
            lines.append("\t".join(row))
        return "\n".join(lines)

    def _copy_clipboard(self) -> None:
        ranges = self._report_table.selectedRanges()
        if ranges:
            min_row = min(r.topRow()     for r in ranges)
            max_row = max(r.bottomRow()  for r in ranges)
            min_col = min(r.leftColumn() for r in ranges)
            max_col = max(r.rightColumn() for r in ranges)
            lines = []
            for row in range(min_row, max_row + 1):
                cells = []
                for col in range(min_col, max_col + 1):
                    in_sel = any(
                        r.topRow() <= row <= r.bottomRow() and
                        r.leftColumn() <= col <= r.rightColumn()
                        for r in ranges
                    )
                    item = self._report_table.item(row, col) if in_sel else None
                    cells.append(item.text() if item else "")
                lines.append("\t".join(cells))
            text = "\n".join(lines)
        else:
            text = self._table_to_text()
        QApplication.clipboard().setText(text)
        QMessageBox.information(self, "Copied", "Copied to clipboard.")

    def _export_excel(self) -> None:
        if not self._models:
            QMessageBox.information(self, "Notice", "No models added.")
            return
        brand = self._engine.brand or "brand"
        model = self._engine.model_name or "model"
        path, _ = QFileDialog.getSaveFileName(
            self, "Save Excel", f"report_{brand}_{model}.xlsx", "Excel (*.xlsx)"
        )
        if path:
            ExcelExporter().export_report_template(brand, model, file_path=path)
            QMessageBox.information(self, "Saved", f"Saved: {path}")

    def _export_ppt(self) -> None:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        except ImportError:
            QMessageBox.critical(
                self, "python-pptx Not Installed",
                "python-pptx is required to save PPT.\n\n  pip install python-pptx"
            )
            return

        if not self._models:
            QMessageBox.information(self, "Notice", "No models added.")
            return

        brand = self._engine.brand or "brand"
        model = self._engine.model_name or "model"
        path, _ = QFileDialog.getSaveFileName(
            self, "Save PPT", f"report_{brand}_{model}.pptx", "PowerPoint (*.pptx)"
        )
        if not path:
            return

        # Linux 에서 필터가 확장자를 자동으로 안 붙이는 경우 대비
        if not path.lower().endswith(".pptx"):
            path += ".pptx"

        import tempfile, os
        from PySide6.QtGui import QImage, QPainter as _QP
        from PySide6.QtCore import QSize

        tmp_imgs: list[str] = []   # finally 에서 항상 참조 가능하도록 미리 선언

        try:
            def _render(widget) -> str:
                """QChartView는 scene.render로 직접 렌더링 (grab()은 OpenGL 내용 미포함)."""
                from PySide6.QtCharts import QChartView as _QCV
                from PySide6.QtWidgets import QApplication as _QApp
                from PySide6.QtCore import QRectF as _QRF

                if isinstance(widget, _QCV):
                    sw, sh = 1300, 800
                    grabbed = QImage(QSize(sw, sh), QImage.Format.Format_RGB32)
                    grabbed.fill(0xFFFFFF)
                    _pp = _QP(grabbed)
                    _pp.setRenderHint(_QP.RenderHint.Antialiasing)
                    chart = widget.chart()
                    scene = widget.scene()
                    _og = chart.geometry()
                    _os = scene.sceneRect()
                    try:
                        chart.setGeometry(_QRF(0, 0, sw, sh))
                        scene.setSceneRect(_QRF(0, 0, sw, sh))
                        _QApp.processEvents()
                        scene.render(_pp, _QRF(0, 0, sw, sh), _QRF(0, 0, sw, sh))
                    finally:
                        chart.setGeometry(_og)
                        scene.setSceneRect(_os)
                        _QApp.processEvents()
                    _pp.end()
                else:
                    sw = max(widget.width(), 700)
                    sh = max(widget.height(), 400)
                    grabbed = QImage(QSize(sw, sh), QImage.Format.Format_RGB32)
                    grabbed.fill(0xFFFFFF)
                    _pp = _QP(grabbed)
                    widget.render(_pp)
                    _pp.end()

                tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                tmp.close()
                grabbed.save(tmp.name)
                return tmp.name

            # ── matplotlib 공통 임포트 + 크로스플랫폼 폰트/DPI 설정 ─────────
            try:
                import matplotlib as _mpl
                _mpl.use("Agg")
                import matplotlib.pyplot as _plt
                import matplotlib.lines as _mlines
                import matplotlib.font_manager as _fm

                # 플랫폼별 한글 폰트 탐색 (우선순위 순)
                _FONT_CANDIDATES = [
                    # Linux — Noto CJK
                    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
                    "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc",
                    "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
                    # Windows — Malgun Gothic (한글 기본 내장 폰트)
                    r"C:\Windows\Fonts\malgun.ttf",
                    r"C:\Windows\Fonts\malgunbd.ttf",
                    # macOS
                    "/System/Library/Fonts/AppleSDGothicNeo.ttc",
                    "/Library/Fonts/NanumGothic.ttf",
                ]
                _FONT_FAMILIES = {
                    "NotoSansCJK-Regular.ttc": "Noto Sans CJK JP",
                    "malgun.ttf":   "Malgun Gothic",
                    "malgunbd.ttf": "Malgun Gothic",
                    "AppleSDGothicNeo.ttc": "Apple SD Gothic Neo",
                    "NanumGothic.ttf": "NanumGothic",
                }
                for _fp in _FONT_CANDIDATES:
                    if os.path.exists(_fp):
                        _fm.fontManager.addfont(_fp)
                        _fname = os.path.basename(_fp)
                        _family = _FONT_FAMILIES.get(_fname, "sans-serif")
                        _mpl.rcParams["font.family"] = _family
                        break

                # 마이너스 기호·DPI 일관성 설정
                _mpl.rcParams["axes.unicode_minus"] = False
                _mpl.rcParams["figure.dpi"] = 100   # 플랫폼 DPI 무관하게 100 고정

                _HAS_MPL = True
            except Exception as _mpl_err:
                import traceback as _tb
                print(f"[PPT] matplotlib 비활성화 — 원인: {_mpl_err}")
                _tb.print_exc()
                _HAS_MPL = False
                QMessageBox.warning(
                    self, "PPT 렌더링 품질 저하",
                    "matplotlib 을 불러올 수 없어 Qt 위젯 캡처로 대체합니다.\n"
                    "차트 해상도·범례 양식이 Linux 빌드와 다를 수 있습니다.\n\n"
                    f"원인: {_mpl_err}\n\n"
                    "해결: pip install matplotlib"
                )

            # ── APL 로딩 차트: matplotlib 고해상도 렌더링 ─────────────────────
            def _render_apl_mpl(mode_key: str, title: str, fallback_widget) -> str:
                """APL 로딩 꺾은선 차트를 matplotlib 으로 고해상도 PNG 생성.

                mode_key: "apl_sdr" 또는 "apl_hdr"
                """
                if not _HAS_MPL:
                    return _render(fallback_widget)

                _fig, _ax = _plt.subplots(figsize=(6.5, 3.8))  # APL 차트 크기 (인치): 가로×세로, dpi=200 → 1300×760 px
                _fig.patch.set_facecolor("white")
                _ax.set_facecolor("white")

                _all_lvs: list[float] = []
                for _mi, _entry in enumerate(self._models[:5]):
                    _apl_data = _entry.get(mode_key, {})
                    if not _apl_data:
                        continue
                    _key = f"{_entry['brand']}_{_entry['model']}"
                    _col = self._model_colors.get(
                        _key, _DEFAULT_MODEL_COLORS[_mi % len(_DEFAULT_MODEL_COLORS)])
                    _label = f"{_entry['brand']} {_entry['model']}"
                    _apls = sorted(_apl_data.keys())
                    _lvs  = [_apl_data[a] for a in _apls]
                    _all_lvs.extend(_lvs)
                    _ax.plot(_apls, _lvs, marker="", linewidth=1.5,
                             color=_col, label=_label)

                # 데이터가 없으면 Qt 위젯 scene.render 로 fallback
                if not _all_lvs:
                    _plt.close(_fig)
                    return _render(fallback_widget)

                # Y축: 0 고정 시작, max + 100 nit 여유 / X축은 autoscale 유지
                _y_top = max(_all_lvs) + 100.0
                _ax.set_xlim(left=0)
                _ax.set_ylim(bottom=0.0, top=_y_top)

                _ax.set_xlabel("APL (%)", fontsize=10)   # X/Y축 레이블 폰트 10pt
                _ax.set_ylabel("Lv (cd/m²)", fontsize=10)
                _ax.set_title(title, fontsize=11, fontweight="bold")  # 제목 11pt bold
                _ax.legend(fontsize=9, loc="upper center",
                           bbox_to_anchor=(0.5, -0.18), ncol=3,
                           frameon=True, edgecolor="#cccccc")  # 범례: 그래프 아래쪽
                _ax.tick_params(labelsize=10)
                _ax.grid(True, linestyle="--", alpha=0.4)  # 그리드: 점선, 불투명도 40%
                _fig.subplots_adjust(left=0.12, right=0.97, top=0.92, bottom=0.30)

                _tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                _tmp.close()
                _fig.savefig(_tmp.name, dpi=200, facecolor='white', edgecolor='none')
                _plt.close(_fig)
                return _tmp.name

            img_sdr = _render_apl_mpl("apl_sdr", "SDR Vivid 모드", self._apl_view_sdr)
            tmp_imgs.append(img_sdr)
            img_hdr = _render_apl_mpl("apl_hdr", "HDR Vivid 모드", self._apl_view_hdr)
            tmp_imgs.append(img_hdr)

            # ── Gamut 차트: matplotlib 고해상도 렌더링 ─────────────────────────
            # - DCI-P3 (#aab0c0 점선) / BT.2020 (#aab0c0 실선) → 보고서 Qt 차트와 동일
            # - 모델 삼각형 → self._model_colors (보고서 지정 색상)
            # - 모델 범례 → 그래프 안 오른쪽 아래 고정
            # - 반환: (이미지 경로, 꼭짓점 픽셀 fraction dict) → 줌 커트 정밀 크롭
            # ── Gamut vfracs 해석적 계산 (matplotlib 유무와 무관) ─────────────────
            # Gamut 차트는 xlim/ylim=(0,0.65), subplots_adjust 고정값으로
            # 데이터 좌표 → figure fraction 을 직접 계산 가능
            def _uv_to_frac(u: float, v: float) -> "tuple[float, float]":
                _adj_l, _adj_r = 0.13, 0.95
                _adj_b, _adj_t = 0.22, 0.91
                xf = _adj_l + (u / 0.65) * (_adj_r - _adj_l)
                yf_b = _adj_b + (v / 0.65) * (_adj_t - _adj_b)
                return max(0.01, min(0.99, xf)), max(0.01, min(0.99, 1.0 - yf_b))

            _gmt_vfracs: dict = {}
            for _ent in self._models[:5]:
                _uv0 = _ent.get("gamut_uv", {})
                _pts0 = {c: _uv0[c] for c in ("red", "green", "blue") if c in _uv0}
                if len(_pts0) == 3:
                    for _c0, (_u0, _v0) in _pts0.items():
                        _gmt_vfracs[_c0] = _uv_to_frac(_u0, _v0)
                    break
            if not _gmt_vfracs:
                # BT.2020 기준 위치 fallback — green v'=0.7952 > ylim(0.65) 클램핑
                for _c0, (_u0, _v0) in {
                    "red":   (0.5556, 0.3344),
                    "green": (0.2095, 0.6400),
                    "blue":  (0.1301, 0.0456),
                }.items():
                    _gmt_vfracs[_c0] = _uv_to_frac(_u0, _v0)

            def _render_gamut_mpl() -> str:
                if not _HAS_MPL:
                    return _render(self._gamut_chart_view)

                from core.gamut_utils import DCI_P3_UV, BT2020_UV

                _DPI = 200              # Gamut 차트 저장 DPI — 높을수록 선명, 200이 파일크기·품질 균형점
                _FW, _FH = 4.5, 4.5   # 정사각형 figure (인치): dpi=200 → 저장 시 900×900 px
                _fig, _ax = _plt.subplots(figsize=(_FW, _FH))
                _fig.patch.set_facecolor("white")
                _ax.set_facecolor("white")
                _ax.set_xlim(0.0, 0.65)
                _ax.set_ylim(0.0, 0.65)
                _ax.set_aspect("equal")
                # subplots_adjust로 하단 범례 공간 확보 (bbox_inches 미사용이므로 필수)
                _fig.subplots_adjust(left=0.13, right=0.95, top=0.91, bottom=0.22)  # 하단 범례(DCI-P3/BT.2020) 공간 확보 — bbox_inches 미사용 시 필수

                # 기준 색역 삼각형 — Qt 보고서와 동일 색상
                #   DCI-P3: #aab0c0 점선,  BT.2020: #aab0c0 실선
                def _rt(pts, ls):
                    xs = [p[0] for p in pts] + [pts[0][0]]
                    ys = [p[1] for p in pts] + [pts[0][1]]
                    _ax.plot(xs, ys, color="#aab0c0", linestyle=ls, linewidth=1.2)  # 기준 색역 선: 회색 #aab0c0, 두께 1.2pt

                _rt(DCI_P3_UV, "--")
                _rt(BT2020_UV, "-")

                # 측정 모델 삼각형 (보고서 지정 색상)
                _handles = []
                _first_verts: dict = {}
                for _mi, _ent in enumerate(self._models[:5]):
                    _uv = _ent.get("gamut_uv", {})
                    if not _uv:
                        continue
                    _k = f"{_ent['brand']}_{_ent['model']}"
                    _col = self._model_colors.get(
                        _k, _DEFAULT_MODEL_COLORS[_mi % len(_DEFAULT_MODEL_COLORS)])
                    _lbl = f"{_ent['brand']} {_ent['model']}"
                    _pts = [(_uv[c][0], _uv[c][1]) for c in ("red","green","blue") if c in _uv]
                    if len(_pts) == 3:
                        _tri = _pts + [_pts[0]]
                        _ax.plot([p[0] for p in _tri], [p[1] for p in _tri],
                                 color=_col, linewidth=1.5)  # 측정 모델 삼각형 선 두께 1.5pt
                        _handles.append(_mlines.Line2D([0],[0], color=_col,
                                                       linewidth=1.5, label=_lbl))
                        if not _first_verts:
                            _first_verts = {c: _uv[c] for c in ("red","green","blue") if c in _uv}
                    # 꼭짓점 점 표시 제거

                # 모델 범례: 그래프 안 오른쪽 아래에 고정
                if _handles:
                    _ax.legend(handles=_handles,
                               loc="lower right",
                               bbox_to_anchor=(0.99, 0.01),   # axes 기준 오른쪽(0.99) 하단(0.01) 고정
                               fontsize=10, frameon=True, framealpha=0.88)  # 범례 배경 불투명도 88%

                _ax.set_xlabel("u'", fontsize=10)
                _ax.set_ylabel("v'", fontsize=10)
                _ax.set_title("u'v' Chromaticity", fontsize=11, fontweight="bold")
                _ax.tick_params(labelsize=10)
                _ax.grid(True, linestyle="--", alpha=0.3)  # 그리드: 점선, 불투명도 30%

                # 하단 별도 범례: DCI-P3(점선) / BT.2020(실선) — 보고서와 동일
                _ref_h = [
                    _mlines.Line2D([0],[0], color="#aab0c0", linestyle="--",
                                   linewidth=1.2, label="DCI-P3"),  # 하단 참조 범례 선 두께
                    _mlines.Line2D([0],[0], color="#aab0c0", linestyle="-",
                                   linewidth=1.2, label="BT.2020"),
                ]
                _fig.legend(handles=_ref_h, loc="lower center", ncol=2,
                            fontsize=10, bbox_to_anchor=(0.5, 0.01),  # 하단 범례: figure 기준 중앙(0.5) 최하단(0.01)
                            frameon=True, edgecolor="#c0c0c0")

                _tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                _tmp.close()
                _fig.savefig(_tmp.name, dpi=_DPI, facecolor='white', edgecolor='none')
                _plt.close(_fig)
                return _tmp.name

            img_gmt = _render_gamut_mpl()
            tmp_imgs.append(img_gmt)

            prs = Presentation()
            # 와이드스크린 16:9  (13.33" × 7.5")
            SW = 13.33  # 슬라이드 가로 (인치) — 16:9 와이드스크린 표준
            SH = 7.5    # 슬라이드 세로 (인치)
            prs.slide_width  = Inches(SW)
            prs.slide_height = Inches(SH)

            BLANK = 6
            for _i, _lay in enumerate(prs.slide_layouts):
                if "blank" in _lay.name.lower():
                    BLANK = _i
                    break

            slide = prs.slides.add_slide(prs.slide_layouts[BLANK])

            from pptx.oxml.ns import qn as _qn
            from lxml import etree as _et

            def _hex(h: str) -> RGBColor:
                h = h.lstrip("#")
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

            def _tb(text: str, x: float, y: float, w: float, h: float,
                    size: int = 10, bold: bool = False,
                    color: str = "000000", align=PP_ALIGN.LEFT) -> None:
                tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                tf = tb.text_frame
                tf.word_wrap = False
                para = tf.paragraphs[0]
                para.alignment = align
                run = para.add_run()
                run.text = text
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = _hex(color)

            def _set_border(cell) -> None:
                """셀 4면 검정 1pt 테두리."""
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for side in ("lnL", "lnR", "lnT", "lnB"):
                    tag = _qn(f"a:{side}")
                    for e in tcPr.findall(tag):
                        tcPr.remove(e)
                    ln = _et.SubElement(tcPr, tag, w="12700")
                    sf = _et.SubElement(ln, _qn("a:solidFill"))
                    _et.SubElement(sf, _qn("a:srgbClr"), val="000000")
                    _et.SubElement(ln, _qn("a:prstDash"), val="solid")

            # ── ① 제목: 왼쪽 배치 + 아래 구분 실선 ──────────────────────
            TITLE_H  = 0.40             # 제목 텍스트박스 높이 (인치)
            LINE_Y   = TITLE_H + 0.08  # 구분선 Y 위치: 제목 아래 0.08인치
            models_str = ",  ".join(
                f"{e['brand']} {e['model']}" for e in self._models[:5]
            )
            # 제목 텍스트 (왼쪽 절반만 사용 — 명확히 왼쪽에 위치)
            # 제목 위치: x=0.02인치, y=0.04인치, 너비=슬라이드 60% / 폰트 18pt bold
            _tb(models_str, 0.02, 0.04, SW * 0.6, TITLE_H,
                size=18, bold=True, color="000000", align=PP_ALIGN.LEFT)

            # 가로 끝에서 끝까지 가는 실선
            from pptx.enum.shapes import MSO_CONNECTOR_TYPE
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR_TYPE.STRAIGHT,
                Inches(0.0), Inches(LINE_Y),
                Inches(SW),  Inches(LINE_Y)
            )
            connector.line.color.rgb = _hex("000000")  # type: ignore[assignment]
            connector.line.width = Pt(0.5)  # 구분선 두께 0.5pt — 얇은 실선

            # ── ② 데이터 테이블 (최대 5개 모델, 항목 열 분할) ──────────────
            row_labels  = self._current_row_labels()
            n_data_rows = len(row_labels)
            _EXTRA_ROWS = ["Reflectance [%]  (SCI/SCE)", "Local Dimming Count  [block]"]
            n_rows      = n_data_rows + 1 + len(_EXTRA_ROWS)
            n_models    = min(len(self._models), 5)
            # 열 구성: 구분(0) | 항목-그룹(1) | 항목-세부(2) | 모델들(3+)
            n_cols      = 3 + n_models

            CAT_W  = 1.1    # 구분 열 너비 (인치)
            GRP_W  = 0.85   # HDR/SDR 그룹 열 너비 (인치)
            DTL_W  = 0.95   # 10%/100% 세부 열 너비 (인치)
            MDL_W  = 1.9    # 모델 데이터 열 너비 (인치) — 5개 모델 기준
            TBL_W  = min(SW - 0.4, CAT_W + GRP_W + DTL_W + n_models * MDL_W)
            TBL_X  = (SW - TBL_W) / 2
            TBL_Y  = LINE_Y + 0.12  # 표 Y 시작: 구분선에서 0.12인치 아래
            ROW_H  = 0.27           # 표 행 높이 (인치)
            TBL_H  = n_rows * ROW_H

            tbl_shape = slide.shapes.add_table(
                n_rows, n_cols,
                Inches(TBL_X), Inches(TBL_Y), Inches(TBL_W), Inches(TBL_H)
            )
            tbl = tbl_shape.table
            tbl.columns[0].width = Inches(CAT_W)
            tbl.columns[1].width = Inches(GRP_W)
            tbl.columns[2].width = Inches(DTL_W)
            for ci in range(3, n_cols):
                tbl.columns[ci].width = Inches(MDL_W)

            # Windows PowerPoint 는 기본 테이블 스타일이 셀 테두리를 덮어씀 → 스타일 제거
            _tblPr = tbl._tbl.find(_qn("a:tblPr"))
            if _tblPr is None:
                _tblPr = _et.SubElement(tbl._tbl, _qn("a:tblPr"))
            _styleId = _tblPr.find(_qn("a:tableStyleId"))
            if _styleId is not None:
                _tblPr.remove(_styleId)

            # HDR/SDR 분할 적용 대상 (비광학 포맷)
            _SPLIT_MAP = {
                "HDR 10%":  ("HDR", "10%"),
                "HDR 100%": ("HDR", "100%"),
                "SDR 10%":  ("SDR", "10%"),
                "SDR 100%": ("SDR", "100%"),
            }
            is_optical = self._is_optical_format()

            def _cell(row: int, col: int, text: str,
                      bold: bool = False, size: int = 12,  # 셀 기본 폰트 12pt
                      align=PP_ALIGN.CENTER, color: str = "000000") -> None:
                c = tbl.cell(row, col)
                c.text = text
                c.fill.background()
                tf2 = c.text_frame
                tf2.paragraphs[0].alignment = align
                run = (tf2.paragraphs[0].runs[0]
                       if tf2.paragraphs[0].runs
                       else tf2.paragraphs[0].add_run())
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = _hex(color)

            def _style_merged(mc, text: str, bold=True,
                              align=PP_ALIGN.CENTER) -> None:
                mc.text = text
                mc.fill.background()
                tf = mc.text_frame
                tf.paragraphs[0].alignment = align
                run = (tf.paragraphs[0].runs[0]
                       if tf.paragraphs[0].runs
                       else tf.paragraphs[0].add_run())
                run.font.size = Pt(12)
                run.font.bold = bold
                run.font.color.rgb = _hex("000000")

            # ── 헤더 ─────────────────────────────────────────────────────
            _cell(0, 0, "Category", bold=True)
            # "Item" 헤더: col 1+2 병합
            tbl.cell(0, 1).merge(tbl.cell(0, 2))
            _style_merged(tbl.cell(0, 1), "Item", align=PP_ALIGN.CENTER)
            for mi, entry in enumerate(self._models[:5]):
                _cell(0, 3 + mi,
                      f"{entry['brand']}_{entry['model']}", bold=True)

            # ── 데이터 행 ────────────────────────────────────────────────
            for ri, (section, item) in enumerate(row_labels, 1):
                _cell(ri, 0, section, bold=True)
                split = _SPLIT_MAP.get(item) if not is_optical else None
                if split:
                    grp, dtl = split
                    _cell(ri, 1, grp, bold=True)
                    _cell(ri, 2, dtl)
                else:
                    # 분할 없는 항목: col 1+2 수평 병합
                    tbl.cell(ri, 1).merge(tbl.cell(ri, 2))
                    _style_merged(tbl.cell(ri, 1), item,
                                  bold=False, align=PP_ALIGN.LEFT)
                for mi, entry in enumerate(self._models[:5]):
                    _cell(ri, 3 + mi,
                          self._cell_value(entry, ri - 1))

            # ── 구분(col 0) 세로 병합 ─────────────────────────────────────
            i = 0
            while i < n_data_rows:
                section = row_labels[i][0]
                j = i + 1
                while j < n_data_rows and row_labels[j][0] == section:
                    j += 1
                if j - i > 1:
                    tbl.cell(i + 1, 0).merge(tbl.cell(j, 0))
                    _style_merged(tbl.cell(i + 1, 0), section)
                i = j

            # ── col 1(그룹) HDR·SDR 세로 병합 (비광학 포맷) ──────────────
            if not is_optical:
                # row_labels의 item은 "HDR 10%"/"SDR 100%" 형태라서
                # 검색 prefix는 HDR/SDR을 쓰고, 병합 셀 표시는 HDR vivid/SDR vivid로 한다.
                for item_prefix, merged_label in (("HDR", "HDR vivid"), ("SDR", "SDR vivid")):
                    grp_rows = [ri + 1 for ri, (_, item) in enumerate(row_labels)
                                if item.startswith(item_prefix + " ")]
                    if len(grp_rows) > 1:
                        tbl.cell(grp_rows[0], 1).merge(tbl.cell(grp_rows[-1], 1))
                        _style_merged(tbl.cell(grp_rows[0], 1), merged_label)

            # ── 하단 추가 2행: col 0~2 수평 병합 ─────────────────────────
            for ei, label in enumerate(_EXTRA_ROWS):
                extra_row = n_data_rows + 1 + ei
                for ci in range(3, n_cols):
                    _cell(extra_row, ci, "")
                tbl.cell(extra_row, 0).merge(tbl.cell(extra_row, 2))
                _style_merged(tbl.cell(extra_row, 0), label,
                              bold=True, align=PP_ALIGN.LEFT)

            # 전체 셀 검정 테두리 + 폰트 12pt 강제
            for ri in range(n_rows):
                for ci in range(n_cols):
                    _set_border(tbl.cell(ri, ci))
                    try:
                        tf3 = tbl.cell(ri, ci).text_frame
                        for para in tf3.paragraphs:
                            if not para.runs:
                                run = para.add_run()
                                run.text = " "
                                run.font.size = Pt(12)  # 빈 셀 공백 채움 폰트 — 행 높이 유지용
                            else:
                                for run in para.runs:
                                    if run.text == "":
                                        run.text = " "
                                    run.font.size = Pt(12)
                    except Exception:
                        pass

            # ── ③ 표 아래 SDR | HDR | Gamut — 외부 레이블 없음(차트 내 제목 사용)
            CHART_AREA_Y = TBL_Y + TBL_H + 0.7           # 차트 시작 Y: 표 아래 0.7인치 간격 — 값 변경으로 차트 위치 조정
            CHART_H      = min(3.0, SH - CHART_AREA_Y - 0.05)  # 차트 높이: 최대 3인치, 하단 여백 0.05인치
            CHART_W      = SW / 3.0                       # SDR/HDR/Gamut 각 차트 슬롯 너비 (슬라이드를 3등분)

            def _add_gamut_zoom_cuts(img_path: str, gx: float, gy: float,
                                        gside: float, slot_w: float,
                                        vfracs: dict) -> None:
                """Gamut 꼭짓점 확대 컷 3개 추가.

                vfracs: {"red":(xf,yf), "green":(xf,yf), "blue":(xf,yf)}
                  xf, yf = 이미지 내 figure fraction (y: 위=0)
                크롭 윈도우 half_zoom 을 꼭짓점 중심으로 계산 → 정확한 위치 커트.
                """
                zoom_side = min(0.95, gside * 0.38)  # 줌 컷 박스 한 변 길이: 메인 차트의 38%, 최대 0.95인치
                zoom_margin = 0.04                   # 줌 컷과 슬롯 모서리 사이 여백 (인치)
                half = 0.13   # 이미지 크롭 반경 (figure fraction): 꼭짓점 ±13% 범위를 표시

                # (PPT 슬라이드 배치 x, y) , 꼭짓점 색상
                corners = [
                    (gx + zoom_margin,                     gy + zoom_margin,                    "green"),
                    (gx + slot_w - zoom_side - zoom_margin, gy + zoom_margin,                   "red"),
                    (gx + zoom_margin,                     gy + gside - zoom_side - zoom_margin, "blue"),
                ]

                for zx, zy, color in corners:
                    if color not in vfracs:
                        continue
                    xc, yc = vfracs[color]   # figure fraction (y: 위=0)
                    crop_l = max(0.0, min(0.85, xc - half))
                    crop_r = max(0.0, min(0.85, 1.0 - xc - half))
                    crop_t = max(0.0, min(0.85, yc - half))
                    crop_b = max(0.0, min(0.85, 1.0 - yc - half))
                    # crop 합이 1 이상이면 유효하지 않으므로 skip
                    if crop_l + crop_r >= 0.98 or crop_t + crop_b >= 0.98:
                        continue
                    pic = slide.shapes.add_picture(
                        img_path, Inches(zx), Inches(zy),
                        Inches(zoom_side), Inches(zoom_side))
                    pic.crop_left   = crop_l
                    pic.crop_top    = crop_t
                    pic.crop_right  = crop_r
                    pic.crop_bottom = crop_b

            for idx, img_path in enumerate([img_sdr, img_hdr, img_gmt]):
                cx = idx * CHART_W
                if idx == 2:
                    # Gamut: 정사각형 슬롯 (figure 도 4.5×4.5 정사각형)
                    gmt_side = min(CHART_W - 0.06, CHART_H)  # Gamut 차트 표시 크기: 슬롯 너비-0.06인치 또는 높이 중 작은 값 (정사각형 유지)
                    gmt_off  = (CHART_W - gmt_side) / 2      # Gamut 차트 좌우 중앙 정렬 오프셋
                    slide.shapes.add_picture(
                        img_path,
                        Inches(cx + gmt_off), Inches(CHART_AREA_Y),
                        Inches(gmt_side), Inches(gmt_side)
                    )
                    _add_gamut_zoom_cuts(
                        img_path, cx, CHART_AREA_Y, gmt_side, CHART_W, _gmt_vfracs
                    )
                else:
                    slide.shapes.add_picture(
                        img_path,
                        Inches(cx + 0.04), Inches(CHART_AREA_Y),
                        Inches(CHART_W - 0.08), Inches(CHART_H)
                    )

            prs.save(path)
            QMessageBox.information(self, "Saved", f"PPT saved:\n{path}")

        except Exception as _err:
            import traceback
            QMessageBox.critical(self, "PPT Save Error",
                                 f"{_err}\n\n{traceback.format_exc()}")
        finally:
            for p in tmp_imgs:
                try:
                    os.unlink(p)
                except OSError:
                    pass

    # ── 파일 불러오기 ────────────────────────────────────────────────────

    def _load_files(self) -> None:
        paths, _ = QFileDialog.getOpenFileNames(
            self, "Select Excel Files", "", "Excel (*.xlsx)"
        )
        if not paths:
            return
        errors: List[str] = []
        for path in paths:
            try:
                self._parse_xlsx(path)
            except Exception as exc:
                errors.append(f"{os.path.basename(path)}: {exc}")
        if errors:
            QMessageBox.warning(self, "File Error", "\n".join(errors))
        self._refresh_report_table()
        self._refresh_apl_chart()
        self._refresh_gamut_chart()

    def _find_or_create_entry(self, brand: str, model: str) -> Dict:
        for entry in self._models:
            if entry["brand"] == brand and entry["model"] == model:
                return entry
        key = f"{brand}_{model}"
        if key not in self._model_colors:
            idx = len(self._models)
            self._model_colors[key] = _DEFAULT_MODEL_COLORS[idx % len(_DEFAULT_MODEL_COLORS)]
        entry: Dict = {
            "brand": brand, "model": model,
            "hdr_10": None, "hdr_100": None,
            "sdr_10": None, "sdr_100": None,
            "contrast_ratio": None, "black_lv": None,
            "dci_overlap": None, "bt2020_overlap": None,
            "apl_hdr": {}, "apl_sdr": {}, "gamut_uv": {}, "gamut_uv_sdr": {}, "gamut_uv_hdr": {},
            # per-mode (광학 측정 데이터 format)
            "sdr_vivid_10": None,    "sdr_vivid_100": None,
            "sdr_standard_10": None, "sdr_standard_100": None,
            "hdr_vivid_10": None,    "hdr_vivid_100": None,
            "hdr_standard_10": None, "hdr_standard_100": None,
            "hdr_cinema_10": None,   "hdr_cinema_100": None,
        }
        self._models.append(entry)
        color = self._model_colors[key]
        self._model_list.addItem(f"  {brand}  {model}")
        item = self._model_list.item(self._model_list.count() - 1)
        item.setForeground(QColor(color))
        return entry

    def _parse_xlsx(self, path: str) -> None:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        try:
            sheet_names = wb.sheetnames
            brand = ""
            model = ""
            sequence = ""

            if "Info" in sheet_names:
                info_ws = wb["Info"]
                for row in info_ws.iter_rows(min_row=1, max_row=5, values_only=True):
                    if not row or row[0] is None:
                        continue
                    key = str(row[0]).strip()
                    val = str(row[1] or "").strip() if len(row) > 1 else ""
                    if key == "Brand":
                        brand = val
                    elif key == "Model":
                        model = val
                    elif key == "Sequence":
                        sequence = val

            if not sequence:
                if "Summary" in sheet_names and any(s.startswith("Raw_") for s in sheet_names):
                    sequence = "Luminance Loading"
                elif "Gamut" in sheet_names:
                    sequence = "Gamut"
                else:
                    raise ValueError(f"Unrecognized file format (sheets: {sheet_names})")

            if not brand or not model:
                basename = os.path.splitext(os.path.basename(path))[0]
                parts = basename.split("_")
                lower = basename.lower()
                if lower.startswith("lum_loading_") and len(parts) >= 5:
                    brand = brand or parts[3]
                    model = model or "_".join(parts[4:])
                elif lower.startswith("gamut_") and len(parts) >= 3:
                    brand = brand or parts[1]
                    model = model or "_".join(parts[2:])
                else:
                    brand = brand or "Unknown"
                    model = model or basename

            # 시퀀스 자동 감지 보완 — all 파일은 여러 시트를 포함
            if not sequence:
                if any(s.startswith("Loading_") for s in sheet_names):
                    sequence = "All Sessions"
                elif "Summary" in sheet_names and any(s.startswith("Raw_") for s in sheet_names):
                    sequence = "Luminance Loading"
                elif any(s.startswith("Gamut") for s in sheet_names):
                    sequence = "Gamut"
                else:
                    raise ValueError(f"Unrecognized file format (sheets: {sheet_names})")

            entry = self._find_or_create_entry(brand, model)

            if "All Sessions" in sequence:
                self._parse_all_sessions_wb(wb, entry)
            elif "Luminance Loading" in sequence:
                basename_lower = os.path.basename(path).lower()
                is_hdr = "hdr" in basename_lower
                mode = ("cinema" if "cinema" in basename_lower
                        else "standard" if ("standard" in basename_lower or "_std_" in basename_lower)
                        else "vivid")
                self._parse_lum_loading_wb(wb, entry, is_hdr, mode=mode)
            elif "Gamut" in sequence:
                self._parse_gamut_wb(wb, entry)
            else:
                raise ValueError(f"Unknown sequence: {sequence!r}")
        finally:
            wb.close()

    def _parse_lum_loading_wb(self, wb: Any, entry: Dict, is_hdr: bool,
                              mode: str = "") -> None:
        sheet_names = wb.sheetnames
        raw_sheets = [s for s in sheet_names if s.startswith("Raw_")]
        agg = self._agg_combo.currentText()  # "최대값" / "중간값" / "최소값"

        apl_dict: Dict[int, float] = {}

        if raw_sheets:
            apl_lv: Dict[int, List[float]] = {}
            for sheet_name in raw_sheets:
                ws = wb[sheet_name]
                for row in ws.iter_rows(min_row=2, values_only=True):
                    if not row or row[0] is None:
                        continue
                    try:
                        apl = round(float(row[0]))
                        lv = float(row[2])  # col3 = Lv (cd/m²)
                    except (TypeError, ValueError, IndexError):
                        continue
                    apl_lv.setdefault(apl, []).append(lv)

            for apl, lvs in apl_lv.items():
                if agg == "Max":
                    apl_dict[apl] = round(max(lvs), 3)
                elif agg == "Min":
                    apl_dict[apl] = round(min(lvs), 3)
                else:
                    apl_dict[apl] = round(statistics.median(lvs), 3)

        elif "Summary" in sheet_names:
            col_idx = {"Max": 2, "Min": 3}.get(agg, 1)  # 1=Avg, 2=Max, 3=Min
            ws = wb["Summary"]
            for row in ws.iter_rows(min_row=2, values_only=True):
                if not row or row[0] is None:
                    continue
                try:
                    apl = round(float(row[0]))
                    lv = float(row[col_idx])
                except (TypeError, ValueError, IndexError):
                    continue
                apl_dict[apl] = round(lv, 3)
        else:
            raise ValueError("Could not find Raw_ or Summary sheet.")

        if is_hdr:
            entry["apl_hdr"].update(apl_dict)
            if 10  in apl_dict: entry["hdr_10"]  = apl_dict[10]
            if 100 in apl_dict: entry["hdr_100"] = apl_dict[100]
        else:
            entry["apl_sdr"].update(apl_dict)
            if 10  in apl_dict: entry["sdr_10"]  = apl_dict[10]
            if 100 in apl_dict: entry["sdr_100"] = apl_dict[100]

        # per-mode keys (광학 측정 데이터 format)
        m = mode.lower()
        if is_hdr:
            if "cinema" in m:
                pfx = "hdr_cinema"
            elif "standard" in m or "std" in m:
                pfx = "hdr_standard"
            else:
                pfx = "hdr_vivid"
        else:
            if "standard" in m or "std" in m:
                pfx = "sdr_standard"
            else:
                pfx = "sdr_vivid"
        if 10  in apl_dict: entry[f"{pfx}_10"]  = apl_dict[10]
        if 100 in apl_dict: entry[f"{pfx}_100"] = apl_dict[100]

    def _parse_gamut_wb(self, wb: Any, entry: Dict) -> None:
        try:
            gamut_ws = wb["Gamut"]
        except KeyError:
            raise ValueError("Could not find Gamut sheet.")

        label_map = {
            "Red": "red", "Green": "green", "Blue": "blue",
            "White": "white", "Black": "black",
        }
        rows = list(gamut_ws.iter_rows(min_row=2, max_row=6, values_only=True))
        for row in rows:
            if row[0] is None:
                continue
            color_label = label_map.get(str(row[0]).strip())
            if color_label is None:
                continue
            try:
                u_prime = float(row[5])
                v_prime = float(row[6])
            except (TypeError, ValueError, IndexError):
                continue
            entry["gamut_uv"][color_label] = (u_prime, v_prime)

        uv = entry["gamut_uv"]
        if "red" in uv and "green" in uv and "blue" in uv:
            stats = calc_gamut_stats(uv["red"], uv["green"], uv["blue"])
            entry["dci_overlap"] = stats.get("dci_overlap")
            entry["bt2020_overlap"] = stats.get("bt2020_overlap")

    def _parse_all_sessions_wb(self, wb: Any, entry: Dict) -> None:
        """_all.xlsx 파싱 — Loading_*, Gamut_* 시트를 자동 감지해 데이터 추출.

        시트별 컬럼 구조 (export_all_session 기준):
          Loading_SDR_Vivid : APL(0) | #(1) | Time(2) | Lv(3) | x(4) | y(5) | u'(6) | ...
          Gamut_SDR/HDR     : Color(0) | Time(1) | Lv(2) | x(3) | y(4) | u'(5) | v'(6) | ...
        """
        sheet_names = wb.sheetnames
        agg = self._agg_combo.currentText()
        label_map = {"Red": "red", "Green": "green", "Blue": "blue",
                     "White": "white", "Black": "black"}

        # ── Loading_* 시트 ────────────────────────────────────────────────
        loading_sheets = [s for s in sheet_names if s.startswith("Loading_Raw_")
                          or (s.startswith("Loading_") and not s == "Loading_Summary")]
        for sheet_name in loading_sheets:
            # 시트명에서 모드 판단: Loading_HDR_Vivid → HDR, Loading_SDR_Vivid → SDR
            upper = sheet_name.upper()
            is_hdr    = "_HDR_" in upper or upper.endswith("_HDR")
            is_cinema = "CINEMA"   in upper
            is_std    = "STANDARD" in upper or "_STD_" in upper

            apl_lv: Dict[int, List[float]] = {}
            ws = wb[sheet_name]
            for row in ws.iter_rows(min_row=2, values_only=True):
                if not row or row[0] is None:
                    continue
                try:
                    apl = round(int(float(row[0])))
                    lv  = float(row[3])   # APL(0) | #(1) | Time(2) | Lv(3)
                except (TypeError, ValueError, IndexError):
                    continue
                apl_lv.setdefault(apl, []).append(lv)

            apl_dict: Dict[int, float] = {}
            for apl, lvs in apl_lv.items():
                if agg == "Max":
                    apl_dict[apl] = round(max(lvs), 3)
                elif agg == "Min":
                    apl_dict[apl] = round(min(lvs), 3)
                else:
                    apl_dict[apl] = round(statistics.median(lvs), 3)

            if is_hdr:
                entry["apl_hdr"].update(apl_dict)
                if 10  in apl_dict: entry["hdr_10"]  = apl_dict[10]
                if 100 in apl_dict: entry["hdr_100"] = apl_dict[100]
                if is_cinema:
                    pfx = "hdr_cinema"
                elif is_std:
                    pfx = "hdr_standard"
                else:
                    pfx = "hdr_vivid"
            else:
                entry["apl_sdr"].update(apl_dict)
                if 10  in apl_dict: entry["sdr_10"]  = apl_dict[10]
                if 100 in apl_dict: entry["sdr_100"] = apl_dict[100]
                pfx = "sdr_standard" if is_std else "sdr_vivid"
            if 10  in apl_dict: entry[f"{pfx}_10"]  = apl_dict[10]
            if 100 in apl_dict: entry[f"{pfx}_100"] = apl_dict[100]

        # ── Gamut_* 시트 ──────────────────────────────────────────────────
        gamut_sheets = [s for s in sheet_names if s.startswith("Gamut_")]
        for sheet_name in gamut_sheets:
            upper = sheet_name.upper()
            is_hdr_gamut = "_HDR" in upper or upper.endswith("HDR")
            uv_key = "gamut_uv_hdr" if is_hdr_gamut else "gamut_uv_sdr"
            ws = wb[sheet_name]
            for row in ws.iter_rows(min_row=2, max_row=7, values_only=True):
                if not row or row[0] is None:
                    continue
                color_label = label_map.get(str(row[0]).strip().capitalize())
                if color_label is None:
                    continue
                try:
                    u_prime = float(row[5])   # Color(0)|Time(1)|Lv(2)|x(3)|y(4)|u'(5)|v'(6)
                    v_prime = float(row[6])
                except (TypeError, ValueError, IndexError):
                    continue
                entry[uv_key][color_label] = (u_prime, v_prime)
                entry["gamut_uv"][color_label] = (u_prime, v_prime)  # fallback 유지

        # gamut_uv_sdr 없으면 gamut_uv 로 채우기
        if not entry["gamut_uv_sdr"] and entry["gamut_uv"]:
            entry["gamut_uv_sdr"] = dict(entry["gamut_uv"])

        uv = entry["gamut_uv"]
        if "red" in uv and "green" in uv and "blue" in uv:
            stats = calc_gamut_stats(uv["red"], uv["green"], uv["blue"])
            entry["dci_overlap"]   = stats.get("dci_overlap")
            entry["bt2020_overlap"] = stats.get("bt2020_overlap")

        # ── Contrast_* 시트 ───────────────────────────────────────────────
        # CR = max(CR at 50%, 20%, 14.1%) — 가장 높은 명암비 값을 Report에 사용
        # 컬럼 구조 (export_all_session 기준):
        #   Black H/V %(0) | Lv(1) | CR(2) | Time(3) | Lv_col(4) | ...
        contrast_sheets = [s for s in sheet_names if s.startswith("Contrast_")]
        for sheet_name in contrast_sheets:
            ws = wb[sheet_name]
            white_lv: float | None = None
            window_lvs: Dict[float, float] = {}
            for row in ws.iter_rows(min_row=2, values_only=True):
                if not row or row[0] is None:
                    continue
                try:
                    label = str(row[0]).strip()
                    lv = float(row[1])
                except (TypeError, ValueError, IndexError):
                    continue
                if label == "Full White":
                    white_lv = lv
                else:
                    try:
                        window_lvs[float(label)] = lv
                    except ValueError:
                        pass
            if white_lv and window_lvs:
                best_cr: float | None = None
                best_blv: float | None = None
                for side_pct in [50.0, 20.0, 14.1]:
                    blv = window_lvs.get(side_pct)
                    if blv and blv > 0:
                        cr = white_lv / blv
                        if best_cr is None or cr > best_cr:
                            best_cr = cr
                            best_blv = blv
                if best_cr is not None:
                    entry["contrast_ratio"] = round(best_cr, 1)
                if best_blv is not None:
                    entry["black_lv"] = round(best_blv, 4)

    # ── 차트 헬퍼 ────────────────────────────────────────────────────────

    def _add_ref_gamuts_to(self, chart: "QChart", ax_u: "QValueAxis", ax_v: "QValueAxis") -> None:
        for ref_pts, color, name, dash in [
            (DCI_P3_UV, QColor("#aab0c0"), "DCI-P3",  True),
            (BT2020_UV, QColor("#aab0c0"), "BT.2020", False),
        ]:
            series = QLineSeries()
            series.setName(name)
            pen = series.pen()
            pen.setColor(color)
            pen.setWidth(1)
            if dash:
                pen.setStyle(Qt.PenStyle.DashLine)
            series.setPen(pen)
            for u, v in ref_pts:
                series.append(u, v)
            series.append(ref_pts[0][0], ref_pts[0][1])
            chart.addSeries(series)
            series.attachAxis(ax_u)
            series.attachAxis(ax_v)

    def _add_ref_gamuts(self) -> None:
        self._add_ref_gamuts_to(self._gamut_chart, self._gamut_axis_u, self._gamut_axis_v)

    def _position_gamut_legend_for(self, chart: "QChart") -> None:
        legend = chart.legend()
        plot = chart.plotArea()
        if plot.isNull():
            return
        legend_w = 118
        legend_h = min(150, max(64, 34 + len(self._models) * 24))
        x = plot.right() - legend_w - 8
        y = plot.bottom() - legend_h - 8
        legend.setGeometry(QRectF(x, y, legend_w, legend_h))
        legend.update()

    def _position_gamut_legend(self) -> None:
        self._position_gamut_legend_for(self._gamut_chart)

    def _refresh_apl_chart(self) -> None:
        self._apl_chart_sdr.removeAllSeries()
        self._apl_chart_hdr.removeAllSeries()
        lv_sdr: List[float] = []
        lv_hdr: List[float] = []

        for entry in self._models:
            model_key = f"{entry['brand']}_{entry['model']}"
            color_hex = self._model_colors.get(model_key, _DEFAULT_MODEL_COLORS[0])
            label = f"{entry['brand']}_{entry['model']}"

            for data_key, chart, axis_x, axis_y, lv_list in [
                ("apl_sdr", self._apl_chart_sdr,
                 self._apl_axis_x_sdr, self._apl_axis_y_sdr, lv_sdr),
                ("apl_hdr", self._apl_chart_hdr,
                 self._apl_axis_x_hdr, self._apl_axis_y_hdr, lv_hdr),
            ]:
                apl_dict: Dict[int, float] = entry.get(data_key, {})
                if not apl_dict or len(apl_dict) <= 2:
                    continue
                series = QLineSeries()
                series.setName(label)
                pen = series.pen()
                pen.setColor(QColor(color_hex))
                pen.setWidth(2)
                series.setPen(pen)
                for apl in sorted(apl_dict):
                    lv = apl_dict[apl]
                    series.append(float(apl), lv)
                    lv_list.append(lv)
                chart.addSeries(series)
                series.attachAxis(axis_x)
                series.attachAxis(axis_y)

        if lv_sdr:
            self._apl_axis_y_sdr.setRange(0, max(lv_sdr) * 1.15)
        if lv_hdr:
            self._apl_axis_y_hdr.setRange(0, max(lv_hdr) * 1.15)

    def _refresh_gamut_chart(self) -> None:
        self._gamut_chart.removeAllSeries()
        self._add_ref_gamuts_to(self._gamut_chart, self._gamut_axis_u, self._gamut_axis_v)

        for entry in self._models:
            model_key = f"{entry['brand']}_{entry['model']}"
            color_hex = self._model_colors.get(model_key, _DEFAULT_MODEL_COLORS[0])
            label = f"{entry['brand']}\n{entry['model']}"
            uv: Dict[str, tuple] = entry.get("gamut_uv", {})
            if not uv:
                continue
            if "red" in uv and "green" in uv and "blue" in uv:
                tri = QLineSeries()
                tri.setName(label)
                pen = tri.pen()
                pen.setColor(QColor(color_hex))
                pen.setWidth(2)
                tri.setPen(pen)
                for key in ("red", "green", "blue", "red"):
                    u, v = uv[key]
                    tri.append(u, v)
                self._gamut_chart.addSeries(tri)
                tri.attachAxis(self._gamut_axis_u)
                tri.attachAxis(self._gamut_axis_v)

        self._position_gamut_legend_for(self._gamut_chart)


# ---------------------------------------------------------------------------
# Main Window
# ---------------------------------------------------------------------------
# Auto-All Panel
# ---------------------------------------------------------------------------

_MODE_ROW_LABELS = ["SDR Vivid", "SDR Standard", "SDR Cinema",
                    "HDR Vivid", "HDR Standard", "HDR Cinema"]
_SEQ_COL_LABELS  = {"lum_swing": "Swing", "lum_loading": "Loading",
                    "gamut": "Gamut", "contrast": "Contrast"}
_SEQ_COL_ORDER   = ["lum_swing", "lum_loading", "gamut", "contrast"]


class CircularProgress(QWidget):
    """세그먼트 분할 도넛 형태의 원형 진행률 위젯."""

    # 완료 세그먼트 색상 그라데이션: 시작(청록) → 끝(파랑)
    _COLOR_START = QColor("#1abc9c")
    _COLOR_END   = QColor("#2196f3")
    _COLOR_BG    = QColor("#dde4f0")

    def __init__(self, total: int = 14, parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self._value = 0
        self._total = total
        self.setFixedSize(110, 110)  # 원형 진행 위젯 크기 110×110 px

    def setValue(self, v: int) -> None:
        self._value = max(0, min(v, self._total))
        self.update()

    def setTotal(self, t: int) -> None:
        self._total = t
        self.update()

    def paintEvent(self, _event: Any) -> None:
        import math
        p = QPainter(self)
        p.setRenderHint(QPainter.RenderHint.Antialiasing)

        w, h = self.width(), self.height()
        seg_thickness = 10   # 세그먼트 호 두께 (px) — 값 키우면 도넛 두꺼워짐
        gap_deg       = 4.5  # 세그먼트 사이 빈 공간 (도) — 값 키우면 간격 넓어짐
        margin        = seg_thickness // 2 + 3  # 캔버스 여백: 세그먼트 두께 절반 + 3px (선이 잘리지 않도록)
        rect = QRectF(margin, margin, w - margin * 2, h - margin * 2)

        n = max(self._total, 1)
        seg_deg = (360.0 / n) - gap_deg   # 각 세그먼트의 호 길이

        for i in range(n):
            # 12시(90°) 기준 시계 방향
            start_angle = 90.0 - i * (360.0 / n)
            span_angle  = -seg_deg  # 음수 = 시계 방향

            done = i < self._value
            if done:
                # 완료 세그먼트: 시작→끝 색상 선형 보간
                t = i / (n - 1) if n > 1 else 0
                r = int(self._COLOR_START.red()   + (self._COLOR_END.red()   - self._COLOR_START.red())   * t)
                g = int(self._COLOR_START.green() + (self._COLOR_END.green() - self._COLOR_START.green()) * t)
                b = int(self._COLOR_START.blue()  + (self._COLOR_END.blue()  - self._COLOR_START.blue())  * t)
                color = QColor(r, g, b)
            else:
                color = self._COLOR_BG

            pen = QPen(color, seg_thickness)
            pen.setCapStyle(Qt.PenCapStyle.RoundCap)
            p.setPen(pen)
            p.drawArc(rect,
                      int(start_angle * 16),
                      int(span_angle  * 16))

        # ── 중앙 텍스트 ──────────────────────────────────────────────
        pct = round(self._value / n * 100)

        # 큰 숫자
        f_big = QFont()
        f_big.setPointSize(20)  # 중앙 퍼센트 숫자 폰트 20pt
        f_big.setBold(True)
        p.setFont(f_big)
        p.setPen(QPen(QColor("#1c2030")))
        p.drawText(QRectF(0, 0, w, h * 0.56),  # 큰 숫자: 위쪽 56% 영역에 하단 정렬
                   Qt.AlignmentFlag.AlignHCenter | Qt.AlignmentFlag.AlignBottom,
                   str(pct))

        # "%" 단위
        f_unit = QFont()
        f_unit.setPointSize(9)  # "%" 단위 폰트 9pt
        p.setFont(f_unit)
        p.setPen(QPen(QColor("#5070a0")))
        p.drawText(QRectF(0, h * 0.54, w, h * 0.16),  # "%" 텍스트: 54%~70% 높이 구간
                   Qt.AlignmentFlag.AlignCenter, "%")

        # "N / total" 작은 텍스트
        f_sub = QFont()
        f_sub.setPointSize(8)  # "N/total" 작은 텍스트 폰트 8pt
        p.setFont(f_sub)
        p.setPen(QPen(QColor("#8090b0")))
        p.drawText(QRectF(0, h * 0.70, w, h * 0.20),  # "N/total": 70%~90% 높이 구간
                   Qt.AlignmentFlag.AlignCenter,
                   f"{self._value} / {n}")

        p.end()


class StepConnector(QWidget):
    """원형 스텝 사이 수평 연결선 — 왼쪽 스텝 완료 시 초록으로 전환."""

    # circle center Y in StepCard: top margin 4 + circle radius 18 = 22
    _LINE_Y = 22

    def __init__(self, parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self.setMinimumWidth(20)
        self._done = False

    def set_done(self, done: bool) -> None:
        self._done = done
        self.update()

    def paintEvent(self, _event: Any) -> None:
        p = QPainter(self)
        p.setRenderHint(QPainter.RenderHint.Antialiasing)
        color = QColor("#28a860") if self._done else QColor("#3a4878")
        pen = QPen(color, 3)
        pen.setCapStyle(Qt.PenCapStyle.RoundCap)
        p.setPen(pen)
        y = self._LINE_Y
        p.drawLine(0, y, self.width(), y)
        p.end()


class StepCard(QWidget):
    """원형 스텝 인디케이터 카드 — 클릭으로 포함/제외 토글, 완료 시 체크 표시."""

    card_clicked = Signal(int)

    # (원 배경색, 원 테두리색, 텍스트색, 아이콘, _unused)
    _THEMES: dict[str, tuple[str, str, str, str, str]] = {
        "pending":  ("#2a3050", "#3a4878", "#6878a0", "",   ""),
        "disabled": ("#1a1e2a", "#c0392b", "#c0392b", "✕",  ""),  # 빨간 X — 제외됨
        "running":  ("#1878d0", "#1060b8", "#ffffff", "▶",  ""),
        "running2": ("#4fa0f0", "#1878d0", "#ffffff", "▶",  ""),
        "done":     ("#28a860", "#1a8048", "#ffffff", "✓",  ""),
    }

    def __init__(self, step_idx: int, seq_label: str,
                 parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self._step_idx  = step_idx
        self._status    = "pending"
        self._enabled   = True          # False = 측정에서 제외
        self._seq_label = seq_label
        # 고정 높이만 두고 너비는 열을 꽉 채워 클릭 영역을 최대화
        self.setFixedHeight(72)
        self.setMinimumWidth(64)
        self.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Fixed)
        self.setCursor(Qt.CursorShape.PointingHandCursor)
        self.setToolTip("Click to include / exclude from measurement")

        outer = QVBoxLayout(self)
        outer.setContentsMargins(0, 4, 0, 0)
        outer.setSpacing(4)
        outer.setAlignment(Qt.AlignmentFlag.AlignTop | Qt.AlignmentFlag.AlignHCenter)

        self._circle = QLabel()
        self._circle.setFixedSize(36, 36)
        self._circle.setAlignment(Qt.AlignmentFlag.AlignCenter)
        outer.addWidget(self._circle, 0, Qt.AlignmentFlag.AlignHCenter)

        self._lbl = QLabel(seq_label)
        self._lbl.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._lbl.setWordWrap(True)
        outer.addWidget(self._lbl, 0, Qt.AlignmentFlag.AlignHCenter)

        self._apply_theme()

    def _effective_status(self) -> str:
        if not self._enabled and self._status == "pending":
            return "disabled"
        return self._status

    def _apply_theme(self) -> None:
        eff = self._effective_status()
        circle_bg, circle_border, txt_color, icon, _ = self._THEMES.get(
            eff, self._THEMES["pending"])

        has_icon = bool(icon)
        font_size = "15px" if has_icon else "13px"
        self._circle.setText(icon if has_icon else str(self._step_idx + 1))
        self._circle.setStyleSheet(
            f"background:{circle_bg};"
            f"border:2px solid {circle_border};"
            f"border-radius:18px;"
            f"color:{txt_color};"
            f"font-size:{font_size};"
            f"font-weight:bold;"
        )

        is_active = eff not in ("pending", "disabled")
        if eff == "disabled":
            lbl_color = "#4a3030"       # 어두운 적갈색 — 제외 항목
        elif is_active:
            lbl_color = "#dde0f0"
        else:
            lbl_color = "#6878a0"
        self._lbl.setStyleSheet(
            f"font-size:12px;"
            f"font-weight:{'bold' if is_active else '500'};"
            f"color:{lbl_color};"
        )

    def set_status(self, status: str) -> None:
        self._status = status
        self._apply_theme()

    def set_enabled(self, enabled: bool) -> None:
        self._enabled = enabled
        self._apply_theme()

    def is_enabled(self) -> bool:
        return self._enabled

    def toggle_blink(self) -> None:
        if self._status == "running":
            self.set_status("running2")
        elif self._status == "running2":
            self.set_status("running")

    def mousePressEvent(self, event: Any) -> None:
        if self._status not in ("running", "running2"):
            self._enabled = not self._enabled
            self._apply_theme()
        self.card_clicked.emit(self._step_idx)
        super().mousePressEvent(event)


class AutoAllPanel(QWidget):
    # Emitted after each step completes: (step_idx, seq_name, key_or_mode)
    # seq_name: "lum_swing" | "lum_loading" | "gamut" | "contrast"
    # key_or_mode: "SDR_Vivid" etc. for swing/loading; "SDR"/"HDR" for gamut/contrast
    step_done:  Signal = Signal(int, str, str)
    step_begin: Signal = Signal(str, str)          # (seq_name, key_or_mode) — step 시작
    live_data:  Signal = Signal(str, str, float, object)  # (seq_name, key_or_mode, pct, data)

    def __init__(self, engine: MeasurementEngine, parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self._engine     = engine
        self._worker: Optional[AutoAllWorker] = None
        self._start_idx  = 0
        self._step_cards: list[Optional[StepCard]] = [None] * 14
        self._running           = False
        self._single_step       = False   # True = re-measure only the selected card
        self._was_single_step   = False   # set in _start, checked in _on_all_finished
        self._current_auto_seq  = ""      # seq_name of the currently running step
        self._current_auto_key  = ""      # key_or_mode of the currently running step

        # running 카드 스트립 블링크 타이머
        self._connectors: dict[tuple[int, int], StepConnector] = {}

        self._blink_timer = QTimer(self)
        self._blink_timer.setInterval(500)  # 블링크 간격 500ms (0.5초마다 깜빡임)
        self._blink_timer.timeout.connect(self._blink_running_cards)

        # 남은 시간 표기용
        self._start_time: float = 0.0
        self._steps_done_in_run: int = 0
        self._total_steps_in_run: int = 14
        self._elapsed_timer = QTimer(self)
        self._elapsed_timer.setInterval(1000)
        self._elapsed_timer.timeout.connect(self._update_eta_label)

        layout = QVBoxLayout(self)
        layout.setSpacing(6)
        layout.setContentsMargins(12, 10, 12, 8)  # 패널 여백: 좌우 12px, 상 10px, 하 8px

        # ── 1행: 타이틀 ──────────────────────────────────────────────
        title = QLabel("🤖 Auto All Measurement")
        title.setStyleSheet("font-size:14px;font-weight:bold;")
        layout.addWidget(title)

        # ── 2행: 브랜드 + LG 대기 시간 (한 줄) ───────────────────────
        row1 = QHBoxLayout(); row1.setSpacing(12)
        row1.addWidget(QLabel("Brand:"))
        self._brand_combo = QComboBox()
        self._brand_combo.addItems(["Competitor", "LG"])
        self._brand_combo.setFixedWidth(110)  # 브랜드 선택 콤보박스 너비 110px
        self._brand_combo.currentTextChanged.connect(self._on_brand_changed)
        row1.addWidget(self._brand_combo)
        row1.addSpacing(16)
        self._lg_wait_lbl = QLabel("Wait after PSM switch:")
        self._lg_wait_lbl.setVisible(False)
        row1.addWidget(self._lg_wait_lbl)
        self._lg_wait_spin = QSpinBox()
        self._lg_wait_spin.setRange(0, 30); self._lg_wait_spin.setValue(3)  # PSM 전환 대기: 기본 3초, 최대 30초
        self._lg_wait_spin.setSuffix(" s"); self._lg_wait_spin.setFixedWidth(72)  # 대기 시간 SpinBox 너비 72px
        self._lg_wait_spin.setVisible(False)
        row1.addWidget(self._lg_wait_spin)
        row1.addSpacing(12)
        self._hdr_detect_lbl = QLabel("HDR detect wait:")
        self._hdr_detect_lbl.setVisible(False)
        self._hdr_detect_lbl.setToolTip("Wait time for the TV to detect the HDR signal when switching to HDR mode")
        row1.addWidget(self._hdr_detect_lbl)
        self._hdr_detect_spin = QSpinBox()
        self._hdr_detect_spin.setRange(0, 15); self._hdr_detect_spin.setValue(3)  # HDR 감지 대기: 기본 3초, 최대 15초
        self._hdr_detect_spin.setSuffix(" s"); self._hdr_detect_spin.setFixedWidth(72)
        self._hdr_detect_spin.setVisible(False)
        self._hdr_detect_spin.setToolTip("Time to wait for the TV to detect the HDR signal before sending the luna command")
        row1.addWidget(self._hdr_detect_spin)
        row1.addStretch()
        layout.addLayout(row1)

        # ── 3행: 측정 설정 (한 줄) ───────────────────────────────────
        row2 = QHBoxLayout(); row2.setSpacing(8)
        row2.addWidget(QLabel("APL:"))
        self._version_combo = QComboBox()
        self._version_combo.addItems(["37-step", "11-step", "10-step", "2-step"])
        self._version_combo.setFixedWidth(80)  # APL 단계 선택 콤보박스 너비 80px
        row2.addWidget(self._version_combo)
        row2.addSpacing(8)
        row2.addWidget(QLabel("Count:"))
        self._meas_count = QSpinBox()
        self._meas_count.setRange(1, 10); self._meas_count.setValue(1)  # 측정 횟수: 기본 1회, 최대 10회
        self._meas_count.setSuffix(" ×"); self._meas_count.setFixedWidth(68)  # 측정 횟수 SpinBox 너비 68px
        row2.addWidget(self._meas_count)
        row2.addSpacing(8)
        self._cooling_check = QCheckBox("Cooling")
        row2.addWidget(self._cooling_check)
        row2.addWidget(QLabel("APL <"))
        self._cool_apl_spin = QSpinBox()
        self._cool_apl_spin.setRange(1, 100); self._cool_apl_spin.setValue(10)  # 쿨링 APL 기준: 기본 10% 미만
        self._cool_apl_spin.setSuffix(" %"); self._cool_apl_spin.setFixedWidth(70)
        row2.addWidget(self._cool_apl_spin)
        row2.addWidget(QLabel("at"))
        self._cool_sec_spin = QSpinBox()
        self._cool_sec_spin.setRange(1, 60); self._cool_sec_spin.setValue(5)  # 쿨링 대기: 기본 5초, 최대 60초
        self._cool_sec_spin.setSuffix(" s"); self._cool_sec_spin.setFixedWidth(64)
        row2.addWidget(self._cool_sec_spin)
        row2.addWidget(QLabel("Black"))
        row2.addSpacing(16)
        row2.addWidget(QLabel("Swing time:"))
        self._swing_total_sec = QSpinBox()
        self._swing_total_sec.setRange(1, 7200); self._swing_total_sec.setValue(301)
        self._swing_total_sec.setSuffix(" s"); self._swing_total_sec.setFixedWidth(90)
        row2.addWidget(self._swing_total_sec)
        row2.addStretch()
        layout.addLayout(row2)

        # ── 4·5행: 원형 진행 + 버튼·상태 ────────────────────────────
        progress_row = QHBoxLayout()
        progress_row.setSpacing(16)
        progress_row.setAlignment(Qt.AlignmentFlag.AlignVCenter)

        # 원형 진행률 위젯
        self._circular_progress = CircularProgress(total=14)  # 전체 측정 단계 수 = 14
        progress_row.addWidget(self._circular_progress)

        # 오른쪽: 버튼 + 상태 + 시퀀스 진행바
        right_col = QVBoxLayout()
        right_col.setSpacing(6)

        btn_row = QHBoxLayout(); btn_row.setSpacing(8)
        self._btn_start = QPushButton("▶  Start")
        self._btn_start.setObjectName("primary")
        self._btn_start.setFixedWidth(90)
        self._btn_start.clicked.connect(self._start)
        btn_row.addWidget(self._btn_start)
        self._btn_stop = QPushButton("■  Stop")
        self._btn_stop.setObjectName("danger")
        self._btn_stop.setEnabled(False)
        self._btn_stop.setFixedWidth(80)
        self._btn_stop.clicked.connect(self._stop)
        btn_row.addWidget(self._btn_stop)
        btn_row.addStretch()
        right_col.addLayout(btn_row)

        self._status_label = QLabel("Idle")
        self._status_label.setObjectName("muted")
        self._status_label.setStyleSheet("font-size:11px;")
        self._status_label.setWordWrap(True)
        right_col.addWidget(self._status_label)

        self._eta_label = QLabel("")
        self._eta_label.setStyleSheet("font-size:11px; color:#5aadff;")
        right_col.addWidget(self._eta_label)

        progress_row.addLayout(right_col, stretch=1)
        layout.addLayout(progress_row)

        # ── 6행: 순서도 그리드 (스크롤 없이) ────────────────────────
        self._hint_label = QLabel("💡 Click a card to include / exclude from measurement  —  14 / 14 selected")
        self._hint_label.setStyleSheet("font-size:10px;color:#888899;")
        layout.addWidget(self._hint_label)

        from PySide6.QtWidgets import QGridLayout
        from .worker import _ALL_STEPS as _STEPS

        grid_widget = QWidget()
        grid_widget.setStyleSheet("background:#1a1c2a;border-radius:10px;")
        grid = QGridLayout(grid_widget)
        grid.setHorizontalSpacing(0)
        grid.setVerticalSpacing(6)
        grid.setContentsMargins(10, 8, 10, 8)

        # 열 구조: col0=행라벨, col1/3/5/7=카드(Expanding), col2/4/6=연결선(stretch)
        # 카드 열: 최소 너비만 설정, 카드 위젯이 Expanding 정책으로 열을 꽉 채움
        # 연결선 열: stretch=1로 남은 공간 배분
        grid.setColumnMinimumWidth(0, 96)
        for card_col in (1, 3, 5, 7):
            grid.setColumnMinimumWidth(card_col, 64)
            grid.setColumnStretch(card_col, 2)
        for conn_col in (2, 4, 6):
            grid.setColumnMinimumWidth(conn_col, 24)
            grid.setColumnStretch(conn_col, 1)

        # 헤더 행
        grid.addWidget(QLabel(""), 0, 0)
        for sci, seq in enumerate(_SEQ_COL_ORDER):
            gcol = sci * 2 + 1   # 1, 3, 5, 7
            h = QLabel(_SEQ_COL_LABELS[seq])
            h.setAlignment(Qt.AlignmentFlag.AlignCenter)
            h.setFixedHeight(26)
            h.setStyleSheet(
                "font-size:13px;color:#aab0cc;font-weight:bold;"
                "border-bottom:1px solid #2d3145;padding-bottom:3px;"
            )
            grid.addWidget(h, 0, gcol)

        # step_idx 조회 맵: (mode_idx, seq_name) → step_idx
        step_map = {(s[0], s[3]): i for i, s in enumerate(_STEPS)}

        for ri, row_label in enumerate(_MODE_ROW_LABELS):
            grow = ri + 1 + (1 if ri >= 3 else 0)

            # SDR/HDR 경계 구분선 (전체 8열 스팬)
            if ri == 3:
                sep = QWidget()
                sep.setFixedHeight(1)
                sep.setStyleSheet("background:#2d3145;")
                grid.addWidget(sep, ri + 1, 0, 1, 8)

            lbl = QLabel(row_label)
            lbl.setAlignment(Qt.AlignmentFlag.AlignVCenter | Qt.AlignmentFlag.AlignLeft)
            lbl.setStyleSheet(
                "font-size:12px;font-weight:bold;color:#dde0f0;"
                "padding-left:6px;"
            )
            grid.addWidget(lbl, grow, 0)

            for sci, seq in enumerate(_SEQ_COL_ORDER):
                card_gcol = sci * 2 + 1   # 1, 3, 5, 7
                si = step_map.get((ri, seq))
                if si is not None:
                    card = StepCard(si, _SEQ_COL_LABELS[seq])
                    card.card_clicked.connect(self._on_card_clicked)
                    self._step_cards[si] = card
                    grid.addWidget(card, grow, card_gcol)

                # 연결선 (마지막 열 뒤에는 없음)
                if sci < len(_SEQ_COL_ORDER) - 1:
                    conn_gcol = sci * 2 + 2   # 2, 4, 6
                    conn = StepConnector()
                    self._connectors[(ri, sci)] = conn
                    grid.addWidget(conn, grow, conn_gcol)

        layout.addWidget(grid_widget, stretch=1)

    # ── 브랜드 전환 ───────────────────────────────────────────────────

    def _on_brand_changed(self, text: str) -> None:
        lg = (text == "LG")
        self._lg_wait_lbl.setVisible(lg)
        self._lg_wait_spin.setVisible(lg)
        self._hdr_detect_lbl.setVisible(lg)
        self._hdr_detect_spin.setVisible(lg)

    # ── 시작 / 중지 ───────────────────────────────────────────────────

    def _start(self) -> None:
        # LG 모드 시 장비연결 패널에서 연결된 포트 확인
        if self._brand_combo.currentText() == "LG":
            if not self._engine.lg_tv_serial or not self._engine.lg_tv_serial.is_open:
                QMessageBox.warning(self, "LG TV Not Connected",
                    "Connect the LG TV serial port in the Equipment panel first.")
                return

        version_map = {"37-step": "37", "11-step": "11", "10-step": "10", "2-step": "2"}
        settings = {
            "version":              version_map[self._version_combo.currentText()],
            "measurements_per_step": self._meas_count.value(),
            "cooling_enabled":      self._cooling_check.isChecked(),
            "cooling_apl_threshold": self._cool_apl_spin.value(),
            "cooling_duration_sec": float(self._cool_sec_spin.value()),
            "swing_sample_count":   self._swing_total_sec.value(),
            "swing_interval_sec":   1.0,
        }

        self._was_single_step = False
        from .worker import _ALL_STEPS as _STEPS_RESET

        # enabled 카드만 pending으로 리셋, 연결선도 리셋
        for i, card in enumerate(self._step_cards):
            if card:
                if card.is_enabled():
                    card.set_status("pending")
                if 0 <= i < len(_STEPS_RESET):
                    mode_idx, _, _, seq_name = _STEPS_RESET[i]
                    sci = _SEQ_COL_ORDER.index(seq_name) if seq_name in _SEQ_COL_ORDER else -1
                    conn = self._connectors.get((mode_idx, sci))
                    if conn is not None:
                        conn.set_done(False)

        skip_indices = {i for i, c in enumerate(self._step_cards) if c and not c.is_enabled()}
        enabled_count = 14 - len(skip_indices)
        if enabled_count == 0:
            QMessageBox.warning(self, "No Steps Selected",
                                "측정할 항목이 없습니다. 최소 한 개 이상 선택해주세요.")
            return

        self._circular_progress.setValue(0)
        self._circular_progress.setTotal(enabled_count)
        self._status_label.setText("Starting…")
        self._running = True

        import time as _time
        self._start_time = _time.monotonic()
        self._steps_done_in_run = 0
        self._total_steps_in_run = enabled_count
        self._eta_label.setText("⏱  Elapsed: 0:00:00")
        self._elapsed_timer.start()

        self._worker = AutoAllWorker(
            self._engine, settings,
            start_idx=0,
            skip_indices=skip_indices,
        )
        self._worker.mode_change_requested.connect(self._on_mode_change_requested)
        self._worker.step_started.connect(self._on_step_started)
        self._worker.progress.connect(self._on_progress)
        self._worker.step_completed.connect(self._on_step_completed)
        self._worker.all_finished.connect(self._on_all_finished)
        self._worker.error.connect(self._on_error)
        wire_worker_cleanup(self._worker, self, '_worker')
        self._worker.start()
        self._blink_timer.start()

        self._btn_start.setEnabled(False)
        self._btn_stop.setEnabled(True)

    def _stop(self) -> None:
        if self._worker is not None:
            self._worker.stop()
        self._running = False
        self._blink_timer.stop()
        self._elapsed_timer.stop()
        self._btn_stop.setEnabled(False)
        self._status_label.setText("Stop requested — halting after current sample")

    # ── 모드 전환 요청 핸들러 ─────────────────────────────────────────

    @Slot(bool, str)
    def _on_mode_change_requested(self, is_hdr: bool, case: str) -> None:
        mode = "HDR" if is_hdr else "SDR"

        if self._brand_combo.currentText() == "Competitor":
            msg = QMessageBox(self)
            msg.setWindowTitle("Mode Switch")
            msg.setText(f"Switch TV PSM to <b>{mode} {case}</b>, then press OK.")
            msg.setIcon(QMessageBox.Icon.Information)
            stop_btn = msg.addButton("■  Stop", QMessageBox.ButtonRole.RejectRole)
            ok_btn   = msg.addButton("✔  OK", QMessageBox.ButtonRole.AcceptRole)
            msg.setDefaultButton(ok_btn)
            msg.exec()
            if msg.clickedButton() is stop_btn:
                self._stop()
            elif self._worker is not None:
                self._worker.confirm()
        else:
            # LG 모드
            wait_ms       = self._lg_wait_spin.value()   * 1000
            hdr_detect_ms = self._hdr_detect_spin.value() * 1000 if is_hdr else 0

            def _confirm():
                if self._worker is not None:
                    self._worker.confirm()

            def _send_and_confirm():
                self._send_psm_command(is_hdr, case)
                self._status_label.setText(
                    f"{'HDR' if is_hdr else 'SDR'} {case} command sent — "
                    f"{'Waiting for PSM…' if wait_ms > 0 else 'Done'}"
                )
                if wait_ms > 0:
                    QTimer.singleShot(wait_ms, _confirm)
                else:
                    _confirm()

            if hdr_detect_ms > 0:
                self._status_label.setText(
                    f"HDR {case} — waiting for TV HDR signal… ({self._hdr_detect_spin.value()} s)"
                )
                QTimer.singleShot(hdr_detect_ms, _send_and_confirm)
            else:
                _send_and_confirm()

    def _send_psm_command(self, is_hdr: bool, case: str) -> None:
        """LG 전용 — PSM 전환 시리얼 명령 전송."""
        cmd_map: dict[tuple[bool, str], str] = {
            (False, "Vivid"):    """luna-send -n 1 -f luna://com.webos.settingsservice/setSystemSettings '{"category":"picture","settings":{"pictureMode":"vivid"}}'""",
            (False, "Standard"): """luna-send -n 1 -f luna://com.webos.settingsservice/setSystemSettings '{"category":"picture","settings":{"pictureMode":"normal"}}'""",
            (False, "Cinema"):   """luna-send -n 1 -f luna://com.webos.settingsservice/setSystemSettings '{"category":"picture","settings":{"pictureMode":"cinema"}}'""",
            (True,  "Vivid"):    """luna-send -n 1 -f luna://com.webos.settingsservice/setSystemSettings '{"category":"picture","settings":{"pictureMode":"hdrVivid"}}'""",
            (True,  "Standard"): """luna-send -n 1 -f luna://com.webos.settingsservice/setSystemSettings '{"category":"picture","settings":{"pictureMode":"hdrStandard"}}'""",
            (True,  "Cinema"):   """luna-send -n 1 -f luna://com.webos.settingsservice/setSystemSettings '{"category":"picture","settings":{"pictureMode":"hdrCinemaBright"}}'""",
        }
        cmd = cmd_map.get((is_hdr, case), "").strip()
        if not cmd:
            return
        ser = self._engine.lg_tv_serial
        if not (ser and ser.is_open):
            return

        # [TX] 로그 — 실제 전송 전에 터미널에 표시
        log_tx = getattr(self._engine, 'lg_log_tx', None)
        if log_tx:
            log_tx(cmd)

        # TV UART FIFO(64 byte) 오버플로우 방지 — 32바이트씩 나눠 전송
        def _chunked_send() -> None:
            import time
            data = (cmd + "\r\n").encode("latin-1")
            for i in range(0, len(data), 32):
                ser.write(data[i:i + 32])
                ser.flush()
                time.sleep(0.03)   # 30 ms 간격

        worker = ConnectWorker(_chunked_send)
        wire_worker_cleanup(worker, self, '_psm_send_worker')
        self._psm_send_worker = worker
        worker.start()

    # ── 진행 상황 슬롯 ────────────────────────────────────────────────

    @Slot(int, str)
    def _on_step_started(self, step_idx: int, label: str) -> None:
        self._set_row_status(step_idx, "running")
        self._status_label.setText(f"Measuring: {label}")
        # 현재 step의 seq/key 를 기록해두고 live_data 포워딩에 사용
        from .worker import _ALL_STEPS
        if 0 <= step_idx < len(_ALL_STEPS):
            _, is_hdr, case, seq_name = _ALL_STEPS[step_idx]
            mode = "HDR" if is_hdr else "SDR"
            key_or_mode = f"{mode}_{case}" if seq_name in ("lum_swing", "lum_loading") else mode
            self._current_auto_seq = seq_name
            self._current_auto_key = key_or_mode
            self.step_begin.emit(seq_name, key_or_mode)
        else:
            self._current_auto_seq = ""
            self._current_auto_key = ""

    @Slot(str, float, object)
    def _on_progress(self, _step: str, pct: float, data: object) -> None:
        if self._current_auto_seq:
            self.live_data.emit(self._current_auto_seq, self._current_auto_key, pct, data)

    @Slot(int)
    def _on_step_completed(self, step_idx: int) -> None:
        from desktop.worker import _ALL_STEPS
        self._set_row_status(step_idx, "done")
        if not self._was_single_step:
            self._circular_progress.setValue(step_idx + 1)
        self._steps_done_in_run += 1
        self._update_eta_label()
        path = _save_all_session(self._engine)
        if path:
            suffix = " (재측정)" if self._was_single_step else ""
            self._status_label.setText(f"Step {step_idx + 1}/14 done{suffix}  |  Saved: {path}")
        # Emit step_done so MainWindow can refresh the relevant panel live
        if 0 <= step_idx < len(_ALL_STEPS):
            mode_idx, is_hdr, case, seq_name = _ALL_STEPS[step_idx]
            mode = "HDR" if is_hdr else "SDR"
            key_or_mode = f"{mode}_{case}" if seq_name in ("lum_swing", "lum_loading") else mode
            self.step_done.emit(step_idx, seq_name, key_or_mode)
            # 해당 스텝 오른쪽 연결선을 완료 색으로 전환
            sci = _SEQ_COL_ORDER.index(seq_name) if seq_name in _SEQ_COL_ORDER else -1
            conn = self._connectors.get((mode_idx, sci))
            if conn is not None:
                conn.set_done(True)

    def _update_eta_label(self) -> None:
        import time as _time
        if not self._running or self._start_time == 0.0:
            return
        elapsed = _time.monotonic() - self._start_time

        def _fmt(secs: float) -> str:
            s = int(secs)
            return f"{s // 3600}:{(s % 3600) // 60:02d}:{s % 60:02d}"

        done  = self._steps_done_in_run
        total = self._total_steps_in_run
        remaining = total - done
        if done > 0 and remaining > 0:
            avg_per_step = elapsed / done
            eta = avg_per_step * remaining
            self._eta_label.setText(
                f"⏱  Elapsed: {_fmt(elapsed)}  |  ETA: ~{_fmt(eta)}  "
                f"({done}/{total} steps done)"
            )
        else:
            self._eta_label.setText(f"⏱  Elapsed: {_fmt(elapsed)}")

    @Slot()
    def _on_all_finished(self) -> None:
        self._running = False
        self._blink_timer.stop()
        self._elapsed_timer.stop()
        self._btn_start.setEnabled(True)
        self._btn_stop.setEnabled(False)
        self._circular_progress.setValue(self._total_steps_in_run)
        self._status_label.setText("✔  All measurements complete")
        self._eta_label.setText("")
        QMessageBox.information(self, "Complete", "Automated measurement complete.")

    @Slot(str)
    def _on_error(self, msg: str) -> None:
        self._running = False
        self._blink_timer.stop()
        self._status_label.setText(f"Error: {msg}")
        self._btn_start.setEnabled(True)
        self._btn_stop.setEnabled(False)
        QMessageBox.critical(self, "Error", msg)

    # ── 헬퍼 ─────────────────────────────────────────────────────────

    def _set_row_status(self, idx: int, status: str) -> None:
        card = self._step_cards[idx] if 0 <= idx < len(self._step_cards) else None
        if card is not None:
            card.set_status(status)

    def _blink_running_cards(self) -> None:
        for card in self._step_cards:
            if card is not None:
                card.toggle_blink()

    @Slot(int)
    def _on_card_clicked(self, _idx: int) -> None:
        if self._running:
            return
        self._update_selection_hint()

    def _update_selection_hint(self) -> None:
        enabled = sum(1 for c in self._step_cards if c and c.is_enabled())
        self._hint_label.setText(
            f"💡 Click a card to include / exclude from measurement  —  {enabled} / 14 selected"
        )


# ---------------------------------------------------------------------------

class MainWindow(QMainWindow):
    def __init__(self) -> None:
        super().__init__()
        self.setWindowTitle("LED Panel Analyzer")
        self.setMinimumSize(1200, 1000)  # 메인 윈도우 최소 크기 1200×1000 px

        self._engine = MeasurementEngine()

        # 테마 설정 로드 (기본값 light)
        from PySide6.QtCore import QSettings
        self._settings = QSettings("LedAnalyzer", "UI")
        self._theme: str = str(self._settings.value("theme", "light"))

        central = QWidget()
        self.setCentralWidget(central)
        root = QVBoxLayout(central)
        root.setContentsMargins(0, 0, 0, 0)
        root.setSpacing(0)

        # ── 상단 헤더 바 ────────────────────────────────────────────────
        self._header_bar = QWidget()
        self._header_bar.setFixedHeight(44)
        hlay = QHBoxLayout(self._header_bar)
        hlay.setContentsMargins(20, 0, 20, 0)
        hlay.setSpacing(10)

        self._header_dot = QLabel("●")
        hlay.addWidget(self._header_dot)
        self._header_logo = QLabel("LED Panel Analyzer")
        self._header_logo.setStyleSheet("font-size:13px;font-weight:700;letter-spacing:0.06em;")
        hlay.addWidget(self._header_logo)
        hlay.addStretch()

        self._header_meter = QLabel("CA: ─")
        self._header_gen   = QLabel("VG: ─")
        for lbl in (self._header_meter, self._header_gen):
            lbl.setObjectName("header_badge")
            hlay.addWidget(lbl)

        self._theme_btn = QPushButton("☀" if self._theme == "dark" else "🌙")
        self._theme_btn.setFixedSize(32, 28)
        self._theme_btn.setToolTip("Toggle Light / Dark theme")
        self._theme_btn.clicked.connect(self._toggle_theme)
        hlay.addWidget(self._theme_btn)

        self._header_ver = QLabel("v 1.0")
        hlay.addWidget(self._header_ver)
        root.addWidget(self._header_bar)

        self._conn_panel = ConnectionPanel(self._engine)
        self._conn_panel.setMaximumHeight(210)  # 연결 패널 최대 높이 210px — 더 크면 레이아웃 낭비
        root.addWidget(self._conn_panel)

        from PySide6.QtWidgets import QFrame
        splitter = QSplitter(Qt.Orientation.Horizontal)
        splitter.setHandleWidth(1)
        splitter.setContentsMargins(0, 0, 0, 0)
        splitter.setFrameShape(QFrame.Shape.NoFrame)

        # ── 사이드바 (QTreeWidget) ─────────────────────────────────────
        self._sidebar = QTreeWidget()
        self._sidebar.setMaximumWidth(160)
        self._sidebar.setHeaderHidden(True)
        self._sidebar.setIndentation(0)
        self._sidebar.setRootIsDecorated(True)
        # sidebar style is inherited from the global _DARK_STYLE QTreeWidget rules
        splitter.addWidget(self._sidebar)

        self._stack = QStackedWidget()
        self._lum_panel = LumLoadingPanel(self._engine)
        self._gamut_panel = GamutPanel(self._engine)
        self._contrast_panel = ContrastPanel(self._engine)
        self._report_panel = ReportPanel(
            self._engine, self._gamut_panel, self._lum_panel, self._contrast_panel
        )
        self._module_panel = ModulePanel(self._engine)
        self._gamma_sub_panel  = GammaSubPanel(self._engine)
        self._color_sub_panel  = ColorSubPanel(self._engine)
        self._calman_sub_panel = _CalmanSweepPanel(self._engine)

        self._setting_panel = SettingPanel(self._engine)
        self._auto_panel   = AutoAllPanel(self._engine)
        self._swing_panel  = LumSwingPanel(self._engine)
        _rival_panels = [
            ("Auto All",     self._auto_panel),
            ("Lum. Swing",   self._swing_panel),
            ("APL Loading",  self._lum_panel),
            ("Gamut",        self._gamut_panel),
            ("Contrast",     self._contrast_panel),
            ("Report",       self._report_panel),
        ]
        _module_panels = [
            ("Module All",   self._module_panel),
            ("Gamma",        self._gamma_sub_panel),
            ("Chromaticity / CA", self._color_sub_panel),
            ("Calman Sweep", self._calman_sub_panel),
        ]

        def _make_item(label: str, idx: int,
                       parent: "QTreeWidgetItem | None" = None) -> QTreeWidgetItem:
            it = QTreeWidgetItem([label]) if parent is None \
                else QTreeWidgetItem(parent, [label])
            it.setData(0, Qt.ItemDataRole.UserRole, idx)
            return it

        stack_idx = 0

        # Setting (Center Align + LG TV Serial) — standalone top-level
        self._setting_panel.setContentsMargins(6, 4, 6, 4)
        self._stack.addWidget(self._setting_panel)
        self._sidebar.addTopLevelItem(_make_item("Setting", stack_idx))
        stack_idx += 1

        # Competitor group (parent points to Auto All)
        rival_parent = QTreeWidgetItem(["Competitor"])
        rival_parent.setData(0, Qt.ItemDataRole.UserRole, stack_idx)
        self._sidebar.addTopLevelItem(rival_parent)
        for label, panel in _rival_panels:
            panel.setContentsMargins(6, 4, 6, 4)
            self._stack.addWidget(panel)
            rival_parent.addChild(_make_item(label, stack_idx, rival_parent))
            stack_idx += 1
        rival_parent.setExpanded(True)

        # Module group
        mod_parent = QTreeWidgetItem(["Module"])
        mod_parent.setData(0, Qt.ItemDataRole.UserRole, stack_idx)
        self._sidebar.addTopLevelItem(mod_parent)
        for label, panel in _module_panels:
            panel.setContentsMargins(6, 4, 6, 4)
            self._stack.addWidget(panel)
            mod_parent.addChild(_make_item(label, stack_idx, mod_parent))
            stack_idx += 1
        mod_parent.setExpanded(True)

        splitter.addWidget(self._stack)
        splitter.setStretchFactor(0, 0)
        splitter.setStretchFactor(1, 1)
        splitter.setSizes([160, 9999])   # 사이드바 160px, 나머지 전부 스택

        root.addWidget(splitter)

        def _on_sidebar_changed(current: QTreeWidgetItem,
                                _prev: QTreeWidgetItem) -> None:
            if current is None:
                return
            idx = current.data(0, Qt.ItemDataRole.UserRole)
            if idx is not None:
                self._stack.setCurrentIndex(idx)

        self._sidebar.currentItemChanged.connect(_on_sidebar_changed)
        first = self._sidebar.topLevelItem(0)
        if first is not None:
            self._sidebar.setCurrentItem(first)

        # Wire measured gamut from module panel → Calman sweep panel
        self._module_panel.gamut_data_ready.connect(
            self._calman_sub_panel.set_gamut_data
        )

        # Wire AutoAll live panel refresh: after each step, update the relevant panel
        self._auto_panel.step_done.connect(self._on_auto_step_done)

        # Wire AutoAll live data → panel live update (step_begin → prepare, live_data → draw)
        self._auto_panel.step_begin.connect(self._on_auto_step_begin)
        self._auto_panel.live_data.connect(self._on_auto_live_data)

        # ── 상태 바 ────────────────────────────────────────────────────
        self.setStatusBar(QStatusBar())
        self.statusBar().showMessage("LED Panel Analyzer  —  Please connect devices")

        # 헤더 배지 주기적 갱신
        from PySide6.QtCore import QTimer as _QT
        self._badge_timer = _QT(self)
        self._badge_timer.timeout.connect(self._update_header_badges)
        self._badge_timer.start(2000)  # 헤더 장비 연결 상태 배지 갱신 주기 2초

        # 초기 테마 적용
        self._theme_btn.setText("☀" if self._theme == "dark" else "🌙")
        self._apply_theme()

    # ── Theme ──────────────────────────────────────────────────────────

    def _apply_theme(self) -> None:
        is_light = (self._theme == "light")
        self.setStyleSheet(_light_style() if is_light else _dark_style())

        # 헤더 바 색상
        if is_light:
            self._header_bar.setStyleSheet(
                "background:#dde1ee; border-bottom:1px solid #c8cedf;")
            self._header_dot.setStyleSheet("font-size:14px;color:#00c9b5;padding-right:2px;")
            self._header_logo.setStyleSheet(
                "font-size:13px;font-weight:700;color:#2c3350;letter-spacing:0.06em;")
            self._header_ver.setStyleSheet("font-size:9px;color:#9ba8c8;padding-left:10px;")
            self._theme_btn.setStyleSheet(
                "background:#eef1f8;border:1px solid #c8cedf;border-radius:6px;"
                "color:#4a5580;font-size:14px;")
        else:
            self._header_bar.setStyleSheet(
                "background:#0d1120; border-bottom:1px solid #1e2a48;")
            self._header_dot.setStyleSheet("font-size:14px;color:#4e8df8;padding-right:2px;")
            self._header_logo.setStyleSheet(
                "font-size:13px;font-weight:700;color:#dde5ff;letter-spacing:0.06em;")
            self._header_ver.setStyleSheet("font-size:9px;color:#2e3858;padding-left:10px;")
            self._theme_btn.setStyleSheet(
                "background:#202840;border:1px solid #3a52a0;border-radius:6px;"
                "color:#a8b8e8;font-size:14px;")

        self._update_header_badges()

    def _toggle_theme(self) -> None:
        self._theme = "dark" if self._theme == "light" else "light"
        self._settings.setValue("theme", self._theme)
        self._theme_btn.setText("☀" if self._theme == "dark" else "🌙")
        self._apply_theme()

    def _update_header_badges(self) -> None:
        m = self._engine.meter
        g = self._engine.generator
        _badge = "font-size:10px;font-weight:600;padding:3px 12px;border-radius:12px;"
        if self._theme == "light":
            _on  = _badge + "color:#ffffff;border:1px solid #00a898;background:#00c9b5;"
            _off = _badge + "color:#7888a8;border:1px solid #d0d5e8;background:#eef1f8;"
        else:
            _on  = _badge + "color:#0a1020;border:1px solid #1dd9a0;background:#1dd9a0;"
            _off = _badge + "color:#4a5878;border:1px solid #263058;background:#141828;"
        self._header_meter.setText("CA: ●" if (m and m.is_connected) else "CA: ─")
        self._header_meter.setStyleSheet(_on if (m and m.is_connected) else _off)
        self._header_gen.setText("VG: ●" if (g and g.is_connected) else "VG: ─")
        self._header_gen.setStyleSheet(_on if (g and g.is_connected) else _off)

    @Slot(str, str)
    def _on_auto_step_begin(self, seq_name: str, key_or_mode: str) -> None:
        """AutoAll step 시작 시 해당 패널을 live 수신 상태로 준비."""
        if seq_name == "lum_swing":
            self._swing_panel.begin_auto_stream(key_or_mode)
        elif seq_name == "lum_loading":
            # 해당 케이스 시리즈 초기화 → 측정 시작 전 이전 데이터 제거
            _is_hdr = key_or_mode.startswith("HDR")
            _case   = key_or_mode.split("_", 1)[1]
            _sm = (self._lum_panel._hdr_apl_series if _is_hdr
                   else self._lum_panel._sdr_apl_series)
            _s = _sm.get(_case)
            if _s is not None:
                _s.clear()
            _ay = (self._lum_panel._ax_y_hdr if _is_hdr
                   else self._lum_panel._ax_y_sdr)
            _ay.setRange(0, 100)

    @Slot(str, str, float, object)
    def _on_auto_live_data(self, seq_name: str, key_or_mode: str, pct: float, data: object) -> None:
        """AutoAll 측정 중 실시간 데이터를 해당 패널로 전달."""
        if seq_name == "lum_swing":
            self._swing_panel._on_progress("auto", pct, data)
        elif seq_name == "lum_loading":
            # APL 단계별 데이터 → LumLoadingPanel 차트에 직접 추가
            if not (isinstance(data, dict) and "apl" in data):
                return
            _apl     = int(data["apl"])
            _results = data.get("results", [])
            if not _results:
                return
            _lv  = sum(r.Lv for r in _results) / len(_results)
            _is_hdr = key_or_mode.startswith("HDR")
            _case   = key_or_mode.split("_", 1)[1]
            _sm = (self._lum_panel._hdr_apl_series if _is_hdr
                   else self._lum_panel._sdr_apl_series)
            _s = _sm.get(_case)
            if _s is None:
                return
            _s.append(float(_apl), _lv)
            _ay = (self._lum_panel._ax_y_hdr if _is_hdr
                   else self._lum_panel._ax_y_sdr)
            if _lv * 1.15 > _ay.max():
                _ay.setRange(0, _lv * 1.15)

    @Slot(int, str, str)
    def _on_auto_step_done(self, _step_idx: int, seq_name: str, key_or_mode: str) -> None:
        """Refresh the relevant measurement panel after each AutoAll step completes."""
        if seq_name == "lum_swing":
            self._swing_panel.refresh_from_engine(key_or_mode)
        elif seq_name == "lum_loading":
            self._lum_panel.refresh_from_engine(key_or_mode)
        elif seq_name == "gamut":
            self._gamut_panel.refresh_from_engine(key_or_mode)
        elif seq_name == "contrast":
            self._contrast_panel.refresh_from_engine(key_or_mode)

    def closeEvent(self, event: Any) -> None:
        self._engine.disconnect_all()
        super().closeEvent(event)
