"""Flask routes: REST API + Server-Sent Events for live measurement progress."""
from __future__ import annotations

import dataclasses
import json
import os
import queue
import threading
from typing import Any, Generator

from flask import Blueprint, Response, current_app, jsonify, request, render_template, send_file

from core.engine import MeasurementEngine
from core.export import ExcelExporter


def _save_all_session(engine: MeasurementEngine) -> str:
    """세션 데이터를 {brand}_{model}_all.xlsx 로 자동 저장. 경로 반환."""
    if not engine.auto_save_dir:
        return ""
    brand = engine.brand or "brand"
    model = engine.model_name or "model"
    path = os.path.join(engine.auto_save_dir, f"{brand}_{model}_all.xlsx")
    try:
        ExcelExporter().export_all_session(
            brand=brand, model=model,
            session_swing=engine.session_swing,
            session_loading=engine.session_loading,
            session_gamut=engine.session_gamut,
            session_contrast=engine.session_contrast,
            file_path=path,
        )
    except Exception as e:
        print(f"[auto_save] 오류: {e}")
        return ""
    return path


def _update_session(engine: MeasurementEngine, seq_name: str,
                    params: dict[str, Any], result: Any) -> None:
    """시퀀스 완료 후 engine 세션 데이터를 갱신한다."""
    mode = "HDR" if params.get("is_hdr") else "SDR"
    case = params.get("case", "")
    if seq_name == "lum_swing" and isinstance(result, dict):
        for k, v in result.items():
            engine.session_swing[f"{mode}_{k}"] = list(v)
    elif seq_name == "lum_loading" and isinstance(result, dict):
        engine.session_loading[f"{mode}_{case}"] = dict(result)
    elif seq_name == "gamut" and isinstance(result, dict):
        engine.session_gamut[mode] = dict(result)
    elif seq_name == "contrast" and isinstance(result, dict):
        engine.session_contrast[mode] = dict(result)


def _to_json_safe(obj: Any) -> Any:
    """Recursively convert dataclasses / lists / dicts to JSON-serialisable types."""
    if dataclasses.is_dataclass(obj) and not isinstance(obj, type):
        return dataclasses.asdict(obj)  # type: ignore[arg-type]
    if isinstance(obj, dict):
        return {k: _to_json_safe(v) for k, v in obj.items()}
    if isinstance(obj, (list, tuple)):
        return [_to_json_safe(v) for v in obj]
    if isinstance(obj, (int, float, str, bool)) or obj is None:
        return obj
    return str(obj)

bp = Blueprint("api", __name__, url_prefix="/api")

ui_bp = Blueprint("ui", __name__)

@ui_bp.route("/")
def index():
    return render_template("index.html")

# Thread-safe queue for SSE progress events
_progress_queue: queue.Queue[str] = queue.Queue(maxsize=500)

# Most-recent full results, keyed by sequence name (for export)
_last_results: dict[str, Any] = {}

# AutoAll 전체 자동화: 모드 전환 게이트 + 중지 이벤트
_auto_all_gate  = threading.Event()
_auto_all_stop  = threading.Event()

# 보고서 템플릿: 파일에서 불러온 모델 데이터 목록 (프로세스 생존 동안 유지)
_report_models: list[dict[str, Any]] = []

# AutoAll 14개 측정 단계 정의 (desktop worker.py 와 동일 순서)
_AUTO_ALL_STEPS: list[tuple[bool, str, str]] = [
    (False, "Vivid",    "lum_swing"),
    (False, "Vivid",    "lum_loading"),
    (False, "Vivid",    "gamut"),
    (False, "Vivid",    "contrast"),
    (False, "Standard", "lum_swing"),
    (False, "Standard", "lum_loading"),
    (False, "Cinema",   "lum_swing"),
    (False, "Cinema",   "lum_loading"),
    (True,  "Vivid",    "lum_swing"),
    (True,  "Vivid",    "lum_loading"),
    (True,  "Standard", "lum_swing"),
    (True,  "Standard", "lum_loading"),
    (True,  "Cinema",   "lum_swing"),
    (True,  "Cinema",   "lum_loading"),
]


def _get_engine() -> MeasurementEngine:
    return current_app.extensions["engine"]


def _sse_event(data: dict[str, Any]) -> str:
    return f"data: {json.dumps(data)}\n\n"


# ---------------------------------------------------------------------------
# Connection
# ---------------------------------------------------------------------------

@bp.route("/connect", methods=["POST"])
def connect():
    """Connect meter and generator.

    JSON body:
        {
            "brand": "...", "model_name": "...",
            "meter_port": "COM3", "meter_model": "CA-410",
            "gen_port": "COM4", "gen_model": "VG-876"
        }
    """
    body = request.get_json(force=True)
    engine = _get_engine()

    engine.brand = body.get("brand", engine.brand)
    engine.model_name = body.get("model_name", engine.model_name)

    errors: list[str] = []

    meter_port = body.get("meter_port")
    if meter_port:
        try:
            engine.connect_meter(meter_port, body.get("meter_model", "CA-410"))
        except Exception as exc:
            errors.append(f"Meter: {exc}")

    gen_port = body.get("gen_port")
    if gen_port:
        try:
            engine.connect_generator(gen_port, body.get("gen_model", "VG-876"))
        except Exception as exc:
            errors.append(f"Generator: {exc}")

    if errors:
        return jsonify({"ok": False, "errors": errors}), 500

    return jsonify({"ok": True, "ready": engine.is_ready})


# ---------------------------------------------------------------------------
# Mock connect / Disconnect
# ---------------------------------------------------------------------------

@bp.route("/mock_connect", methods=["POST"])
def mock_connect():
    """Connect mock meter and generator for offline testing."""
    engine = _get_engine()
    from core.equipment.mock import MockMeter, MockGenerator
    engine.meter = MockMeter()
    engine.generator = MockGenerator()
    return jsonify({"ok": True, "ready": engine.is_ready})


@bp.route("/disconnect", methods=["POST"])
def disconnect_all():
    """Disconnect all equipment."""
    engine = _get_engine()
    engine.disconnect_all()
    engine.meter = None
    engine.generator = None
    return jsonify({"ok": True})


# ---------------------------------------------------------------------------
# Port scan
# ---------------------------------------------------------------------------

@bp.route("/ports", methods=["GET"])
def list_ports():
    """Return a sorted list of available serial port device names."""
    try:
        import serial.tools.list_ports
        ports = sorted(p.device for p in serial.tools.list_ports.comports())
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc), "ports": []}), 500
    return jsonify({"ok": True, "ports": ports})


# ---------------------------------------------------------------------------
# Status
# ---------------------------------------------------------------------------

@bp.route("/status", methods=["GET"])
def status():
    engine = _get_engine()
    return jsonify({
        "ready": engine.is_ready,
        "brand": engine.brand,
        "model_name": engine.model_name,
        "meter_connected": engine.meter is not None and engine.meter.is_connected,
        "generator_connected": engine.generator is not None and engine.generator.is_connected,
    })


@bp.route("/hdr", methods=["POST"])
def set_hdr():
    """Switch the connected generator between SDR and HDR immediately."""
    engine = _get_engine()
    gen = engine.generator
    if gen is None or not gen.is_connected:
        return jsonify({"ok": False, "error": "Generator is not connected"}), 400

    body = request.get_json(force=True) or {}
    enabled = bool(body.get("enabled", False))
    try:
        gen.set_hdr(enabled)
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 500

    return jsonify({"ok": True, "hdr": enabled})


# ---------------------------------------------------------------------------
# Run sequence
# ---------------------------------------------------------------------------

@bp.route("/run/<seq_name>", methods=["POST"])
def run_sequence(seq_name: str):
    """Start a named measurement sequence in a background thread.

    Sequence parameters are passed as JSON body.  Progress is streamed via
    GET /api/progress (SSE).
    """
    engine = _get_engine()
    if not engine.is_ready:
        return jsonify({"ok": False, "error": "Engine not ready"}), 400

    # 이전 측정의 잔류 이벤트 제거 — 새 SSE 연결이 구 데이터를 받지 않도록
    while not _progress_queue.empty():
        try:
            _progress_queue.get_nowait()
        except queue.Empty:
            break

    params: dict[str, Any] = request.get_json(force=True) or {}

    def _progress(step: str, pct: float, data: Any) -> None:
        try:
            serialisable_data: Any = _to_json_safe(data)
        except Exception:
            serialisable_data = None

        event = _sse_event({"step": step, "progress": pct, "data": serialisable_data})
        try:
            _progress_queue.put_nowait(event)
        except queue.Full:
            pass  # drop if consumer is too slow

    engine.on_progress = _progress

    def _run() -> None:
        try:
            result = engine.run_sequence(seq_name, **params)
            _last_results[seq_name] = result
            _update_session(engine, seq_name, params, result)
            save_path = _save_all_session(engine)
            _progress_queue.put_nowait(_sse_event({
                "step": seq_name, "progress": 1.0, "done": True,
                "auto_save_path": save_path,
            }))
        except Exception as exc:
            _progress_queue.put_nowait(_sse_event({"error": str(exc)}))

    thread = threading.Thread(target=_run, daemon=True)
    thread.start()

    return jsonify({"ok": True, "seq_name": seq_name})


# ---------------------------------------------------------------------------
# Auto-save directory
# ---------------------------------------------------------------------------

@bp.route("/auto_save_dir", methods=["GET"])
def get_auto_save_dir():
    engine = _get_engine()
    return jsonify({"ok": True, "path": engine.auto_save_dir})


@bp.route("/auto_save_dir", methods=["POST"])
def set_auto_save_dir():
    body = request.get_json(force=True) or {}
    path = body.get("path", "").strip()
    if path:
        try:
            os.makedirs(path, exist_ok=True)
        except Exception as exc:
            return jsonify({"ok": False, "error": str(exc)}), 400
    _get_engine().auto_save_dir = path
    return jsonify({"ok": True, "path": path})


# ---------------------------------------------------------------------------
# Stop
# ---------------------------------------------------------------------------

@bp.route("/stop", methods=["POST"])
def stop_sequence():
    """Signal the currently running sequence to stop."""
    _auto_all_stop.set()   # AutoAll 도 함께 정지
    _auto_all_gate.set()   # 대기 중인 게이트 해제
    _get_engine().stop_sequence()
    return jsonify({"ok": True})


# ---------------------------------------------------------------------------
# AutoAll — 전체 자동화 측정
# ---------------------------------------------------------------------------

@bp.route("/auto_all/run", methods=["POST"])
def run_auto_all():
    """14개 단계 전체 자동 측정. 모드 전환 시 SSE로 mode_change 이벤트를 보내고
    /api/auto_all/confirm 이 호출될 때까지 대기한다."""
    engine = _get_engine()
    if not engine.is_ready:
        return jsonify({"ok": False, "error": "Engine not ready"}), 400

    while not _progress_queue.empty():
        try:
            _progress_queue.get_nowait()
        except queue.Empty:
            break

    params: dict[str, Any] = request.get_json(force=True) or {}
    version               = str(params.get("version", "37"))
    cooling_enabled       = bool(params.get("cooling_enabled", False))
    cooling_apl_threshold = int(params.get("cooling_apl_threshold", 10))
    cooling_duration_sec  = float(params.get("cooling_duration_sec", 5))
    measurements_per_step = int(params.get("measurements_per_step", 1))

    _auto_all_stop.clear()
    _auto_all_gate.clear()
    total = len(_AUTO_ALL_STEPS)

    def _run() -> None:
        prev_mode_key: tuple | None = None
        for step_idx, (is_hdr, case, seq_name) in enumerate(_AUTO_ALL_STEPS):
            if _auto_all_stop.is_set():
                break
            mode    = "HDR" if is_hdr else "SDR"
            key     = f"{mode}_{case}"
            mk      = (is_hdr, case)

            # 모드가 바뀔 때: ① 신호 먼저 전환 → ② UI에 PSM 변경 요청 → ③ 사용자 확인 대기
            if prev_mode_key is None or mk != prev_mode_key:
                # ① 패턴 제너레이터 신호를 먼저 전환 (idempotent — 이미 같은 상태면 무시)
                #    HDR 신호가 TV에 도달한 뒤 사용자가 PSM 모드를 바꿀 수 있도록
                if engine.generator is not None:
                    if is_hdr:
                        engine.generator.set_hdr(True)
                    else:
                        engine.generator.set_sdr()

                # ② UI에 모드 전환 요청 이벤트 발송
                _auto_all_gate.clear()
                _progress_queue.put_nowait(_sse_event({
                    "mode_change": True,
                    "is_hdr": is_hdr, "case": case,
                    "step_idx": step_idx, "total": total,
                }))
                # ③ 사용자가 "전환 완료" 클릭 시까지 대기 (최대 10분)
                _auto_all_gate.wait(timeout=600)
                if _auto_all_stop.is_set():
                    break
                prev_mode_key = mk

            _SEQ_KR = {
                "lum_swing": "휘도 스윙", "lum_loading": "APL 로딩",
                "gamut": "색재현율",      "contrast": "명암비",
            }
            label = f"{mode} {case} — {_SEQ_KR.get(seq_name, seq_name)}"
            _progress_queue.put_nowait(_sse_event({
                "step": f"auto_{step_idx}",
                "progress": step_idx / total,
                "data": {"step_idx": step_idx, "label": label, "total": total},
            }))

            def _progress(s: str, p: float, d: Any, si: int = step_idx) -> None:
                try:
                    _progress_queue.put_nowait(_sse_event({
                        "step": f"auto_{si}/{s}",
                        "progress": (si + p) / total,
                        "data": _to_json_safe(d),
                    }))
                except queue.Full:
                    pass

            engine.on_progress = _progress

            try:
                if seq_name == "lum_swing":
                    result = engine.run_sequence("lum_swing", case=case, is_hdr=is_hdr)
                    engine.session_swing[key] = list(result.get(case, []))

                elif seq_name == "lum_loading":
                    result = engine.run_sequence(
                        "lum_loading",
                        version=version, case=case, is_hdr=is_hdr,
                        cooling_enabled=cooling_enabled,
                        cooling_apl_threshold=cooling_apl_threshold,
                        cooling_duration_sec=cooling_duration_sec,
                        measurements_per_step=measurements_per_step,
                    )
                    engine.session_loading[key] = dict(result)

                elif seq_name == "gamut":
                    result = engine.run_sequence("gamut", is_hdr=is_hdr)
                    engine.session_gamut[mode] = dict(result)

                elif seq_name == "contrast":
                    result = engine.run_sequence("contrast", is_hdr=is_hdr)
                    engine.session_contrast[mode] = dict(result)

            except Exception as exc:
                _progress_queue.put_nowait(_sse_event({
                    "step": f"auto_{step_idx}",
                    "progress": step_idx / total,
                    "data": {"step_idx": step_idx, "error": str(exc)},
                }))
                continue

            if _auto_all_stop.is_set():
                break

            _progress_queue.put_nowait(_sse_event({
                "step": f"auto_{step_idx}",
                "progress": (step_idx + 1) / total,
                "data": {"step_idx": step_idx, "completed": True, "label": label},
            }))

        save_path = _save_all_session(engine)
        _progress_queue.put_nowait(_sse_event({
            "step": "auto_all", "progress": 1.0, "done": True,
            "auto_save_path": save_path,
        }))

    threading.Thread(target=_run, daemon=True).start()
    return jsonify({"ok": True, "seq_name": "auto_all"})


@bp.route("/auto_all/confirm", methods=["POST"])
def auto_all_confirm():
    """모드 전환 완료 — UI가 호출해서 AutoAll 게이트를 해제한다."""
    _auto_all_gate.set()
    return jsonify({"ok": True})


# ---------------------------------------------------------------------------
# Report template — 다중 xlsx 파일 파싱 + 모델 비교
# ---------------------------------------------------------------------------

def _parse_report_xlsx(path: str, brand_hint: str = "", model_hint: str = "") -> dict[str, Any]:
    """all-session xlsx 파일에서 모델 비교 데이터를 추출한다.

    export_all_session() 의 시트 구조:
      Info          : Brand / Model / Sequence
      Loading_Summary : APL(%) | {case} Avg | {case} Max | {case} Min | ...
      Loading_{case}  : APL(%) | #(idx) | Time | Lv | x | y | u' | v' | ...
                         → Lv 는 row[3] (0-indexed), row[1] = # (측정 순서)
      Gamut_SDR / Gamut : Color | Time | Lv | x | y | u'(row[5]) | v'(row[6]) | ...
      Contrast_SDR / Contrast_HDR : Black H/V(%) | Lv | CR | ...
                         → row[0] = "Full White" 또는 float(100.0/50.0/…), row[1] = Lv
    """
    import openpyxl as _opx
    import statistics as _stat

    wb = _opx.load_workbook(path, read_only=True, data_only=True)
    try:
        sheet_names = wb.sheetnames
        brand = brand_hint
        model = model_hint

        if "Info" in sheet_names:
            ws = wb["Info"]
            for row in ws.iter_rows(min_row=1, max_row=10, values_only=True):
                if not row or row[0] is None:
                    continue
                k = str(row[0]).strip()
                v = str(row[1] or "").strip() if len(row) > 1 else ""
                if k == "Brand":   brand = brand or v
                elif k == "Model": model = model or v

        if not brand or not model:
            base = os.path.splitext(os.path.basename(path))[0]
            parts = base.split("_")
            brand = brand or (parts[0] if parts else "Unknown")
            model = model or (base if len(parts) < 2 else "_".join(parts[1:]))

        entry: dict[str, Any] = {
            "brand": brand, "model": model,
            "hdr_10": None, "hdr_100": None,
            "sdr_10": None, "sdr_100": None,
            "contrast_ratio": None, "black_lv": None,
            "dci_overlap": None, "bt2020_overlap": None,
            "sdr_vivid_10":    None, "sdr_vivid_100":    None,
            "sdr_standard_10": None, "sdr_standard_100": None,
            "hdr_vivid_10":    None, "hdr_vivid_100":    None,
            "hdr_standard_10": None, "hdr_standard_100": None,
            "hdr_cinema_10":   None, "hdr_cinema_100":   None,
            "apl_sdr": {}, "apl_hdr": {},
            "gamut_uv": {},
        }

        # ── ① Loading_Summary 시트 우선 파싱 ────────────────────────────────
        # 헤더: ["APL (%)", "SDR_Vivid Avg", "SDR_Vivid Max", "SDR_Vivid Min", ...]
        if "Loading_Summary" in sheet_names:
            ws_sum = wb["Loading_Summary"]
            all_rows = list(ws_sum.iter_rows(min_row=1, values_only=True))
            if len(all_rows) >= 2:
                hdr_row = all_rows[0]
                # case_key → col index (Avg 컬럼만)
                case_col: dict[str, int] = {}
                for ci, h in enumerate(hdr_row):
                    if h and str(h).endswith(" Avg"):
                        case_col[str(h)[:-4]] = ci  # "SDR_Vivid" → col index

                apl_data: dict[str, dict[int, float]] = {ck: {} for ck in case_col}
                for row in all_rows[1:]:
                    if not row or row[0] is None:
                        continue
                    try:
                        apl_i = int(float(str(row[0])))
                    except (TypeError, ValueError):
                        continue
                    for ck, ci in case_col.items():
                        if ci < len(row) and row[ci] is not None:
                            try:
                                apl_data[ck][apl_i] = float(str(row[ci]))
                            except (TypeError, ValueError):
                                pass

                for case_key, apl_map in apl_data.items():
                    # case_key: "SDR_Vivid", "HDR_Cinema" 등
                    parts_ck = case_key.split("_", 1)
                    if len(parts_ck) < 2:
                        continue
                    mode_s = parts_ck[0].upper()    # "SDR" / "HDR"
                    case_s = parts_ck[1]             # "Vivid" / "Standard" / "Cinema"
                    k_pfx  = f"{mode_s.lower()}_{case_s.lower()}"  # "sdr_vivid"

                    chart_key = "apl_hdr" if mode_s == "HDR" else "apl_sdr"
                    entry[chart_key][case_s] = [
                        {"apl": a, "Lv": lv} for a, lv in sorted(apl_map.items())
                    ]
                    for sfx, apl_val in (("_10", 10), ("_100", 100)):
                        ek = f"{k_pfx}{sfx}"
                        if ek in entry and apl_val in apl_map:
                            entry[ek] = apl_map[apl_val]
                    if case_s == "Vivid":
                        if mode_s == "HDR":
                            if 10  in apl_map: entry["hdr_10"]  = apl_map[10]
                            if 100 in apl_map: entry["hdr_100"] = apl_map[100]
                        else:
                            if 10  in apl_map: entry["sdr_10"]  = apl_map[10]
                            if 100 in apl_map: entry["sdr_100"] = apl_map[100]

        # ── ② Loading_Summary 없을 때 개별 Loading_{case} 시트 파싱 ────────
        # 시트 구조: APL(%) | #(idx) | Time | Lv | x | y | u' | v' | ...
        #            row[0]  row[1]   row[2]  row[3] ← Lv 위치
        else:
            for sn in sheet_names:
                if not sn.startswith("Loading_"):
                    continue
                parts_sn = sn.split("_", 2)  # ["Loading", "SDR", "Vivid"]
                if len(parts_sn) < 3:
                    continue
                mode_str = parts_sn[1].upper()
                case_str = parts_sn[2]
                k_pfx    = f"{mode_str.lower()}_{case_str.lower()}"

                ws2 = wb[sn]
                # APL별로 여러 측정이 있을 수 있으므로 리스트로 수집 후 평균
                apl_lists: dict[int, list[float]] = {}
                for row in ws2.iter_rows(min_row=2, values_only=True):
                    try:
                        apl_v = row[0]
                        lv_v  = row[3]   # ← 컬럼 4 (0-indexed 3): Lv
                        if apl_v is None or lv_v is None:
                            continue
                        apl_i = int(float(str(apl_v)))
                        lv_f  = float(str(lv_v))
                        apl_lists.setdefault(apl_i, []).append(lv_f)
                    except (TypeError, ValueError, IndexError):
                        continue

                apl_map2 = {
                    a: round(_stat.mean(lvs), 3) for a, lvs in apl_lists.items()
                }
                chart_key2 = "apl_hdr" if mode_str == "HDR" else "apl_sdr"
                entry[chart_key2][case_str] = [
                    {"apl": a, "Lv": lv} for a, lv in sorted(apl_map2.items())
                ]
                for sfx, apl_val in (("_10", 10), ("_100", 100)):
                    ek = f"{k_pfx}{sfx}"
                    if ek in entry and apl_val in apl_map2:
                        entry[ek] = apl_map2[apl_val]
                if case_str == "Vivid":
                    if mode_str == "HDR":
                        if 10  in apl_map2: entry["hdr_10"]  = apl_map2[10]
                        if 100 in apl_map2: entry["hdr_100"] = apl_map2[100]
                    else:
                        if 10  in apl_map2: entry["sdr_10"]  = apl_map2[10]
                        if 100 in apl_map2: entry["sdr_100"] = apl_map2[100]

        # ── ③ Gamut 시트 파싱 ────────────────────────────────────────────────
        # 시트 구조: Color | Time | Lv | x | y | u'(row[5]) | v'(row[6]) | ...
        for gn in ("Gamut_SDR", "Gamut_HDR", "Gamut"):
            if gn not in sheet_names:
                continue
            gws = wb[gn]
            r_uv = g_uv = b_uv = None
            for row in gws.iter_rows(min_row=2, values_only=True):
                if not row or row[0] is None:
                    continue
                color = str(row[0]).strip().lower()
                try:
                    u = float(str(row[5])) if len(row) > 5 and row[5] is not None else None
                    v = float(str(row[6])) if len(row) > 6 and row[6] is not None else None
                    if u is None or v is None:
                        continue
                    entry["gamut_uv"][color] = (u, v)
                    if color in ("red",   "r"): r_uv = (u, v)
                    elif color in ("green","g"): g_uv = (u, v)
                    elif color in ("blue", "b"): b_uv = (u, v)
                except (TypeError, ValueError):
                    continue
            if r_uv and g_uv and b_uv:
                try:
                    from core.gamut_utils import calc_gamut_stats
                    st = calc_gamut_stats(r_uv, g_uv, b_uv)
                    entry["dci_overlap"]    = st["dci_overlap"]
                    entry["bt2020_overlap"] = st["bt2020_overlap"]
                except Exception:
                    pass
            break  # SDR Gamut 우선

        # ── ④ Contrast 시트 파싱 ─────────────────────────────────────────────
        # 시트 구조: Black H/V(%) | Lv(cd/m²) | CR(White/Lv) | ...
        #   row[0] = "Full White" 또는 float(100.0 / 50.0 / …)
        #   row[1] = Lv  ← 올바른 컬럼
        for cn in ("Contrast_SDR", "Contrast_HDR", "Contrast"):
            if cn not in sheet_names:
                continue
            cws = wb[cn]
            white_lv = black_lv = None
            for row in cws.iter_rows(min_row=2, values_only=True):
                if not row or row[0] is None or row[1] is None:
                    continue
                try:
                    lv_val = float(str(row[1]))
                except (TypeError, ValueError):
                    continue
                label = str(row[0]).strip().lower()
                if "full" in label or "white" in label:
                    # Full White 행
                    white_lv = lv_val
                else:
                    # 숫자 레이블 (window size %): 100 이 Full Black 윈도우
                    try:
                        side_pct = float(str(row[0]))
                        if side_pct >= 100.0:
                            black_lv = lv_val
                    except (TypeError, ValueError):
                        pass
            if white_lv and black_lv and black_lv > 0:
                entry["contrast_ratio"] = round(white_lv / black_lv, 1)
                entry["black_lv"]       = black_lv
            break

        return entry
    finally:
        wb.close()


@bp.route("/report/load", methods=["POST"])
def report_load():
    """업로드된 xlsx 파일에서 모델 데이터를 파싱해 _report_models 에 추가."""
    from werkzeug.utils import secure_filename
    import tempfile

    files = request.files.getlist("files")
    if not files:
        return jsonify({"ok": False, "error": "No files"}), 400

    errors: list[str] = []
    added: list[str] = []
    for f in files:
        fname = secure_filename(f.filename or "file.xlsx")
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            f.save(tmp.name)
            tmp_path = tmp.name
        try:
            entry = _parse_report_xlsx(tmp_path)
            key = f"{entry['brand']}_{entry['model']}"
            # 이미 있으면 덮어씌움
            for i, m in enumerate(_report_models):
                if f"{m['brand']}_{m['model']}" == key:
                    _report_models[i] = entry
                    added.append(key)
                    break
            else:
                _report_models.append(entry)
                added.append(key)
        except Exception as exc:
            errors.append(f"{fname}: {exc}")
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    return jsonify({"ok": True, "added": added, "errors": errors,
                    "models": [{"brand": m["brand"], "model": m["model"]} for m in _report_models]})


@bp.route("/report/models", methods=["GET"])
def report_models():
    return jsonify({"ok": True,
                    "models": [{"brand": m["brand"], "model": m["model"]} for m in _report_models]})


@bp.route("/report/clear", methods=["POST"])
def report_clear():
    _report_models.clear()
    return jsonify({"ok": True})


@bp.route("/report/delete", methods=["POST"])
def report_delete():
    body = request.get_json(force=True) or {}
    key = body.get("key", "")
    for i, m in enumerate(_report_models):
        if f"{m['brand']}_{m['model']}" == key:
            _report_models.pop(i)
            return jsonify({"ok": True})
    return jsonify({"ok": False, "error": "Not found"}), 404


@bp.route("/report/table", methods=["GET"])
def report_table():
    """현재 모델 목록을 비교 테이블 데이터로 반환."""
    fmt = request.args.get("fmt", "competitor")  # competitor | optical

    if fmt == "competitor":
        ROW_DEFS = [
            ("White 휘도[nit]", "HDR 10%",        "hdr_10"),
            ("White 휘도[nit]", "HDR 100%",       "hdr_100"),
            ("White 휘도[nit]", "SDR 10%",        "sdr_10"),
            ("White 휘도[nit]", "SDR 100%",       "sdr_100"),
            ("White 휘도[nit]", "Contrast Ratio", "contrast_ratio"),
            ("White 휘도[nit]", "Black",           "black_lv"),
            ("Color Gamut[%]",  "DCI-P3 (%)",     "dci_overlap"),
            ("Color Gamut[%]",  "BT.2020 (%)",    "bt2020_overlap"),
        ]
    else:
        ROW_DEFS = [
            ("휘도", "Vivid SDR 10%",        "sdr_vivid_10"),
            ("휘도", "Vivid SDR 100%",       "sdr_vivid_100"),
            ("휘도", "Standard SDR 10%",     "sdr_standard_10"),
            ("휘도", "Standard SDR 100%",    "sdr_standard_100"),
            ("휘도", "Vivid HDR 10%",        "hdr_vivid_10"),
            ("휘도", "Vivid HDR 100%",       "hdr_vivid_100"),
            ("휘도", "Standard HDR 10%",     "hdr_standard_10"),
            ("휘도", "Standard HDR 100%",    "hdr_standard_100"),
            ("휘도", "Cinema HDR 10%",       "hdr_cinema_10"),
            ("휘도", "Cinema HDR 100%",      "hdr_cinema_100"),
            ("Contrast", "Black (Ratio)",    "contrast_ratio"),
            ("Color Gamut", "DCI-P3 (%)",   "dci_overlap"),
            ("Color Gamut", "BT.2020 (%)",  "bt2020_overlap"),
        ]

    def _fmt_val(v: Any, key: str) -> str:
        if v is None:
            return "—"
        INT_KEYS = {
            "hdr_10","hdr_100","sdr_10","sdr_100",
            "sdr_vivid_10","sdr_vivid_100","sdr_standard_10","sdr_standard_100",
            "hdr_vivid_10","hdr_vivid_100","hdr_standard_10","hdr_standard_100",
            "hdr_cinema_10","hdr_cinema_100",
        }
        if key in INT_KEYS:
            return str(round(float(v)))
        if isinstance(v, float):
            return f"{v:.2f}"
        return str(v)

    rows = []
    for section, item, data_key in ROW_DEFS:
        cells = [section, item]
        for m in _report_models:
            cells.append(_fmt_val(m.get(data_key), data_key))
        rows.append(cells)

    headers = ["구분", "항목"] + [f"{m['brand']}_{m['model']}" for m in _report_models]
    return jsonify({"ok": True, "headers": headers, "rows": rows})


# ---------------------------------------------------------------------------
# SSE progress stream
# ---------------------------------------------------------------------------

@bp.route("/progress", methods=["GET"])
def progress_stream():
    """Server-Sent Events endpoint.  Connect once; receive events until done."""

    def _generate() -> Generator[str, None, None]:
        yield _sse_event({"connected": True})
        while True:
            try:
                event = _progress_queue.get(timeout=30)
                yield event
                if '"done": true' in event or '"error":' in event:
                    break
            except queue.Empty:
                # Keep-alive comment to prevent proxy timeouts
                yield ": keep-alive\n\n"

    return Response(
        _generate(),
        mimetype="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


# ---------------------------------------------------------------------------
# Export
# ---------------------------------------------------------------------------

@bp.route("/export/ppt", methods=["GET"])
def export_ppt():
    """현재 세션 데이터로 PPT 보고서 생성 후 다운로드."""
    import io
    import os
    import statistics
    import tempfile

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE
        from pptx.enum.text import PP_ALIGN
        from pptx.oxml.ns import qn as _qn
        from pptx.util import Inches, Pt
        from lxml import etree as _et
    except ImportError:
        return jsonify({"error": "pip install python-pptx lxml"}), 500

    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.font_manager as _fm_web

        # 한글 폰트 설정 (깨짐 방지)
        _KR = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
        if os.path.exists(_KR):
            _fm_web.fontManager.addfont(_KR)
            matplotlib.rcParams["font.family"] = "Noto Sans CJK JP"
        matplotlib.rcParams["axes.unicode_minus"] = False
    except ImportError:
        return jsonify({"error": "pip install matplotlib"}), 500

    engine = _get_engine()
    brand  = engine.brand      or "brand"
    model  = engine.model_name or "model"

    # ── 세션 데이터 → 표시값 ────────────────────────────────────────────
    def _lv(session_key: str, apl: int) -> float | None:
        data = engine.session_loading.get(session_key, {}).get(apl)
        if not data:
            return None
        return max(r.Lv for r in data)

    def _fmt(v: Any) -> str:
        if v is None:
            return "—"
        if isinstance(v, float) and v < 1.0:
            return f"{v:.3f}"
        return f"{round(v)}"

    ROW_LABELS = [
        ("White 휘도[nit]", "HDR 10%"),
        ("White 휘도[nit]", "HDR 100%"),
        ("White 휘도[nit]", "SDR 10%"),
        ("White 휘도[nit]", "SDR 100%"),
        ("White 휘도[nit]", "Contrast Ratio"),
        ("White 휘도[nit]", "Black"),
        ("Color Gamut[%]",  "DCI-P3 (%)"),
        ("Color Gamut[%]",  "BT.2020 (%)"),
    ]
    SPLIT_MAP = {
        "HDR 10%": ("HDR", "10%"), "HDR 100%": ("HDR", "100%"),
        "SDR 10%": ("SDR", "10%"), "SDR 100%": ("SDR", "100%"),
    }

    cr_data   = engine.session_contrast.get("SDR", {})
    white_r   = cr_data.get(0.0)
    black100  = cr_data.get(100.0)
    cr_val    = (white_r.Lv / black100.Lv) if (white_r and black100 and black100.Lv > 0) else None
    blk_val   = black100.Lv if black100 else None

    dci_overlap = bt2020_overlap = None
    gamut_data = engine.session_gamut.get("SDR", {})
    if all(k in gamut_data for k in ("red", "green", "blue")):
        try:
            from core.gamut_utils import calc_gamut_stats
            stats = calc_gamut_stats(
                (gamut_data["red"].u_prime,   gamut_data["red"].v_prime),
                (gamut_data["green"].u_prime, gamut_data["green"].v_prime),
                (gamut_data["blue"].u_prime,  gamut_data["blue"].v_prime),
            )
            dci_overlap    = stats["dci_overlap"]
            bt2020_overlap = stats["bt2020_overlap"]
        except Exception:
            pass

    row_values = [
        _fmt(_lv("HDR_Vivid", 10)),  _fmt(_lv("HDR_Vivid", 100)),
        _fmt(_lv("SDR_Vivid", 10)),  _fmt(_lv("SDR_Vivid", 100)),
        _fmt(cr_val), _fmt(blk_val),
        f"{dci_overlap:.1f}"    if dci_overlap    is not None else "—",
        f"{bt2020_overlap:.1f}" if bt2020_overlap is not None else "—",
    ]

    # ── 그래프 PNG 생성 (matplotlib) ─────────────────────────────────────
    tmp_files: list[str] = []

    def _save_tmp(fig, dpi: int = 200) -> str:
        """figure를 임시 PNG로 저장. bbox_inches 미사용 → figsize×dpi 픽셀 고정."""
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        tmp.close()
        fig.savefig(tmp.name, dpi=dpi)  # bbox_inches 없음 — 정확한 픽셀 크기 보장
        plt.close(fig)
        tmp_files.append(tmp.name)
        return tmp.name

    def _make_apl_chart(title: str, modes: list[str]) -> str:
        """APL 로딩 꺾은선 그래프 — 데스크탑과 동일 스타일."""
        fig, ax = plt.subplots(figsize=(6.5, 3.8))   # 6.5×3.8인치, dpi=200 → 1300×760 px
        fig.patch.set_facecolor("white")
        ax.set_facecolor("white")
        palette = ["#4f8ef7", "#e74c3c", "#27ae60", "#f39c12", "#9b59b6",
                   "#1abc9c", "#e67e22", "#34495e"]   # 데스크탑 _DEFAULT_MODEL_COLORS 와 동일
        for i, key in enumerate(modes):
            data = engine.session_loading.get(key, {})
            if not data:
                continue
            apls = sorted(data.keys())
            lvs  = [max(r.Lv for r in data[a]) for a in apls]
            ax.plot(apls, lvs, marker="", linewidth=1.5,   # marker="" — 점 표시 제거
                    color=palette[i % len(palette)],
                    label=key.replace("_", " "))
        ax.set_xlabel("APL (%)", fontsize=10)
        ax.set_ylabel("Lv (cd/m²)", fontsize=10)
        ax.set_title(title, fontsize=11, fontweight="bold")
        ax.set_ylim(bottom=0)                         # Y축 0부터 표기
        ax.legend(fontsize=10, loc="upper right")     # 범례: 오른쪽 상단
        ax.tick_params(labelsize=10)
        ax.grid(True, linestyle="--", alpha=0.4)      # 점선 그리드, 불투명도 40%
        fig.tight_layout()
        return _save_tmp(fig, dpi=200)                # 200 DPI 고해상도

    def _make_gamut_chart() -> tuple[str, dict]:
        """u'v' 색도 차트 — 데스크탑과 동일 스타일.

        반환: (이미지 경로, BT.2020 꼭짓점 figure fraction dict)
          vfracs = {"red":(xf,yf), "green":(xf,yf), "blue":(xf,yf)}  (y: 위=0)
        """
        import matplotlib.lines as mlines
        from core.gamut_utils import DCI_P3_UV, BT2020_UV

        _DPI = 200                          # 저장 DPI — 900×900 px (figsize 4.5×4.5인치)
        _FW, _FH = 4.5, 4.5                 # 정사각형 figure (인치)
        fig, ax = plt.subplots(figsize=(_FW, _FH))
        fig.patch.set_facecolor("white")
        ax.set_facecolor("white")
        ax.set_xlim(0.0, 0.65)
        ax.set_ylim(0.0, 0.65)
        ax.set_aspect("equal")
        # bbox_inches 미사용이므로 subplots_adjust로 하단 범례 공간 확보
        fig.subplots_adjust(left=0.13, right=0.95, top=0.91, bottom=0.22)

        # 기준 색역 삼각형 — 보고서 Qt 차트와 동일: 둘 다 #aab0c0 회색
        #   DCI-P3: 점선("--"),  BT.2020: 실선("-")
        def _rt(pts, ls):
            xs = [p[0] for p in pts] + [pts[0][0]]
            ys = [p[1] for p in pts] + [pts[0][1]]
            ax.plot(xs, ys, color="#aab0c0", linestyle=ls, linewidth=1.2)

        _rt(DCI_P3_UV, "--")
        _rt(BT2020_UV, "-")

        # 측정 모델 삼각형 — 꼭짓점 점 표시 제거, 삼각형 라인만
        mode_palette = ["#4f8ef7", "#e74c3c", "#27ae60", "#f39c12", "#9b59b6",
                        "#1abc9c", "#e67e22", "#34495e"]   # 데스크탑 _DEFAULT_MODEL_COLORS
        model_handles = []
        for mi, (mode, gdata) in enumerate(engine.session_gamut.items()):
            col = mode_palette[mi % len(mode_palette)]
            tri = [(gdata[c].u_prime, gdata[c].v_prime)
                   for c in ("red", "green", "blue") if c in gdata]
            if len(tri) == 3:
                xs = [p[0] for p in tri] + [tri[0][0]]
                ys = [p[1] for p in tri] + [tri[0][1]]
                ax.plot(xs, ys, color=col, linewidth=1.5)
                model_handles.append(
                    mlines.Line2D([0], [0], color=col, linewidth=1.5,
                                  label=mode.replace("_", " ")))
            # 꼭짓점 점 표시 제거

        # 모델 범례: 그래프 안 오른쪽 아래 강제 고정
        if model_handles:
            ax.legend(handles=model_handles,
                      loc="lower right",
                      bbox_to_anchor=(0.99, 0.01),   # axes 기준 오른쪽(0.99) 하단(0.01)
                      fontsize=10, frameon=True, framealpha=0.88)

        ax.set_xlabel("u'", fontsize=10)
        ax.set_ylabel("v'", fontsize=10)
        ax.set_title("u'v' 색도", fontsize=11, fontweight="bold")
        ax.tick_params(labelsize=10)
        ax.grid(True, linestyle="--", alpha=0.3)

        # 하단 별도 범례: DCI-P3(점선) / BT.2020(실선) — figure 기준 중앙 하단
        ref_handles = [
            mlines.Line2D([0], [0], color="#aab0c0", linestyle="--",
                          linewidth=1.2, label="DCI-P3"),
            mlines.Line2D([0], [0], color="#aab0c0", linestyle="-",
                          linewidth=1.2, label="BT.2020"),
        ]
        fig.legend(handles=ref_handles, loc="lower center", ncol=2,
                   fontsize=10, bbox_to_anchor=(0.5, 0.01),
                   frameon=True, edgecolor="#c0c0c0")

        # BT.2020 꼭짓점 pixel fraction 계산 (줌 커트용)
        # transData.transform()은 canvas DPI(fig.dpi) 기준 픽셀 반환 →
        # fig.dpi로 나눠야 fraction 정확 (저장 DPI와 다름)
        fig.canvas.draw()
        _canvas_dpi = fig.dpi                   # canvas 실제 DPI (≈100)
        _w_px = _FW * _canvas_dpi
        _h_px = _FH * _canvas_dpi
        _bt2020_verts = {
            "red":   (BT2020_UV[0][0], BT2020_UV[0][1]),
            "green": (BT2020_UV[1][0], BT2020_UV[1][1]),
            "blue":  (BT2020_UV[2][0], BT2020_UV[2][1]),
        }
        vfracs: dict = {}
        for _c, (_u, _v) in _bt2020_verts.items():
            _px, _py = ax.transData.transform((_u, _v))
            vfracs[_c] = (_px / _w_px, 1.0 - _py / _h_px)  # y: 위=0 (PPT 기준)

        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        tmp.close()
        fig.savefig(tmp.name, dpi=_DPI)  # bbox_inches 없음 → 정확히 900×900 px
        plt.close(fig)
        tmp_files.append(tmp.name)
        return tmp.name, vfracs

    sdr_modes = [k for k in engine.session_loading if "SDR" in k]
    hdr_modes = [k for k in engine.session_loading if "HDR" in k]

    img_sdr = _make_apl_chart("SDR APL Loading", sdr_modes)
    img_hdr = _make_apl_chart("HDR APL Loading", hdr_modes)
    img_gmt, _gmt_vfracs = _make_gamut_chart()  # (경로, BT.2020 꼭짓점 fraction)

    # ── PPT 생성 ─────────────────────────────────────────────────────────
    try:
        prs = Presentation()
        SW, SH = 13.33, 7.5
        prs.slide_width  = Inches(SW)
        prs.slide_height = Inches(SH)

        BLANK = 6
        for i, lay in enumerate(prs.slide_layouts):
            if "blank" in lay.name.lower():
                BLANK = i; break
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK])

        def _hex(h: str) -> RGBColor:
            h = h.lstrip("#")
            return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        def _tb(text: str, x: float, y: float, w: float, h: float,
                size: int = 10, bold: bool = False,
                color: str = "000000", align=PP_ALIGN.LEFT) -> None:
            tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = tb.text_frame; tf.word_wrap = False
            para = tf.paragraphs[0]; para.alignment = align
            run = para.add_run(); run.text = text
            run.font.size = Pt(size); run.font.bold = bold
            run.font.color.rgb = _hex(color)

        def _set_border(cell) -> None:
            """셀 4면 검정 1pt 테두리 — Windows PowerPoint 호환 (prstDash 필수)."""
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for side in ("lnL", "lnR", "lnT", "lnB"):
                tag = _qn(f"a:{side}")
                for e in tcPr.findall(tag):
                    tcPr.remove(e)
                ln = _et.SubElement(tcPr, tag, w="12700")
                sf = _et.SubElement(ln, _qn("a:solidFill"))
                _et.SubElement(sf, _qn("a:srgbClr"), val="000000")
                _et.SubElement(ln, _qn("a:prstDash"), val="solid")

        # ① 제목 + 구분선
        TITLE_H = 0.40
        LINE_Y  = TITLE_H + 0.08
        _tb(f"{brand}  {model}", 0.02, 0.04, SW * 0.6, TITLE_H,
            size=18, bold=True, color="000000", align=PP_ALIGN.LEFT)
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            Inches(0.0), Inches(LINE_Y), Inches(SW), Inches(LINE_Y))
        conn.line.color.rgb = _hex("000000")  # type: ignore[assignment]
        conn.line.width = Pt(0.5)

        # ② 데이터 테이블
        _EXTRA_ROWS = ["반사율[%]  (SCI/SCE)", "Local Dimming 개수  [block]"]
        n_data_rows = len(ROW_LABELS)
        n_rows = n_data_rows + 1 + len(_EXTRA_ROWS)
        n_cols = 4   # 구분 | 항목-그룹 | 항목-세부 | 모델

        CAT_W = 1.1; GRP_W = 0.85; DTL_W = 0.95; MDL_W = 1.9
        TBL_W = CAT_W + GRP_W + DTL_W + MDL_W
        TBL_X = (SW - TBL_W) / 2
        TBL_Y = LINE_Y + 0.12
        ROW_H = 0.27
        TBL_H = n_rows * ROW_H

        tbl_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(TBL_X), Inches(TBL_Y), Inches(TBL_W), Inches(TBL_H))
        tbl = tbl_shape.table
        tbl.columns[0].width = Inches(CAT_W)
        tbl.columns[1].width = Inches(GRP_W)
        tbl.columns[2].width = Inches(DTL_W)
        tbl.columns[3].width = Inches(MDL_W)

        # Windows PowerPoint 기본 테이블 스타일이 셀 테두리를 덮어씀 → 제거
        _tblPr = tbl._tbl.find(_qn("a:tblPr"))
        if _tblPr is None:
            _tblPr = _et.SubElement(tbl._tbl, _qn("a:tblPr"))
        _styleId = _tblPr.find(_qn("a:tableStyleId"))
        if _styleId is not None:
            _tblPr.remove(_styleId)

        def _cell_write(row: int, col: int, text: str, bold: bool = False,
                        size: int = 12, align=PP_ALIGN.CENTER,
                        color: str = "000000") -> None:
            c = tbl.cell(row, col)
            c.text = text; c.fill.background()
            tf = c.text_frame; tf.paragraphs[0].alignment = align
            run = (tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs
                   else tf.paragraphs[0].add_run())
            run.font.size = Pt(size); run.font.bold = bold
            run.font.color.rgb = _hex(color)

        def _style_merged(mc, text: str, bold: bool = True,
                          align=PP_ALIGN.CENTER) -> None:
            mc.text = text; mc.fill.background()
            tf = mc.text_frame; tf.paragraphs[0].alignment = align
            run = (tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs
                   else tf.paragraphs[0].add_run())
            run.font.size = Pt(12); run.font.bold = bold
            run.font.color.rgb = _hex("000000")

        # 헤더 행
        _cell_write(0, 0, "구분", bold=True)
        tbl.cell(0, 1).merge(tbl.cell(0, 2))
        _style_merged(tbl.cell(0, 1), "항목", align=PP_ALIGN.CENTER)
        _cell_write(0, 3, f"{brand}_{model}", bold=True)

        # 데이터 행
        for ri, ((section, item), val) in enumerate(zip(ROW_LABELS, row_values), 1):
            _cell_write(ri, 0, section, bold=True)
            split = SPLIT_MAP.get(item)
            if split:
                grp, dtl = split
                _cell_write(ri, 1, grp, bold=True)
                _cell_write(ri, 2, dtl)
            else:
                tbl.cell(ri, 1).merge(tbl.cell(ri, 2))
                _style_merged(tbl.cell(ri, 1), item, bold=False, align=PP_ALIGN.LEFT)
            _cell_write(ri, 3, val)

        # 구분(col 0) 세로 병합
        idx = 0
        while idx < n_data_rows:
            section = ROW_LABELS[idx][0]
            end = idx + 1
            while end < n_data_rows and ROW_LABELS[end][0] == section:
                end += 1
            if end - idx > 1:
                tbl.cell(idx + 1, 0).merge(tbl.cell(end, 0))
                _style_merged(tbl.cell(idx + 1, 0), section)
            idx = end

        # col 1 HDR/SDR 세로 병합
        for prefix, label in (("HDR", "HDR vivid"), ("SDR", "SDR vivid")):
            grp_rows = [ri + 1 for ri, (_, item) in enumerate(ROW_LABELS)
                        if item.startswith(prefix + " ")]
            if len(grp_rows) > 1:
                tbl.cell(grp_rows[0], 1).merge(tbl.cell(grp_rows[-1], 1))
                _style_merged(tbl.cell(grp_rows[0], 1), label)

        # 하단 추가 행
        for ei, label in enumerate(_EXTRA_ROWS):
            extra_row = n_data_rows + 1 + ei
            _cell_write(extra_row, 3, "")
            tbl.cell(extra_row, 0).merge(tbl.cell(extra_row, 2))
            _style_merged(tbl.cell(extra_row, 0), label, bold=True, align=PP_ALIGN.LEFT)

        # 전체 셀 테두리 + 12pt 강제
        for ri in range(n_rows):
            for ci in range(n_cols):
                _set_border(tbl.cell(ri, ci))
                try:
                    tf3 = tbl.cell(ri, ci).text_frame
                    for para in tf3.paragraphs:
                        if not para.runs:
                            run = para.add_run()
                            run.text = " "; run.font.size = Pt(12)
                        else:
                            for run in para.runs:
                                if run.text == "":
                                    run.text = " "
                                run.font.size = Pt(12)
                except Exception:
                    pass

        # ③ 차트 이미지 삽입
        CHART_AREA_Y = TBL_Y + TBL_H + 0.7    # 표 아래 0.7인치 — 데스크탑과 동일
        CHART_H = min(3.0, SH - CHART_AREA_Y - 0.05)
        CHART_W = SW / 3.0

        def _add_gamut_zoom_cuts(img_path: str, gx: float, gy: float,
                                  gside: float, slot_w: float,
                                  vfracs: dict) -> None:
            """BT.2020 꼭짓점 확대 컷 3개 — 데스크탑과 동일 로직."""
            zoom_side   = min(0.95, gside * 0.38)  # 줌 컷 크기: 차트의 38%, 최대 0.95인치
            zoom_margin = 0.04                      # 슬롯 모서리 여백 (인치)
            half        = 0.13                      # 크롭 반경: 꼭짓점 ±13% (클수록 줌 배율 감소)
            corners = [
                (gx + zoom_margin,                      gy + zoom_margin,                     "green"),
                (gx + slot_w - zoom_side - zoom_margin,  gy + zoom_margin,                     "red"),
                (gx + zoom_margin,                      gy + gside - zoom_side - zoom_margin,  "blue"),
            ]
            for zx, zy, color in corners:
                if color not in vfracs:
                    continue
                xc, yc = vfracs[color]
                crop_l = max(0.0, xc - half)
                crop_r = max(0.0, 1.0 - xc - half)
                crop_t = max(0.0, yc - half)
                crop_b = max(0.0, 1.0 - yc - half)
                pic = slide.shapes.add_picture(
                    img_path, Inches(zx), Inches(zy),
                    Inches(zoom_side), Inches(zoom_side))
                pic.crop_left   = crop_l
                pic.crop_top    = crop_t
                pic.crop_right  = crop_r
                pic.crop_bottom = crop_b

        for chart_idx, img_path in enumerate([img_sdr, img_hdr, img_gmt]):
            cx = chart_idx * CHART_W
            if chart_idx == 2:
                gmt_side = min(CHART_W - 0.06, CHART_H)  # 정사각형 유지
                gmt_off  = (CHART_W - gmt_side) / 2       # 좌우 중앙 정렬
                slide.shapes.add_picture(
                    img_path, Inches(cx + gmt_off), Inches(CHART_AREA_Y),
                    Inches(gmt_side), Inches(gmt_side))
                _add_gamut_zoom_cuts(
                    img_path, cx, CHART_AREA_Y, gmt_side, CHART_W, _gmt_vfracs)
            else:
                slide.shapes.add_picture(
                    img_path, Inches(cx + 0.04), Inches(CHART_AREA_Y),
                    Inches(CHART_W - 0.08), Inches(CHART_H))

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return send_file(
            output,
            as_attachment=True,
            download_name=f"report_{brand}_{model}.pptx",
            mimetype=(
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation"
            ),
        )

    finally:
        for p in tmp_files:
            try:
                os.unlink(p)
            except OSError:
                pass


@bp.route("/export/<export_type>", methods=["GET"])
def export_results(export_type: str):
    """Download an Excel file for the requested result type.

    export_type: "lum_swing" | "lum_loading" | "gamut" | "contrast"
    """
    engine = _get_engine()
    exporter = ExcelExporter()
    use_avg = request.args.get("use_avg", "true").lower() == "true"

    try:
        if export_type == "lum_swing":
            data = _last_results.get("lum_swing")
            if not data:
                return jsonify({"error": "No lum_swing results available"}), 404
            path = exporter.export_lum_swing(data, engine.brand, engine.model_name)

        elif export_type == "lum_loading":
            data = _last_results.get("lum_loading")
            if not data:
                return jsonify({"error": "No lum_loading results available"}), 404
            path = exporter.export_lum_loading(data, engine.brand, engine.model_name, use_avg=use_avg)

        elif export_type == "gamut":
            data = _last_results.get("gamut")
            if not data:
                return jsonify({"error": "No gamut results available"}), 404
            path = exporter.export_gamut(data, engine.brand, engine.model_name)

        elif export_type == "contrast":
            data = _last_results.get("contrast")
            if not data:
                return jsonify({"error": "No contrast results available"}), 404
            path = exporter.export_contrast(data, engine.brand, engine.model_name)

        else:
            return jsonify({"error": f"Unknown export type: {export_type}"}), 400

    except Exception as exc:
        return jsonify({"error": str(exc)}), 500

    return send_file(path, as_attachment=True, download_name=path.split("/")[-1])
